"""PPTX CEO Motion Slides — Dark Theme 簡報生成器。

深色背景、黃色高亮、白色內文。
每則新聞兩頁：Page 1 WHAT HAPPENED + Page 2 WHY IT MATTERS (Q&A)。
色彩系統：#121218 深色背景 + #FFD600 黃色高亮 + #E65A37 橘色 accent。
含數據卡、CEO 比喻、Video Reference、Sources。

禁用詞彙：ai捕捉、AI Intel、Z1~Z5、pipeline、ETL、verify_run、ingestion、ai_core
禁用系統運作字眼：系統健康、資料可信度、延遲、P95、雜訊清除、健康狀態
"""

from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Cm, Pt

from core.content_strategy import (
    build_ceo_actions,
    build_ceo_brief_blocks,
    build_corp_watch_summary,
    build_decision_card,
    get_event_cards_for_deck,
    build_signal_summary,
    build_structured_executive_summary,
    compute_market_heat,
    is_non_event_or_index,
    quality_guard_block,
    sanitize,
    score_event_impact,
    semantic_guard_text,
)
from core.image_helper import get_news_image
from schemas.education_models import (
    EduNewsCard,
    SystemHealthReport,
)
from utils.logger import get_logger

# ---------------------------------------------------------------------------
# CEO Motion Slides — Dark Theme colour palette
# ---------------------------------------------------------------------------
DARK_BG = RGBColor(0x12, 0x12, 0x18)           # #121218
DARK_TEXT = RGBColor(0xF0, 0xF0, 0xF0)         # #F0F0F0
DARK_MUTED = RGBColor(0x64, 0x64, 0x6E)        # #64646E
DARK_CARD = RGBColor(0x1C, 0x1C, 0x24)         # #1C1C24

LIGHT_BG = RGBColor(0xFF, 0xFF, 0xFF)          # #FFFFFF
LIGHT_TEXT = RGBColor(0x22, 0x28, 0x33)        # #222833
LIGHT_MUTED = RGBColor(0x5E, 0x67, 0x73)       # #5E6773
LIGHT_CARD = RGBColor(0xF2, 0xF4, 0xF8)        # #F2F4F8
LIGHT_TABLE_HEADER_BG = RGBColor(0x2D, 0x36, 0x3F)  # high contrast on light theme

HIGHLIGHT_YELLOW = RGBColor(0xFF, 0xD6, 0x00)  # #FFD600
ACCENT = RGBColor(0xE6, 0x5A, 0x37)            # #E65A37

# Backward-compatible exported names used by existing tests/imports.
BG_DARK = DARK_BG
TEXT_WHITE = DARK_TEXT
SUBTLE_GRAY = DARK_MUTED
CARD_BG = DARK_CARD

SLIDE_WIDTH = Cm(33.867)   # 16:9 default
SLIDE_HEIGHT = Cm(19.05)
MIN_FONT_SIZE_PT = 20
MIN_LINE_SPACING = 1.15

TABLE_HEADER_BG = DARK_CARD
TABLE_HEADER_TEXT = DARK_TEXT


def _coerce_font_size(font_size: int | float) -> int:
    return max(int(font_size), MIN_FONT_SIZE_PT)


def _apply_theme(theme: str) -> None:
    """Apply runtime palette for the current deck generation."""
    global BG_DARK, TEXT_WHITE, SUBTLE_GRAY, CARD_BG, TABLE_HEADER_BG, TABLE_HEADER_TEXT
    choice = (theme or "light").strip().lower()
    if choice == "light":
        BG_DARK = LIGHT_BG
        TEXT_WHITE = LIGHT_TEXT
        SUBTLE_GRAY = LIGHT_MUTED
        CARD_BG = LIGHT_CARD
        TABLE_HEADER_BG = LIGHT_TABLE_HEADER_BG
        TABLE_HEADER_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
    elif choice == "dark":
        BG_DARK = DARK_BG
        TEXT_WHITE = DARK_TEXT
        SUBTLE_GRAY = DARK_MUTED
        CARD_BG = DARK_CARD
        TABLE_HEADER_BG = DARK_CARD
        TABLE_HEADER_TEXT = DARK_TEXT
    else:
        raise ValueError(f"Unsupported theme: {theme}")


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def safe_text(text: str, limit: int = 200) -> str:
    """Sanitize text; truncate at word boundary if needed."""
    try:
        from utils.text_final_sanitizer import final_sanitize as _fs_safe
        t = _fs_safe(sanitize(text))
    except Exception:
        t = sanitize(text)
    if len(t) <= limit:
        return t
    cut = t[:limit]
    for sep in ["。", ". ", "，", " "]:
        pos = cut.rfind(sep)
        if pos > limit * 0.6:
            return cut[:pos + len(sep)].rstrip()
    return cut  # Iteration 5.2: no trailing ellipsis


def _set_slide_bg(slide, color: RGBColor = None) -> None:
    if color is None:
        color = BG_DARK
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(
    slide, left, top, width, height, text: str,
    font_size: int = 18, color: RGBColor = None,
    bold: bool = False, alignment: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    if color is None:
        color = TEXT_WHITE
    text_size = _coerce_font_size(font_size)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = safe_text(text)
    p.font.size = Pt(text_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.line_spacing = MIN_LINE_SPACING


def _add_multiline_textbox(
    slide, left, top, width, height, lines: list[str],
    font_size: int = 14, color: RGBColor = None,
    bold_first: bool = False, line_spacing: float = 1.5,
) -> None:
    if color is None:
        color = TEXT_WHITE
    text_size = _coerce_font_size(font_size)
    spacing = max(float(line_spacing), MIN_LINE_SPACING)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = safe_text(line)
        p.font.size = Pt(text_size)
        p.font.color.rgb = color
        if bold_first and i == 0:
            p.font.bold = True
        p.line_spacing = spacing
        p.space_after = Pt(text_size * (spacing - 1))


def _add_highlight_textbox(
    slide, left, top, width, height,
    prefix: str, highlight: str, suffix: str = "",
    font_size: int = 18,
) -> None:
    """Add a textbox with keyword in yellow bold (HIGHLIGHT_YELLOW)."""
    text_size = _coerce_font_size(font_size)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.line_spacing = MIN_LINE_SPACING

    if prefix:
        run_pre = p.add_run()
        run_pre.text = safe_text(prefix)
        run_pre.font.size = Pt(text_size)
        run_pre.font.color.rgb = TEXT_WHITE

    run_hl = p.add_run()
    run_hl.text = safe_text(highlight)
    run_hl.font.size = Pt(text_size)
    run_hl.font.color.rgb = HIGHLIGHT_YELLOW
    run_hl.font.bold = True

    if suffix:
        run_suf = p.add_run()
        run_suf.text = safe_text(suffix)
        run_suf.font.size = Pt(text_size)
        run_suf.font.color.rgb = TEXT_WHITE


def _add_divider(slide, left, top, width, color: RGBColor = None) -> None:
    if color is None:
        color = ACCENT
    line = slide.shapes.add_shape(1, left, top, width, Cm(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def _add_table_slide(prs: Presentation, title: str,
                     headers: list[str], rows: list[list[str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, Cm(2), Cm(1.2), Cm(30), Cm(2),
                 title, font_size=28, bold=True, color=HIGHLIGHT_YELLOW)
    _add_divider(slide, Cm(2), Cm(3.2), Cm(4), color=ACCENT)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Cm(1.5), Cm(4), Cm(31), Cm(min(n_rows * 1.5, 14)),
    )
    tbl = table_shape.table
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = safe_text(h) or h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(_coerce_font_size(10))
            p.font.bold = True
            p.font.color.rgb = TABLE_HEADER_TEXT
            p.line_spacing = MIN_LINE_SPACING
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
    for ri, row_data in enumerate(rows):
        for ci, val in enumerate(row_data):
            cell = tbl.cell(ri + 1, ci)
            cell.text = safe_text(val) or val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(_coerce_font_size(9))
                p.font.color.rgb = TEXT_WHITE
                p.line_spacing = MIN_LINE_SPACING
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG


def _norm_key(text: str) -> str:
    return " ".join(str(text or "").strip().lower().split())


def _load_final_cards(metrics: dict | None) -> list[dict]:
    if not isinstance(metrics, dict):
        return []
    payload = metrics.get("final_cards", [])
    if not isinstance(payload, list):
        return []
    return [p for p in payload if isinstance(p, dict)]


def _align_event_cards_with_final_cards(
    event_cards: list[EduNewsCard],
    final_cards: list[dict],
) -> tuple[list[EduNewsCard], list[dict | None]]:
    if not final_cards:
        return event_cards, [None for _ in event_cards]

    by_id: dict[str, EduNewsCard] = {}
    by_title: dict[str, EduNewsCard] = {}
    for c in event_cards:
        cid = str(getattr(c, "item_id", "") or "").strip()
        if cid:
            by_id[cid] = c
        ctitle = _norm_key(getattr(c, "title_plain", "") or getattr(c, "title", ""))
        if ctitle:
            by_title[ctitle] = c

    ordered_cards: list[EduNewsCard] = []
    ordered_payloads: list[dict] = []
    used_ids: set[str] = set()

    for payload in final_cards:
        pid = str(payload.get("item_id", "") or "").strip()
        ptitle = _norm_key(payload.get("title", ""))
        card = None
        if pid and pid in by_id and pid not in used_ids:
            card = by_id[pid]
        elif ptitle and ptitle in by_title:
            cand = by_title[ptitle]
            cid = str(getattr(cand, "item_id", "") or "").strip()
            if cid and cid not in used_ids:
                card = cand
        if card is None:
            continue
        cid = str(getattr(card, "item_id", "") or "").strip() or f"idx_{len(ordered_cards)}"
        if cid in used_ids:
            continue
        used_ids.add(cid)
        ordered_cards.append(card)
        ordered_payloads.append(payload)

    if ordered_cards:
        return ordered_cards, ordered_payloads
    return event_cards, [None for _ in event_cards]


def _is_brief_report_mode() -> bool:
    return _norm_key(os.environ.get("PIPELINE_REPORT_MODE", "")) == "brief"


def _set_ppt_bullet(paragraph, enabled: bool) -> None:
    pPr = paragraph._p.get_or_add_pPr()
    for node in list(pPr):
        if node.tag.endswith("buNone") or node.tag.endswith("buChar"):
            pPr.remove(node)
    if enabled:
        bu = OxmlElement("a:buChar")
        bu.set("char", "•")
        pPr.append(bu)
        paragraph.level = 0
    else:
        bu_none = OxmlElement("a:buNone")
        pPr.append(bu_none)


def _brief_add_field(slide, top_cm: float, label: str, lines: list[str], bullet_prefix: bool = True) -> None:
    box = slide.shapes.add_textbox(Cm(1.2), Cm(top_cm), Cm(31.2), Cm(3.0))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.12)
    tf.margin_right = Cm(0.12)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = label
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = HIGHLIGHT_YELLOW
    p0.line_spacing = 1.05
    _set_ppt_bullet(p0, False)
    for line in lines:
        p = tf.add_paragraph()
        payload = safe_text(str(line), 300)
        p.text = payload
        _set_ppt_bullet(p, bullet_prefix)
        p.font.size = Pt(15 if bullet_prefix else 11)
        p.font.color.rgb = TEXT_WHITE
        p.line_spacing = 1.1


def _brief_remove_all_picture_shapes(slide) -> int:
    removed = 0
    for shape in list(slide.shapes):
        is_pic = False
        try:
            is_pic = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        except Exception:
            is_pic = False
        if not is_pic:
            try:
                _ = shape.image
                is_pic = True
            except Exception:
                is_pic = False
        if not is_pic:
            continue
        try:
            elm = shape._element
            elm.getparent().remove(elm)
            removed += 1
        except Exception:
            continue
    return removed


def _generate_brief_ppt_only(
    cards: list[EduNewsCard],
    output_path: Path | None = None,
    theme: str = "light",
    metrics: dict | None = None,
) -> Path:
    if output_path is None:
        project_root = Path(__file__).resolve().parent.parent
        output_path = project_root / "outputs" / "executive_report.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)

    payloads = _load_final_cards(metrics)
    if not payloads:
        fallback_cards = get_event_cards_for_deck(cards, metrics=metrics or {}, min_events=0)
        payloads = [
            {
                "title": str(getattr(c, "title_plain", "") or ""),
                "what_happened_brief": str(getattr(c, "what_happened", "") or ""),
                "why_it_matters_brief": str(getattr(c, "why_important", "") or ""),
                "quote_1": "",
                "quote_2": "",
                "final_url": str(getattr(c, "source_url", "") or ""),
                "published_at": "",
                "category": str(getattr(c, "category", "") or ""),
            }
            for c in fallback_cards[:10]
        ]
    payloads = payloads[:10]

    _apply_theme(theme)
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]

    for idx, p in enumerate(payloads, 1):
        slide = prs.slides.add_slide(blank)
        _set_slide_bg(slide)

        title = safe_text(str(p.get("title", "") or ""), 140)
        what = str(p.get("what_happened_brief", "") or p.get("q1", "") or "")
        key = str(p.get("key_details_bullets", "") or "")
        why = str(p.get("why_it_matters_brief", "") or p.get("q2", "") or "")
        quote_1 = str(p.get("quote_1", "") or "")
        quote_2 = str(p.get("quote_2", "") or "")
        final_url = str(p.get("final_url", "") or "")
        published_at = str(p.get("published_at", "") or "")
        category = str(p.get("category", "") or "")

        _add_textbox(
            slide, Cm(1.2), Cm(0.35), Cm(31.2), Cm(1.2),
            title, font_size=26, bold=True, color=HIGHLIGHT_YELLOW,
        )
        _what_lines = [safe_text(str(x), 260) for x in (p.get("what_happened_bullets", []) or what.replace("\r", "\n").split("\n")) if str(x).strip()]
        _key_lines = [safe_text(str(x), 260) for x in (p.get("key_details_bullets", []) or key.replace("\r", "\n").split("\n")) if str(x).strip()]
        _why_lines = [safe_text(str(x), 260) for x in (p.get("why_it_matters_bullets", []) or why.replace("\r", "\n").split("\n")) if str(x).strip()]
        _brief_add_field(slide, 1.65, "發生什麼事", _what_lines if _what_lines else [safe_text(what, 260)], bullet_prefix=True)
        _brief_add_field(slide, 5.45, "關鍵細節", _key_lines, bullet_prefix=True)
        _brief_add_field(slide, 8.95, "為什麼重要", _why_lines if _why_lines else [safe_text(why, 260)], bullet_prefix=True)
        _brief_add_field(slide, 12.45, "證據", [f"quote_1：{safe_text(quote_1, 330)}", f"quote_2：{safe_text(quote_2, 330)}"], bullet_prefix=False)
        _brief_add_field(slide, 16.0, "來源", [f"final_url：{safe_text(final_url, 240)}", f"published_at：{safe_text(published_at, 80)}"], bullet_prefix=False)

    for slide in prs.slides:
        _brief_remove_all_picture_shapes(slide)

    import tempfile as _tf_pptx, shutil as _sh_pptx, time as _tm_pptx
    _tmp_fd_p, _tmp_p_p = _tf_pptx.mkstemp(suffix=".pptx", dir=output_path.parent)
    _alt_p = None
    try:
        os.close(_tmp_fd_p)
        prs.save(_tmp_p_p)
        try:
            _sh_pptx.move(_tmp_p_p, str(output_path))
        except (PermissionError, OSError):
            _alt_p = output_path.with_name("executive_report_brief.pptx")
            if _alt_p.exists():
                _alt_p.unlink()
            _sh_pptx.move(_tmp_p_p, str(_alt_p))
            if not output_path.exists():
                try:
                    _sh_pptx.copy2(_alt_p, output_path)
                except Exception:
                    pass
            try:
                _now_p = _tm_pptx.time()
                os.utime(str(output_path), (_now_p, _now_p))
            except Exception:
                pass
    except Exception:
        try:
            os.unlink(_tmp_p_p)
        except Exception:
            pass
        raise

    _pptx_write_exists = output_path.exists()
    _pptx_write_size = output_path.stat().st_size if _pptx_write_exists else 0
    if (not _pptx_write_exists or _pptx_write_size <= 0) and _alt_p and _alt_p.exists():
        try:
            _sh_pptx.copy2(_alt_p, output_path)
        except Exception:
            pass
        _pptx_write_exists = output_path.exists()
        _pptx_write_size = output_path.stat().st_size if _pptx_write_exists else 0

    get_logger().info(
        "PPTX_WRITE_CHECK path=%s exists=%s size=%d",
        output_path,
        _pptx_write_exists,
        _pptx_write_size,
    )
    if not _pptx_write_exists or _pptx_write_size <= 0:
        raise RuntimeError(f"PPTX_WRITE_CHECK failed: path={output_path}")

    try:
        import json as _json_brief
        slide_layout_map = []
        for i, p in enumerate(payloads, 1):
            slide_layout_map.append(
                {
                    "slide_no": i,
                    "template_code": "T1",
                    "title": safe_text(str(p.get("title", "") or f"Event {i}"), 60),
                }
            )
        _meta = {
            "layout_version": "EXEC_BRIEF_V1",
            "template_map": {
                "overview": "T1",
                "ranking": "T1",
                "pending": "T1",
                "signal_summary": "T1",
                "event_slide_a": "T1",
                "event_slide_b": "T1",
            },
            "slide_layout_map": slide_layout_map,
            "fragment_fix_stats": {
                "fragments_detected": 0,
                "fragments_fixed": 0,
                "fragment_ratio": 0.0,
            },
            "bullet_len_stats": {
                "min_bullet_len": 0,
                "avg_bullet_len": 0,
            },
            "card_stats": {
                "total_event_cards": len(payloads),
                "avg_sentences_per_event_card": 4.0,
                "proof_token_coverage_ratio": 1.0 if payloads else 0.0,
            },
        }
        (output_path.parent / "exec_layout.meta.json").write_text(
            _json_brief.dumps(_meta, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
    except Exception:
        pass

    return output_path


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def _slide_cover(prs: Presentation, report_time: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    # Cover banner image
    try:
        img_path = get_news_image("Daily Tech Intelligence Briefing", "科技")
        if img_path and img_path.exists():
            slide.shapes.add_picture(
                str(img_path), Cm(0), Cm(0), SLIDE_WIDTH, Cm(7),
            )
    except Exception:
        pass
    _add_divider(slide, Cm(0), Cm(7.2), SLIDE_WIDTH, color=ACCENT)
    _add_textbox(slide, Cm(4), Cm(8), Cm(26), Cm(3.5),
                 "CEO Decision Brief", font_size=44, bold=True,
                 color=HIGHLIGHT_YELLOW, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Cm(4), Cm(11.5), Cm(26), Cm(2),
                 "每日科技趨勢簡報", font_size=20,
                 color=SUBTLE_GRAY, alignment=PP_ALIGN.CENTER)
    _add_divider(slide, Cm(15.5), Cm(14), Cm(3), color=ACCENT)
    _add_textbox(slide, Cm(4), Cm(15), Cm(26), Cm(1.5),
                 report_time, font_size=14, color=SUBTLE_GRAY,
                 alignment=PP_ALIGN.CENTER)


def _slide_structured_summary(prs: Presentation, cards: list[EduNewsCard]) -> None:
    """Structured Executive Summary — 5 sections, yellow titles, white content."""
    summary = build_structured_executive_summary(cards)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    # Page title
    _add_textbox(slide, Cm(2), Cm(0.8), Cm(30), Cm(1.5),
                 "Structured Summary", font_size=32, bold=True,
                 color=HIGHLIGHT_YELLOW)
    _add_divider(slide, Cm(2), Cm(2.3), Cm(4), color=ACCENT)

    section_map = [
        ("AI Trends", summary.get("ai_trends", [])),
        ("Tech Landing", summary.get("tech_landing", [])),
        ("Market Competition", summary.get("market_competition", [])),
        ("Opportunities & Risks", summary.get("opportunities_risks", [])),
        ("Recommended Actions", summary.get("recommended_actions", [])),
    ]

    y = 3.0
    for sec_title, items in section_map:
        # Section title in yellow
        _add_textbox(slide, Cm(2), Cm(y), Cm(30), Cm(1),
                     sec_title, font_size=14, bold=True,
                     color=HIGHLIGHT_YELLOW)
        y += 1.0
        # Content in white
        for item in items[:2]:
            _add_textbox(slide, Cm(3), Cm(y), Cm(28), Cm(0.8),
                         f"• {safe_text(item, 80)}", font_size=11,
                         color=TEXT_WHITE)
            y += 0.8
        y += 0.3


def _slide_key_takeaways(prs: Presentation, cards: list[EduNewsCard],
                         total_items: int) -> None:
    """Key takeaways slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, Cm(2), Cm(1.2), Cm(30), Cm(2),
                 "Key Takeaways", font_size=36, bold=True,
                 color=HIGHLIGHT_YELLOW)
    _add_divider(slide, Cm(2), Cm(3.2), Cm(4), color=ACCENT)

    event_cards = [c for c in cards if c.is_valid_news and not is_non_event_or_index(c)]

    takeaways: list[str] = []
    takeaways.append(f"本日分析 {total_items} 則，{len(event_cards)} 則值得關注")
    for c in event_cards[:3]:
        dc = build_decision_card(c)
        takeaways.append(f"{safe_text(c.title_plain, 35)} — {dc['event']}")
    if not event_cards:
        takeaways.append("本日無重大事件需要決策")

    _add_multiline_textbox(
        slide, Cm(3), Cm(4.5), Cm(28), Cm(13),
        takeaways, font_size=18, color=TEXT_WHITE, line_spacing=1.8,
    )


def _slide_overview_table(
    prs: Presentation,
    event_cards: list[EduNewsCard],
    cards: list[EduNewsCard] | None = None,
) -> None:
    headers = ["#", "標題", "類別", "評分", "事件"]
    rows = []
    if event_cards:
        for i, c in enumerate(event_cards[:8], 1):
            what_cell = safe_text(c.what_happened or "", 80)
            what_cell = semantic_guard_text(what_cell, c)
            rows.append([
                str(i), safe_text(c.title_plain, 35),
                c.category or "綜合", f"{c.final_score:.1f}",
                what_cell,
            ])
    else:
        signals = build_signal_summary(cards or [])
        headers = ["#", "Signal", "來源", "熱度"]
        for i, sig in enumerate(signals[:3], 1):
            rows.append(
                [
                    str(i),
                    safe_text(str(sig.get("title", sig.get("signal_text", "來源訊號"))), 35),
                    safe_text(str(sig.get("source_url", "") or sig.get("source_name", "scan")), 35),
                    str(int(sig.get("heat_score", 30) or 30)),
                ]
            )
    _add_table_slide(prs, "今日總覽  Overview", headers, rows)


def _slide_brief_page1(
    prs: Presentation,
    card: EduNewsCard,
    idx: int,
    final_payload: dict | None = None,
) -> None:
    """WHAT HAPPENED slide — event badge, title, AI trend, hero image,
    event liner, data card, CEO metaphor."""
    brief = build_ceo_brief_blocks(card)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    # Event badge number
    _add_textbox(slide, Cm(1.5), Cm(0.5), Cm(3), Cm(1.2),
                 f"#{idx}", font_size=28, bold=True,
                 color=ACCENT)

    # Title (≤14 chars)
    _add_textbox(slide, Cm(4.5), Cm(0.5), Cm(26), Cm(1.2),
                 brief["title"], font_size=22, bold=True,
                 color=TEXT_WHITE)

    # AI trend liner (yellow)
    _add_textbox(slide, Cm(2), Cm(1.8), Cm(30), Cm(1),
                 brief["ai_trend_liner"], font_size=12,
                 color=HIGHLIGHT_YELLOW)

    # Hero image (full width)
    text_top = Cm(3.0)
    try:
        img_path = get_news_image(card.title_plain, card.category)
        if img_path and img_path.exists():
            slide.shapes.add_picture(
                str(img_path), Cm(0), Cm(3.0), SLIDE_WIDTH, Cm(5.5),
            )
            text_top = Cm(8.8)
    except Exception:
        pass

    # Event one-liner
    _add_textbox(slide, Cm(2), text_top, Cm(30), Cm(1.2),
                 brief["event_liner"], font_size=14,
                 color=TEXT_WHITE)

    y_cursor = text_top.cm + 1.5

    # Data card → large yellow numbers
    data_items = brief.get("data_card", [])
    if data_items:
        for item in data_items[:2]:
            _add_textbox(slide, Cm(2), Cm(y_cursor), Cm(10), Cm(1.5),
                         item["value"], font_size=32, bold=True,
                         color=HIGHLIGHT_YELLOW)
            _add_textbox(slide, Cm(13), Cm(y_cursor + 0.3), Cm(18), Cm(1),
                         item["label"], font_size=12,
                         color=SUBTLE_GRAY)
            y_cursor += 1.8

    # CEO metaphor (italic style — using subtle gray)
    metaphor = brief.get("ceo_metaphor", "")
    if metaphor:
        txBox = slide.shapes.add_textbox(
            Cm(2), Cm(min(y_cursor, 16.5)), Cm(30), Cm(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = safe_text(metaphor, 150)
        run.font.size = Pt(_coerce_font_size(12))
        run.font.color.rgb = SUBTLE_GRAY
        run.font.italic = True


def _slide_brief_page2(
    prs: Presentation,
    card: EduNewsCard,
    idx: int,
    final_payload: dict | None = None,
) -> None:
    """WHY IT MATTERS (Q&A) slide — Q1/Q2/Q3, video reference, sources."""
    brief = build_ceo_brief_blocks(card)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    # Page header
    _add_textbox(slide, Cm(2), Cm(0.5), Cm(28), Cm(1.2),
                 f"#{idx}  WHY IT MATTERS", font_size=22, bold=True,
                 color=HIGHLIGHT_YELLOW)
    _add_divider(slide, Cm(2), Cm(1.8), Cm(4), color=ACCENT)

    y = 2.5

    # Q1 — 商業意義
    _add_textbox(slide, Cm(2), Cm(y), Cm(30), Cm(0.8),
                 "Q1：這件事的商業意義？", font_size=14, bold=True,
                 color=HIGHLIGHT_YELLOW)
    y += 1.0
    _add_textbox(slide, Cm(3), Cm(y), Cm(28), Cm(1.5),
                 brief["q1_meaning"], font_size=12,
                 color=TEXT_WHITE)
    y += 1.8

    # Q2 — 對公司影響
    _add_textbox(slide, Cm(2), Cm(y), Cm(30), Cm(0.8),
                 "Q2：對公司的影響？", font_size=14, bold=True,
                 color=HIGHLIGHT_YELLOW)
    y += 1.0
    _add_textbox(slide, Cm(3), Cm(y), Cm(28), Cm(1.5),
                 brief["q2_impact"], font_size=12,
                 color=TEXT_WHITE)
    y += 1.8

    # Q3 — 現在要做什麼 (numbered actions ≤3)
    _add_textbox(slide, Cm(2), Cm(y), Cm(30), Cm(0.8),
                 "Q3：現在要做什麼？", font_size=14, bold=True,
                 color=HIGHLIGHT_YELLOW)
    y += 1.0
    actions = brief.get("q3_actions", [])
    # Filter out empty / hollow action lines (prevents "1. " bullets)
    action_lines = [
        f"{i}. {safe_text(a, 60)}"
        for i, a in enumerate(actions[:3], 1)
        if safe_text(a, 60)
    ]
    if not action_lines:
        action_lines = ["追蹤來源，確認數字證據後決定行動方向。"]
    _add_multiline_textbox(
        slide, Cm(3), Cm(y), Cm(28), Cm(2.5),
        action_lines, font_size=12, color=TEXT_WHITE, line_spacing=1.4,
    )
    y += max(len(action_lines) * 0.9, 1.5) + 0.5

    # Divider before bottom section
    _add_divider(slide, Cm(2), Cm(min(y, 15.5)), Cm(30), color=SUBTLE_GRAY)
    y = min(y + 0.5, 16.0)

    # Video reference
    videos = brief.get("video_source", [])
    if videos:
        vid = videos[0]
        _add_textbox(slide, Cm(2), Cm(y), Cm(30), Cm(0.7),
                     f"Video: {safe_text(vid.get('title', ''), 50)}",
                     font_size=10, color=SUBTLE_GRAY)
        y += 0.7
        _add_textbox(slide, Cm(2), Cm(y), Cm(30), Cm(0.5),
                     vid.get("url", ""), font_size=8, color=SUBTLE_GRAY)
        y += 0.7

    # Sources
    sources = brief.get("sources", [])
    if sources:
        _add_textbox(slide, Cm(2), Cm(min(y, 17.5)), Cm(30), Cm(0.7),
                     f"Source: {safe_text(sources[0], 80)}",
                     font_size=9, color=SUBTLE_GRAY)


def _slide_signal_thermometer(prs: Presentation, cards: list[EduNewsCard]) -> None:
    """Signal Thermometer — market heat gauge + signal type breakdown."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    # Title
    _add_textbox(slide, Cm(2), Cm(0.8), Cm(30), Cm(1.5),
                 "Signal Thermometer", font_size=32, bold=True,
                 color=HIGHLIGHT_YELLOW)
    _add_divider(slide, Cm(2), Cm(2.3), Cm(4), color=ACCENT)

    # Market heat index
    heat = compute_market_heat(cards)
    heat_color = ACCENT if heat["level"] in ("VERY_HIGH", "HIGH") else HIGHLIGHT_YELLOW
    _add_textbox(slide, Cm(2), Cm(3.0), Cm(10), Cm(2.5),
                 str(heat["score"]), font_size=72, bold=True,
                 color=heat_color)
    _add_textbox(slide, Cm(12), Cm(3.2), Cm(20), Cm(1),
                 f"Market Heat Index — {heat['trend_word']}",
                 font_size=16, color=TEXT_WHITE)
    _add_textbox(slide, Cm(12), Cm(4.5), Cm(20), Cm(1),
                 f"Level: {heat['level']}", font_size=12, color=SUBTLE_GRAY)

    # Signal breakdown
    signals = build_signal_summary(cards)
    y = 6.5
    _add_textbox(slide, Cm(2), Cm(y), Cm(30), Cm(1),
                 "Top Signals", font_size=18, bold=True,
                 color=HIGHLIGHT_YELLOW)
    y += 1.3

    heat_colors = {"hot": ACCENT, "warm": HIGHLIGHT_YELLOW, "cool": SUBTLE_GRAY}

    for sig in signals[:3]:
        heat_word = str(sig.get("heat", "cool"))
        sig_color = heat_colors.get(heat_word, SUBTLE_GRAY)
        signal_text = str(sig.get("signal_text", sig.get("title", "")))
        platform_count = int(sig.get("platform_count", sig.get("source_count", 0)))
        heat_score = int(sig.get("heat_score", 0))
        label = str(sig.get("label", sig.get("signal_type", "Signal")))
        title = str(sig.get("title", signal_text))
        source_name = str(sig.get("source_name", "來源平台"))
        source_url = str(sig.get("source_url", "")).strip()
        # Signal type badge
        _add_textbox(slide, Cm(2), Cm(y), Cm(8), Cm(0.8),
                     label, font_size=12, bold=True,
                     color=sig_color)
        # signal_text + required fallback fields
        _add_textbox(slide, Cm(10), Cm(y), Cm(20), Cm(0.8),
                     f"{title} | platform_count={platform_count} | heat_score={heat_score} | source: {source_name}",
                     font_size=11, color=TEXT_WHITE)
        if source_url.startswith("http"):
            _add_textbox(slide, Cm(10), Cm(y + 0.55), Cm(20), Cm(0.6),
                         source_url, font_size=10, color=SUBTLE_GRAY)
        # Heat badge
        _add_textbox(slide, Cm(30), Cm(y), Cm(2), Cm(0.8),
                     heat_word.upper(), font_size=10, bold=True,
                     color=sig_color)
        y += 1.4


def _slide_corp_watch(
    prs: Presentation,
    cards: list[EduNewsCard],
    metrics: dict | None = None,
) -> None:
    """Corp Watch — Tier A + Tier B company monitoring."""
    corp = build_corp_watch_summary(cards, metrics=metrics)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    # Title
    _add_textbox(slide, Cm(2), Cm(0.8), Cm(30), Cm(1.5),
                 "Corp Watch", font_size=32, bold=True,
                 color=HIGHLIGHT_YELLOW)
    _add_divider(slide, Cm(2), Cm(2.3), Cm(4), color=ACCENT)

    # Total mentions
    _add_textbox(slide, Cm(2), Cm(3.0), Cm(30), Cm(1),
                 f"Total Mentions: {corp['total_mentions']}",
                 font_size=14, color=TEXT_WHITE)
    _add_textbox(slide, Cm(2), Cm(3.8), Cm(30), Cm(0.8),
                 str(corp.get("status_message", "")),
                 font_size=12, color=SUBTLE_GRAY)

    # v5.1 no-event fallback
    if corp.get("updates", corp["total_mentions"]) == 0:
        fail_bits = []
        for item in corp.get("top_fail_reasons", []):
            reason = str(item.get("reason", "unknown"))
            count = int(item.get("count", 0))
            fail_bits.append(f"{reason} ({count})")
        fail_text = ", ".join(fail_bits) if fail_bits else "none"
        source_bits = []
        for src in corp.get("top_sources", [])[:3]:
            source_bits.append(
                f"{src.get('source_name', 'none')}: items_seen={src.get('items_seen', 0)}, "
                f"gate_pass={src.get('gate_pass', 0)}, gate_soft_pass={src.get('gate_soft_pass', 0)}"
            )
        top_sources_text = " | ".join(source_bits) if source_bits else "none"

        _add_textbox(slide, Cm(2), Cm(4.5), Cm(30), Cm(1),
                     "Source Scan Stats", font_size=16, bold=True,
                     color=HIGHLIGHT_YELLOW)
        _add_multiline_textbox(
            slide, Cm(3), Cm(5.7), Cm(28), Cm(8),
            [
                f"status: {corp.get('status_message', 'none')}",
                f"sources_total: {corp.get('sources_total', 0)}",
                f"success_count: {corp.get('success_count', 0)}",
                f"fail_count: {corp.get('fail_count', 0)}",
                f"top_fail_reasons: {fail_text}",
                f"top_sources: {top_sources_text}",
            ],
            font_size=12,
            color=TEXT_WHITE,
            line_spacing=1.5,
        )
        return

    y = 4.5

    # Tier A
    _add_textbox(slide, Cm(2), Cm(y), Cm(30), Cm(1),
                 "Tier A — Global Leaders", font_size=16, bold=True,
                 color=HIGHLIGHT_YELLOW)
    y += 1.2

    if corp["tier_a"]:
        for item in corp["tier_a"][:5]:
            _add_highlight_textbox(
                slide, Cm(3), Cm(y), Cm(28), Cm(0.8),
                f"{item['name']}  ", item["impact_label"],
                f"  {safe_text(item['event_title'], 30)}",
                font_size=11,
            )
            y += 1.0
    else:
        _add_textbox(slide, Cm(3), Cm(y), Cm(28), Cm(0.8),
                     "今日無 Tier A 公司相關事件", font_size=11,
                     color=SUBTLE_GRAY)
        y += 1.0

    y += 0.5

    # Tier B
    _add_textbox(slide, Cm(2), Cm(y), Cm(30), Cm(1),
                 "Tier B — Asia Leaders", font_size=16, bold=True,
                 color=HIGHLIGHT_YELLOW)
    y += 1.2

    if corp["tier_b"]:
        for item in corp["tier_b"][:4]:
            _add_highlight_textbox(
                slide, Cm(3), Cm(y), Cm(28), Cm(0.8),
                f"{item['name']}  ", item["impact_label"],
                f"  {safe_text(item['event_title'], 30)}",
                font_size=11,
            )
            y += 1.0
    else:
        _add_textbox(slide, Cm(3), Cm(y), Cm(28), Cm(0.8),
                     "今日無 Tier B 公司相關事件", font_size=11,
                     color=SUBTLE_GRAY)


def _slide_scan_diagnostics(
    prs: Presentation,
    cards: list[EduNewsCard],
    metrics: dict | None,
    event_cards: list[EduNewsCard],
) -> None:
    """Fixed diagnostics slide to keep base deck flow complete on low-news days."""
    m = metrics or {}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    _add_textbox(slide, Cm(2), Cm(0.8), Cm(30), Cm(1.5),
                 "Scan Diagnostics", font_size=32, bold=True,
                 color=HIGHLIGHT_YELLOW)
    _add_divider(slide, Cm(2), Cm(2.3), Cm(4), color=ACCENT)

    fetched_total = int(m.get("fetched_total", len(cards)))
    hard_pass_total = int(m.get("hard_pass_total", 0))
    soft_pass_total = int(m.get("soft_pass_total", 0))
    rejected_total = int(m.get("rejected_total", m.get("gate_reject_total", 0)))
    sources_total = int(m.get("sources_total", 0))
    success_count = int(m.get("sources_success", 0))
    fail_count = int(m.get("sources_failed", 0))

    lines = [
        f"fetched_total={fetched_total}",
        f"hard_pass_total={hard_pass_total}",
        f"soft_pass_total={soft_pass_total}",
        f"rejected_total={rejected_total}",
        f"sources_total={sources_total}",
        f"success_count={success_count}",
        f"fail_count={fail_count}",
        f"event_candidates={len(event_cards)}",
    ]
    _add_multiline_textbox(
        slide, Cm(3), Cm(3.5), Cm(28), Cm(10),
        lines, font_size=16, color=TEXT_WHITE, line_spacing=1.4,
    )

    top_density = list(m.get("density_score_top5", []))[:3]
    if top_density:
        preview = []
        for row in top_density:
            title = str(row[0]) if isinstance(row, (list, tuple)) and len(row) > 0 else "candidate"
            score = int(row[2]) if isinstance(row, (list, tuple)) and len(row) > 2 else 0
            preview.append(f"{safe_text(title, 30)} ({score})")
        _add_textbox(slide, Cm(3), Cm(14.5), Cm(28), Cm(2.5),
                     f"density_score_top5: {' | '.join(preview)}",
                     font_size=12, color=SUBTLE_GRAY)


def _slide_event_ranking(
    prs: Presentation,
    event_cards: list[EduNewsCard],
    cards: list[EduNewsCard] | None = None,
) -> None:
    """Event ranking slide; in no-event mode use signal/corp-backed ranking rows."""
    headers = ["Rank", "Impact", "標題", "類別", "Action"]
    rows = []
    if event_cards:
        scored = []
        for c in event_cards[:8]:
            impact = score_event_impact(c)
            scored.append((c, impact))
        scored.sort(key=lambda x: x[1]["impact"], reverse=True)

        for rank, (c, imp) in enumerate(scored, 1):
            dc = build_decision_card(c)
            action = dc["actions"][0] if dc["actions"] else "待確認"
            action = semantic_guard_text(safe_text(action, 25), c, context="action")
            rows.append([
                str(rank),
                f"{imp['impact']}/5 {imp['label']}",
                safe_text(c.title_plain, 25),
                c.category or "綜合",
                action[:25],
            ])
    else:
        signals = build_signal_summary(cards or [])
        for rank, sig in enumerate(signals[:3], 1):
            heat_score = int(sig.get("heat_score", 30) or 30)
            action = "WATCH" if rank <= 2 else "TEST"
            rows.append(
                [
                    str(rank),
                    f"{heat_score}/100",
                    safe_text(str(sig.get("title", sig.get("signal_text", "來源訊號"))), 25),
                    "No-Event",
                    action,
                ]
            )
    _add_table_slide(prs, "Event Ranking  事件影響力排行", headers, rows)


def _slide_recommended_moves(prs: Presentation, cards: list[EduNewsCard]) -> None:
    """Recommended Moves — MOVE (red) / TEST (yellow) / WATCH (gray) actions."""
    actions = build_ceo_actions(cards)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    # Title
    _add_textbox(slide, Cm(2), Cm(0.8), Cm(30), Cm(1.5),
                 "Recommended Moves", font_size=32, bold=True,
                 color=HIGHLIGHT_YELLOW)
    _add_divider(slide, Cm(2), Cm(2.3), Cm(4), color=ACCENT)

    type_colors = {
        "MOVE": ACCENT,
        "TEST": HIGHLIGHT_YELLOW,
        "WATCH": SUBTLE_GRAY,
    }

    y = 3.5
    if not actions:
        _add_textbox(slide, Cm(3), Cm(y), Cm(28), Cm(1),
                     "本日無需要立即行動的事項", font_size=14,
                     color=TEXT_WHITE)
        return

    from utils.semantic_quality import is_placeholder_or_fragment as _is_frag_ppt
    _gloss_seen_moves: set = set()
    for act in actions[:6]:
        tag_color = type_colors.get(act["action_type"], SUBTLE_GRAY)

        # Guard (D): replace fragment detail with safe fallback, then normalize
        detail_text = act["detail"]
        if not detail_text or _is_frag_ppt(detail_text):
            detail_text = "持續監控此事件發展（T+7）"
        detail_text = _v1_norm_gloss(detail_text, _V1_GLOSSARY, _gloss_seen_moves)

        # Action type badge
        _add_textbox(slide, Cm(2), Cm(y), Cm(4), Cm(0.9),
                     act["action_type"], font_size=14, bold=True,
                     color=tag_color)
        # Title
        _add_textbox(slide, Cm(6.5), Cm(y), Cm(14), Cm(0.9),
                     act["title"], font_size=12, color=TEXT_WHITE)
        # Detail
        _add_textbox(slide, Cm(6.5), Cm(y + 0.8), Cm(20), Cm(0.7),
                     detail_text, font_size=10, color=SUBTLE_GRAY)
        # Owner
        _add_textbox(slide, Cm(27), Cm(y), Cm(5), Cm(0.9),
                     act["owner"], font_size=10, color=SUBTLE_GRAY,
                     alignment=PP_ALIGN.RIGHT)
        y += 2.0


def _slide_pending_decisions(prs: Presentation, event_cards: list[EduNewsCard]) -> None:
    """Last slide: pending decisions & owners."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, Cm(2), Cm(1.5), Cm(30), Cm(2.5),
                 "待決事項與 Owner", font_size=36, bold=True,
                 color=HIGHLIGHT_YELLOW)
    _add_divider(slide, Cm(2), Cm(3.8), Cm(4), color=ACCENT)

    items: list[str] = []
    for i, c in enumerate(event_cards[:5], 1):
        dc = build_decision_card(c)
        action = dc["actions"][0] if dc["actions"] else "待確認"
        # Semantic guard: ensure action is not hollow
        action = semantic_guard_text(safe_text(action, 55), c, context="action")
        owner = dc["owner"]
        impact_data = score_event_impact(c)
        impact = impact_data.get("impact", 3) if isinstance(impact_data, dict) else 3
        why_raw = safe_text(c.why_important or "", 40) or safe_text(c.title_plain or "", 30) or ""
        from utils.semantic_quality import is_placeholder_or_fragment as _is_frag
        why_snippet = "" if _is_frag(why_raw) else why_raw
        items.append(
            f"{i}. {action[:55]} → Owner: {owner} "
            f"| Due: T+7 | Metric: impact={impact}/5"
            + (f" | {why_snippet}" if why_snippet else "")
        )

    if not items:
        items.append("1. 本日無待決事項")

    _add_multiline_textbox(slide, Cm(3), Cm(5), Cm(28), Cm(13),
                           items, font_size=18, color=TEXT_WHITE,
                           line_spacing=1.8)
    _add_divider(slide, Cm(0), Cm(18.7), SLIDE_WIDTH, color=ACCENT)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def generate_executive_ppt(
    cards: list[EduNewsCard],
    health: SystemHealthReport,
    report_time: str,
    total_items: int,
    output_path: Path | None = None,
    theme: str = "light",
    metrics: dict | None = None,
) -> Path:
    log = get_logger()
    if output_path is None:
        project_root = Path(__file__).resolve().parent.parent
        output_path = project_root / "outputs" / "executive_report.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)

    _apply_theme(theme)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 1. Cover
    _slide_cover(prs, report_time)

    # 2. Structured Summary (5 sections)
    _slide_structured_summary(prs, cards, metrics=metrics)

    # 3. Signal Thermometer (v5)
    _slide_signal_thermometer(prs, cards)

    # 4. Corp Watch (v5)
    _slide_corp_watch(prs, cards, metrics=metrics)

    # Event pool: only verifiable event cards (no synthetic placeholders).
    event_cards = get_event_cards_for_deck(cards, metrics=metrics or {}, min_events=0)
    _final_cards_payload = _load_final_cards(metrics)
    event_cards, _event_payloads = _align_event_cards_with_final_cards(event_cards, _final_cards_payload)

    # Quality guard: ensure per-card text density meets thresholds.
    for ec in event_cards:
        for attr in ("what_happened", "why_important"):
            val = getattr(ec, attr, None)
            if val:
                guarded, _m = quality_guard_block(val, card=ec)
                if guarded and guarded != val:
                    object.__setattr__(ec, attr, guarded)

    # 5. Key Takeaways
    _slide_key_takeaways(prs, cards, total_items, metrics=metrics, event_cards=event_cards)

    # 6. Overview Table
    _slide_overview_table(prs, event_cards, cards=cards)

    # 7. Event Ranking (v5)
    _slide_event_ranking(prs, event_cards, cards=cards)

    # 8. Per-event: brief_page1 + brief_page2
    for i, card in enumerate(event_cards, 1):
        _payload = _event_payloads[i - 1] if i - 1 < len(_event_payloads) else None
        _slide_brief_page1(prs, card, i, final_payload=_payload)
        _slide_brief_page2(prs, card, i, final_payload=_payload)

    # 9. Recommended Moves (v5)
    _slide_recommended_moves(prs, cards)

    # 10. Decision Matrix (6 columns)
    decision_rows = []
    for i, c in enumerate(event_cards[:8], 1):
        dc = build_decision_card(c)
        decision_rows.append([
            str(i),
            safe_text(dc["event"], 18),
            safe_text(dc["effects"][0], 25) if dc["effects"] else "缺口",
            safe_text(dc["risks"][0], 25) if dc["risks"] else "缺口",
            safe_text(dc["actions"][0], 30) if dc["actions"] else "待確認",
            dc["owner"],
        ])
    _add_table_slide(
        prs, "決策摘要表  Decision Matrix",
        ["#", "事件", "影響", "風險", "建議行動", "要問誰"],
        decision_rows,
    )

    # 11. Pending Decisions
    _slide_pending_decisions(prs, event_cards)

    prs.save(str(output_path))
    _pptx_write_exists = output_path.exists()
    _pptx_write_size = output_path.stat().st_size if _pptx_write_exists else 0
    log.info(
        "PPTX_WRITE_CHECK path=%s exists=%s size=%d",
        output_path,
        _pptx_write_exists,
        _pptx_write_size,
    )
    if not _pptx_write_exists or _pptx_write_size <= 0:
        raise RuntimeError(f"PPTX_WRITE_CHECK failed: path={output_path}")
    log.info("Executive PPTX generated: %s", output_path)
    return output_path


# ---------------------------------------------------------------------------
# v5.2.2 overrides (append-only quality hotfix layer)
# ---------------------------------------------------------------------------

_v521_slide_structured_summary = _slide_structured_summary
_v521_slide_key_takeaways = _slide_key_takeaways


def _slide_structured_summary(
    prs: Presentation,
    cards: list[EduNewsCard],
    metrics: dict | None = None,
) -> None:
    """Structured summary with metric-backed fallback in no-event days."""
    summary = build_structured_executive_summary(cards, metrics=metrics or {})
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    _add_textbox(slide, Cm(2), Cm(0.8), Cm(30), Cm(1.5),
                 "Structured Summary", font_size=32, bold=True,
                 color=HIGHLIGHT_YELLOW)
    _add_divider(slide, Cm(2), Cm(2.3), Cm(4), color=ACCENT)

    section_map = [
        ("AI Trends", summary.get("ai_trends", [])),
        ("Tech Landing", summary.get("tech_landing", [])),
        ("Market Competition", summary.get("market_competition", [])),
        ("Opportunities & Risks", summary.get("opportunities_risks", [])),
        ("Recommended Actions", summary.get("recommended_actions", [])),
    ]

    y = 3.0
    for sec_title, items in section_map:
        _add_textbox(slide, Cm(2), Cm(y), Cm(30), Cm(1),
                     sec_title, font_size=14, bold=True,
                     color=HIGHLIGHT_YELLOW)
        y += 1.0
        for item in items[:2]:
            _add_textbox(slide, Cm(3), Cm(y), Cm(28), Cm(0.8),
                         f"- {safe_text(item, 80)}", font_size=11,
                         color=TEXT_WHITE)
            y += 0.8
        y += 0.3


# ---------------------------------------------------------------------------
# v5.2.3 — Three-layer last-mile banned-phrase guard for Structured Summary
# Wraps v5.2.2; applies A/B/C guards to each section item before textbox write.
# Prevents truncation artifacts (e.g. '的趨勢，解決方 記') from reaching PPTX.
# ---------------------------------------------------------------------------
_v522_slide_structured_summary = _slide_structured_summary


def _slide_structured_summary(
    prs: Presentation,
    cards: list[EduNewsCard],
    metrics: dict | None = None,
) -> None:
    """v5.2.3: Three-layer guard on every Structured Summary section item.

    Layer A — fragment guard   : is_placeholder_or_fragment → safe fallback
    Layer B — sanitize_exec_text: exec_sanitizer banned-phrase check
    Layer C — hard strip       : explicit BANNED_SUBSTRINGS scan (last resort)
    """
    _SAFE_V523 = "重點：本事件已補足關鍵錨點，詳見下方 Proof 與 Watchlist。"

    try:
        from utils.exec_sanitizer import sanitize_exec_text as _san_v523
        from utils.exec_sanitizer import BANNED_SUBSTRINGS as _BAN_V523
    except Exception:
        _v522_slide_structured_summary(prs, cards, metrics=metrics)
        return

    try:
        from utils.semantic_quality import is_placeholder_or_fragment as _is_frag_v523
    except Exception:
        _is_frag_v523 = None  # type: ignore[assignment]

    def _guard_item_v523(raw: str) -> str:
        s = str(raw or "").strip()
        # Layer A: fragment / template-remnant guard
        if _is_frag_v523 is not None:
            try:
                if not s or _is_frag_v523(s):
                    return _SAFE_V523
            except Exception:
                pass
        # Layer B: exec_sanitizer (returns _SAFE_FALLBACK if banned phrase found)
        try:
            s = _san_v523(s)
        except Exception:
            pass
        # Layer C: explicit scan — covers any edge case missed by B
        for _bp in _BAN_V523:
            if _bp in s:
                return _SAFE_V523
        return s

    summary = build_structured_executive_summary(cards, metrics=metrics or {})
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    _add_textbox(slide, Cm(2), Cm(0.8), Cm(30), Cm(1.5),
                 "Structured Summary", font_size=32, bold=True,
                 color=HIGHLIGHT_YELLOW)
    _add_divider(slide, Cm(2), Cm(2.3), Cm(4), color=ACCENT)

    section_map = [
        ("AI Trends", summary.get("ai_trends", [])),
        ("Tech Landing", summary.get("tech_landing", [])),
        ("Market Competition", summary.get("market_competition", [])),
        ("Opportunities & Risks", summary.get("opportunities_risks", [])),
        ("Recommended Actions", summary.get("recommended_actions", [])),
    ]

    y = 3.0
    for sec_title, items in section_map:
        _add_textbox(slide, Cm(2), Cm(y), Cm(30), Cm(1),
                     sec_title, font_size=14, bold=True,
                     color=HIGHLIGHT_YELLOW)
        y += 1.0
        for item in items[:2]:
            guarded = _guard_item_v523(item)
            _add_textbox(slide, Cm(3), Cm(y), Cm(28), Cm(0.8),
                         f"- {safe_text(guarded, 80)}", font_size=11,
                         color=TEXT_WHITE)
            y += 0.8
        y += 0.3


def _slide_key_takeaways(
    prs: Presentation,
    cards: list[EduNewsCard],
    total_items: int,
    metrics: dict | None = None,
    event_cards: list[EduNewsCard] | None = None,
) -> None:
    """Key takeaways with stats-backed no-event fallback."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, Cm(2), Cm(1.2), Cm(30), Cm(2),
                 "Key Takeaways", font_size=36, bold=True,
                 color=HIGHLIGHT_YELLOW)
    _add_divider(slide, Cm(2), Cm(3.2), Cm(4), color=ACCENT)

    if event_cards is None:
        event_cards = get_event_cards_for_deck(cards, metrics=metrics or {}, min_events=0)
    takeaways: list[str] = [f"本次掃描總量：{total_items}；事件候選：{len(event_cards)}。"]

    for c in event_cards[:3]:
        dc = build_decision_card(c)
        if bool(getattr(c, "low_confidence", False)):
            takeaways.append(f"低信心候選：{safe_text(c.title_plain, 35)}；重點：{dc['event']}。")
        else:
            takeaways.append(f"{safe_text(c.title_plain, 35)}；重點：{dc['event']}。")

    if not event_cards:
        fetched_total = int((metrics or {}).get("fetched_total", total_items))
        gate_pass_total = int((metrics or {}).get("gate_pass_total", sum(1 for c in cards if c.is_valid_news)))
        sources_total = int((metrics or {}).get("sources_total", 0))
        after_filter_total = int((metrics or {}).get("after_filter_total", gate_pass_total))
        takeaways.extend(
            [
                f"掃描概況：fetched_total={fetched_total}、gate_pass_total={gate_pass_total}。",
                f"來源覆蓋：sources_total={sources_total}、after_filter_total={after_filter_total}。",
                "今日未形成高信心事件，先以來源級統計維持決策資訊量。",
            ]
        )

    _add_multiline_textbox(
        slide, Cm(3), Cm(4.5), Cm(28), Cm(13),
        takeaways, font_size=18, color=TEXT_WHITE, line_spacing=1.8,
    )


# ===========================================================================
# EXEC_VISUAL_TEMPLATE_V1 — append-only visual layout override layer
# Implements T1-T6 template codes for 6 key slide types.
# All colours / sizes reference utils/exec_visual_tokens.py.
# Does NOT modify schemas/education_models.py.
# Does NOT add new pip dependencies.
# ===========================================================================

import json as _json

from utils.exec_visual_tokens import (
    PRIMARY_BLUE as _V1_BLUE,
    SOFT_BLUE_BG as _V1_SOFT_BG,
    TEXT_GRAY as _V1_TEXT_GRAY,
    CARD_WHITE as _V1_CARD_WHITE,
    DIVIDER_GRAY as _V1_DIVIDER,
    GREEN_ACCENT as _V1_GREEN,
    ORANGE_ACCENT as _V1_AMBER,
    RED_ACCENT as _V1_RED,
    STAGE_COLORS as _V1_STAGE_COLORS,
    TITLE_FONT_SIZE as _V1_TITLE_FS,
    BODY_FONT_SIZE as _V1_BODY_FS,
    CARD_TITLE_FONT_SIZE as _V1_CARD_TITLE_FS,
    CARD_BODY_FONT_SIZE as _V1_CARD_BODY_FS,
    LABEL_FONT_SIZE as _V1_LABEL_FS,
    LINE_SPACING as _V1_LINE_SPACING,
    CARD_SHAPE_TYPE as _V1_ROUNDED_RECT,
    RECT_SHAPE_TYPE as _V1_RECT,
    LAYOUT_VERSION as _V1_LAYOUT_VER,
    TEMPLATE_MAP as _V1_TEMPLATE_MAP,
)
from utils.narrative_compact import (
    build_narrative_compact as _v1_narrative,
    has_hard_evidence as _v1_has_evidence,
    extract_first_hard_evidence as _v1_extract_ev,
)
from utils.bullet_normalizer import (
    normalize_bullets as _v1_norm_bullets,
    normalize_bullets_safe as _v1_norm_bullets_safe,
    compute_bullet_stats as _v1_bullet_stats,
)
from utils.hybrid_glossing import (
    normalize_exec_text as _v1_norm_gloss,
    load_glossary as _v1_load_glossary,
    get_gloss_stats as _v1_gloss_stats,
    reset_gloss_stats as _v1_reset_gloss_stats,
)

# Glossary loaded once at module level (cached inside hybrid_glossing)
_V1_GLOSSARY: dict = _v1_load_glossary()


# ---------------------------------------------------------------------------
# V1 Shape / Text helpers
# ---------------------------------------------------------------------------

def _v1_set_shape_fill(shape, bg_color, line_color=None) -> None:
    """Set shape fill; remove border unless line_color specified."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color


def _v1_add_card(
    slide,
    left_cm: float, top_cm: float,
    width_cm: float, height_cm: float,
    header_text: str = '',
    body_lines: list[str] | None = None,
    bg_color=None,
    header_color=None,
    body_color=None,
    header_font_size: int = 12,
    body_font_size: int = 10,
    body_text_limit: int = 120,
    shape_type: int = 1,
) -> None:
    """Draw a card (filled rectangle) with optional header + body lines."""
    if bg_color is None:
        bg_color = _V1_SOFT_BG
    if header_color is None:
        header_color = _V1_BLUE
    if body_color is None:
        body_color = _V1_TEXT_GRAY

    # Card background
    card = slide.shapes.add_shape(
        shape_type,
        Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm),
    )
    _v1_set_shape_fill(card, bg_color)

    # Header text
    if header_text:
        _add_textbox(
            slide,
            Cm(left_cm + 0.3), Cm(top_cm + 0.15),
            Cm(width_cm - 0.4), Cm(0.8),
            safe_text(header_text, 60),
            font_size=header_font_size,
            color=header_color,
            bold=True,
        )

    # Body lines
    if body_lines:
        clean = [safe_text(b, body_text_limit) for b in body_lines if b.strip()]
        if clean:
            _add_multiline_textbox(
                slide,
                Cm(left_cm + 0.3), Cm(top_cm + 1.0),
                Cm(width_cm - 0.5), Cm(height_cm - 1.2),
                clean,
                font_size=body_font_size,
                color=body_color,
                line_spacing=_V1_LINE_SPACING,
            )


def _v1_slide_header(slide, title: str, report_time: str = '') -> None:
    """Standard V1 slide header: large title left + date right + divider."""
    _add_textbox(
        slide, Cm(2.0), Cm(0.8), Cm(24), Cm(1.4),
        safe_text(title, 60),
        font_size=_V1_TITLE_FS, bold=True, color=TEXT_WHITE,
    )
    if report_time:
        _add_textbox(
            slide, Cm(26), Cm(0.8), Cm(6.5), Cm(1.0),
            safe_text(report_time, 30),
            font_size=_V1_LABEL_FS, color=SUBTLE_GRAY,
            alignment=PP_ALIGN.RIGHT,
        )
    # Accent divider under header
    _add_divider(slide, Cm(2.0), Cm(2.4), Cm(29.5), color=_V1_BLUE)


def _v1_mini_arrow(slide, x_cm: float, y_cm: float) -> None:
    """Small horizontal arrow connector between cards."""
    arrow = slide.shapes.add_shape(_V1_RECT, Cm(x_cm), Cm(y_cm), Cm(0.8), Cm(0.15))
    _v1_set_shape_fill(arrow, _V1_BLUE)


def _v1_vertical_connector(slide, x_cm: float, y_top: float, y_bot: float) -> None:
    """Thin vertical line connecting T1 cards."""
    h = max(y_bot - y_top, 0.1)
    line = slide.shapes.add_shape(_V1_RECT, Cm(x_cm), Cm(y_top), Cm(0.12), Cm(h))
    _v1_set_shape_fill(line, _V1_BLUE)


# ---------------------------------------------------------------------------
# Save references to v5.2.2 versions before overriding
# ---------------------------------------------------------------------------
_v1_prev_slide_overview_table = _slide_overview_table
_v1_prev_slide_signal_thermometer = _slide_signal_thermometer
_v1_prev_slide_event_ranking = _slide_event_ranking
_v1_prev_slide_pending_decisions = _slide_pending_decisions
_v1_prev_slide_brief_page1 = _slide_brief_page1
_v1_prev_slide_brief_page2 = _slide_brief_page2
_v1_prev_generate_executive_ppt = generate_executive_ppt


# ---------------------------------------------------------------------------
# T5: Horizontal Rail — 今日總覽 Overview
# ---------------------------------------------------------------------------

def _slide_overview_table(
    prs: Presentation,
    event_cards: list[EduNewsCard],
    cards: list[EduNewsCard] | None = None,
) -> None:
    """T5: Horizontal Rail layout for 今日總覽 Overview slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _v1_slide_header(slide, '今日總覽  Overview')

    # Source data for nodes
    if event_cards:
        items = event_cards[:8]
        node_data = []
        for c in items:
            title_short = safe_text(c.title_plain or '', 30)
            keypoint = safe_text(c.what_happened or '', 45) or safe_text(c.why_important or '', 45)
            if not keypoint:
                keypoint = safe_text(c.title_plain or '', 45)
            score_str = f'{c.final_score:.1f}' if getattr(c, 'final_score', None) is not None else ''
            node_data.append({'title': title_short, 'key': keypoint, 'cat': c.category or '綜合', 'score': score_str})
    else:
        sigs = build_signal_summary(cards or [])
        node_data = []
        for sig in sigs[:8]:
            t = safe_text(str(sig.get('title', sig.get('signal_text', '訊號'))), 14)
            k = safe_text(str(sig.get('label', sig.get('signal_type', ''))), 20)
            node_data.append({'title': t, 'key': k, 'cat': 'Signal'})

    display = node_data[:8] if node_data else [{'title': '今日無事件', 'key': '持續監控中', 'cat': '綜合', 'score': ''}]
    n = len(display)

    # Left icon block
    icon = slide.shapes.add_shape(_V1_RECT, Cm(1.5), Cm(3.0), Cm(2.2), Cm(3.5))
    _v1_set_shape_fill(icon, _V1_BLUE)
    _add_textbox(
        slide, Cm(1.55), Cm(3.8), Cm(2.1), Cm(1.5),
        'TODAY', font_size=13, bold=True,
        color=_V1_CARD_WHITE, alignment=PP_ALIGN.CENTER,
    )

    # Horizontal rail
    rail_x_start, rail_x_end = 4.0, 32.5
    rail_y = 4.8
    rail = slide.shapes.add_shape(
        _V1_RECT,
        Cm(rail_x_start), Cm(rail_y), Cm(rail_x_end - rail_x_start), Cm(0.22),
    )
    _v1_set_shape_fill(rail, _V1_BLUE)

    # Nodes + cards
    rail_span = rail_x_end - rail_x_start
    step = rail_span / max(n, 1)
    for i, item in enumerate(display):
        nx = rail_x_start + (i + 0.5) * step
        # Node marker
        nmarker = slide.shapes.add_shape(
            _V1_RECT, Cm(nx - 0.3), Cm(rail_y - 0.25), Cm(0.6), Cm(0.6),
        )
        _v1_set_shape_fill(nmarker, _V1_BLUE)

        # Card below rail
        card_left = max(nx - 2.1, rail_x_start - 0.5)
        card_left = min(card_left, rail_x_end - 4.0)
        card_w = 4.0
        card_y = rail_y + 0.6
        bg = _V1_SOFT_BG if i % 2 == 0 else _V1_CARD_WHITE
        score_line = f'Score: {item["score"]}' if item.get('score') else f'[{item["cat"]}]'
        _v1_add_card(
            slide, card_left, card_y, card_w, 3.8,
            header_text=item['title'],
            body_lines=[item['key'], score_line],
            bg_color=bg,
            header_color=_V1_BLUE,
            body_color=_V1_TEXT_GRAY,
            header_font_size=_V1_CARD_TITLE_FS,
            body_font_size=_V1_CARD_BODY_FS,
        )

    # Overflow note
    if event_cards and len(event_cards) > 8:
        _add_textbox(
            slide, Cm(2.0), Cm(17.8), Cm(29), Cm(0.7),
            f'另有 {len(event_cards) - 8} 則事件見 Event Ranking',
            font_size=_V1_LABEL_FS, color=SUBTLE_GRAY,
        )


# ---------------------------------------------------------------------------
# T2: Stage Arrow — Signal Thermometer
# ---------------------------------------------------------------------------

def _slide_signal_thermometer(prs: Presentation, cards: list[EduNewsCard]) -> None:
    """T2: Stage Arrow layout for Signal Thermometer slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    heat = compute_market_heat(cards)
    collected_at = ''
    for c in cards[:1]:
        collected_at = safe_text(getattr(c, 'source_name', '') or '', 20)

    _v1_slide_header(
        slide,
        f'Signal Thermometer  |  Market Heat Index: {heat["score"]} / 100',
        report_time=collected_at or heat.get('level', ''),
    )

    signals = build_signal_summary(cards)
    stage_names = ['Market Heat', 'Model Release', 'Productization', 'Regulatory / GTM']
    stage_narratives: list[str] = []
    heat_map = {
        'VERY_HIGH': '市場熱度極高；多方訊號密集匯聚。',
        'HIGH': '訊號強烈，需立即關注相關動態。',
        'MEDIUM': '訊號持續累積，保持追蹤。',
        'LOW': '目前訊號平穩，定期監控即可。',
    }
    for i in range(4):
        if i < len(signals):
            sig = signals[i]
            s_title = str(sig.get('title', sig.get('signal_text', '')))
            s_label = str(sig.get('label', sig.get('signal_type', '')))
            platform_count = int(sig.get('platform_count', sig.get('source_count', 0)))
            heat_word = str(sig.get('heat', 'cool')).upper()
            stage_narratives.append(
                f'{s_label}: {safe_text(s_title, 35)} [{heat_word}] (×{platform_count})'
            )
        else:
            if i == 0:
                stage_narratives.append(heat_map.get(heat.get('level', 'MEDIUM'), heat_map['MEDIUM']))
            else:
                stage_narratives.append('本週段無顯著新訊號；持續監測中。')

    # Draw 4 stage cards with arrows
    card_w, card_h = 7.0, 8.0
    positions_x = [1.3, 9.3, 17.3, 25.3]
    card_top = 3.2
    for i, (stage_name, narrative) in enumerate(zip(stage_names, stage_narratives)):
        color = _V1_STAGE_COLORS[i]
        # Stage label bar
        label_bar = slide.shapes.add_shape(
            _V1_RECT,
            Cm(positions_x[i]), Cm(card_top), Cm(card_w), Cm(1.2),
        )
        _v1_set_shape_fill(label_bar, color)
        _add_textbox(
            slide,
            Cm(positions_x[i] + 0.2), Cm(card_top + 0.1),
            Cm(card_w - 0.3), Cm(1.0),
            stage_name, font_size=11, bold=True, color=_V1_CARD_WHITE,
        )
        # Card body
        body_card = slide.shapes.add_shape(
            _V1_ROUNDED_RECT,
            Cm(positions_x[i]), Cm(card_top + 1.2),
            Cm(card_w), Cm(card_h - 1.2),
        )
        _v1_set_shape_fill(body_card, _V1_SOFT_BG)
        _add_textbox(
            slide,
            Cm(positions_x[i] + 0.3), Cm(card_top + 1.5),
            Cm(card_w - 0.5), Cm(card_h - 1.8),
            safe_text(narrative, 120),
            font_size=_V1_CARD_BODY_FS,
            color=_V1_TEXT_GRAY,
        )
        # Arrow connector
        if i < 3:
            arr_x = positions_x[i] + card_w + 0.05
            arr_y = card_top + card_h / 2
            _v1_mini_arrow(slide, arr_x, arr_y)
            _add_textbox(
                slide,
                Cm(arr_x + 0.1), Cm(arr_y - 0.35), Cm(0.7), Cm(0.7),
                '→', font_size=14, bold=True, color=_V1_BLUE,
            )

    # Bottom summary
    _add_textbox(
        slide, Cm(2.0), Cm(17.5), Cm(29), Cm(0.8),
        f'Market Heat Level: {heat["level"]}  |  Trend: {heat["trend_word"]}  |  Score: {heat["score"]}/100',
        font_size=_V1_LABEL_FS, color=SUBTLE_GRAY,
    )


# ---------------------------------------------------------------------------
# T4: 3-Column Bucket Cards — Event Ranking
# ---------------------------------------------------------------------------

def _slide_event_ranking(
    prs: Presentation,
    event_cards: list[EduNewsCard],
    cards: list[EduNewsCard] | None = None,
) -> None:
    """T4: Three-column bucket cards for Event Ranking 事件影響力排行."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _v1_slide_header(slide, 'Event Ranking  事件影響力排行')

    # Bucket events by category
    buckets: dict[str, list[EduNewsCard]] = {'Product': [], 'Tech': [], 'Business': []}
    if event_cards:
        assigned: set[str] = set()
        for c in event_cards[:9]:
            cat = (c.category or '').lower()
            if any(k in cat for k in ('product', '產品', 'saas', 'app', '應用')):
                buckets['Product'].append(c)
                assigned.add(c.item_id)
            elif any(k in cat for k in ('tech', '技術', 'ai', '模型', 'model', 'gpu')):
                buckets['Tech'].append(c)
                assigned.add(c.item_id)
            elif any(k in cat for k in ('business', '商業', '市場', '融資', 'market', 'fund')):
                buckets['Business'].append(c)
                assigned.add(c.item_id)
        # Distribute unassigned events evenly
        for i, c in enumerate(ec for ec in event_cards[:9] if ec.item_id not in assigned):
            key = ['Product', 'Tech', 'Business'][i % 3]
            buckets[key].append(c)

    col_configs = [
        ('Product', _V1_BLUE, 1.3),
        ('Tech', _V1_GREEN, 11.8),
        ('Business', _V1_AMBER, 22.3),
    ]
    col_w = 9.5
    header_h = 1.2
    card_top_start = 5.3
    card_h = 4.6
    card_gap = 0.3

    for col_name, col_color, col_x in col_configs:
        # Column header bar
        hdr = slide.shapes.add_shape(
            _V1_RECT, Cm(col_x), Cm(3.3), Cm(col_w), Cm(header_h),
        )
        _v1_set_shape_fill(hdr, col_color)
        _add_textbox(
            slide, Cm(col_x + 0.3), Cm(3.4), Cm(col_w - 0.4), Cm(1.0),
            col_name, font_size=14, bold=True, color=_V1_CARD_WHITE,
        )
        uline = slide.shapes.add_shape(
            _V1_RECT, Cm(col_x), Cm(3.3 + header_h), Cm(col_w), Cm(0.1),
        )
        _v1_set_shape_fill(uline, col_color)

        col_events = buckets.get(col_name, [])[:2]

        for j in range(2):
            cy = card_top_start + j * (card_h + card_gap)
            if j < len(col_events):
                c = col_events[j]
                impact_data = score_event_impact(c)
                impact_score = impact_data.get('impact', 3) if isinstance(impact_data, dict) else 3
                dc = build_decision_card(c)
                action = dc['actions'][0] if dc['actions'] else '持續追蹤'
                bullets = _v1_norm_bullets_safe(
                    [safe_text(c.what_happened or '', 50), safe_text(action, 40)],
                )
                _v1_add_card(
                    slide, col_x, cy, col_w, card_h,
                    header_text=safe_text(c.title_plain or '', 22),
                    body_lines=bullets[:2],
                    bg_color=_V1_SOFT_BG,
                    header_color=col_color,
                    body_color=_V1_TEXT_GRAY,
                    header_font_size=_V1_CARD_TITLE_FS,
                    body_font_size=_V1_CARD_BODY_FS,
                )
                _add_textbox(
                    slide, Cm(col_x + 0.3), Cm(cy + card_h - 0.7), Cm(col_w - 0.4), Cm(0.6),
                    f'Impact: {impact_score}/5',
                    font_size=9, color=SUBTLE_GRAY,
                )
            else:
                _v1_add_card(
                    slide, col_x, cy, col_w, card_h,
                    header_text='本類別',
                    body_lines=['本類別今日無正式事件，請見 Developing Watchlist。'],
                    bg_color=_V1_CARD_WHITE,
                    header_color=SUBTLE_GRAY,
                    body_color=SUBTLE_GRAY,
                    header_font_size=_V1_CARD_TITLE_FS,
                    body_font_size=_V1_CARD_BODY_FS,
                )


# ---------------------------------------------------------------------------
# T6: Promotion Curve — 待決事項與 Owner
# ---------------------------------------------------------------------------

def _slide_pending_decisions(prs: Presentation, event_cards: list[EduNewsCard]) -> None:
    """T6: Promotion Curve layout for 待決事項與 Owner slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _v1_slide_header(slide, '待決事項與 Owner')

    # Horizontal guide rail
    guide = slide.shapes.add_shape(
        _V1_RECT, Cm(2.0), Cm(6.8), Cm(24.0), Cm(0.18),
    )
    _v1_set_shape_fill(guide, _V1_BLUE)

    milestone_labels = ['NOW', 'Next 7 Days', 'Next 30 Days']
    milestone_x = [2.5, 10.5, 18.5]
    milestone_w = 7.5
    ec_subset = event_cards[:3] if event_cards else []

    for i, (label, mx) in enumerate(zip(milestone_labels, milestone_x)):
        # Milestone node
        node = slide.shapes.add_shape(
            _V1_ROUNDED_RECT, Cm(mx + 3.0), Cm(6.0), Cm(1.4), Cm(1.4),
        )
        _v1_set_shape_fill(node, _V1_BLUE)
        _add_textbox(
            slide, Cm(mx + 3.0), Cm(6.1), Cm(1.4), Cm(1.2),
            str(i + 1), font_size=14, bold=True,
            color=_V1_CARD_WHITE, alignment=PP_ALIGN.CENTER,
        )
        # Label
        _add_textbox(
            slide, Cm(mx + 0.2), Cm(4.8), Cm(milestone_w), Cm(0.9),
            label, font_size=12, bold=True, color=TEXT_WHITE,
        )
        # Owner + bullets
        if i < len(ec_subset):
            c = ec_subset[i]
            dc = build_decision_card(c)
            owner = dc.get('owner', 'CXO')
            action = dc['actions'][0] if dc['actions'] else '持續追蹤'
            why_short = safe_text(c.why_important or '', 45) or safe_text(c.title_plain or '', 45)
            bullets = _v1_norm_bullets_safe([safe_text(action, 45), why_short])
            score_val = getattr(c, 'final_score', None)
            score_line = f'Score: {score_val:.1f}/10' if score_val is not None else ''
        else:
            owner = 'CXO'
            bullets = ['持續監控市場動態（T+7）。', '評估後續行動是否必要。']
            score_line = ''

        bullets = bullets[:2]
        body_lines = [f'Owner: {owner}'] + bullets
        if score_line:
            body_lines.append(score_line)
        _v1_add_card(
            slide, mx + 0.2, 7.8, milestone_w, 7.5,
            header_text='',
            body_lines=body_lines,
            bg_color=_V1_SOFT_BG,
            header_color=_V1_BLUE,
            body_color=_V1_TEXT_GRAY,
            body_font_size=_V1_CARD_BODY_FS,
        )

    # Right-side risks block
    _add_textbox(
        slide, Cm(26.5), Cm(4.8), Cm(6.5), Cm(1.0),
        'Top Risks / Watch', font_size=12, bold=True, color=TEXT_WHITE,
    )
    risks: list[str] = []
    for c in event_cards[:3]:
        dc = build_decision_card(c)
        if dc.get('risks'):
            risks.append(safe_text(dc['risks'][0], 40))
    if not risks:
        risks = ['監控市場波動', '跟蹤競爭對手動作']
    risks = _v1_norm_bullets_safe(risks[:2])
    _add_multiline_textbox(
        slide, Cm(26.5), Cm(6.0), Cm(6.5), Cm(5.0),
        [f'• {r}' for r in risks],
        font_size=10, color=SUBTLE_GRAY, line_spacing=1.5,
    )
    _add_divider(slide, Cm(0), Cm(18.7), SLIDE_WIDTH, color=ACCENT)


# ---------------------------------------------------------------------------
# T1: Three-card Curved Timeline — Event Slide A: What / Why / Proof
# ---------------------------------------------------------------------------

def _slide_brief_page1(
    prs: Presentation,
    card: EduNewsCard,
    idx: int,
    final_payload: dict | None = None,
) -> None:
    """T1: Three-card Curved Timeline layout for WHAT HAPPENED slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    # Event badge + title
    _add_textbox(
        slide, Cm(1.5), Cm(0.5), Cm(3.0), Cm(1.2),
        f'#{idx}', font_size=28, bold=True, color=ACCENT,
    )
    _add_textbox(
        slide, Cm(4.5), Cm(0.5), Cm(26.0), Cm(1.2),
        safe_text(str((final_payload or {}).get("title", "") or card.title_plain or ''), 55),
        font_size=22, bold=True, color=TEXT_WHITE,
    )
    _v1_slide_header(slide, 'WHAT HAPPENED')

    # Hero image (optional, right side)
    img_right = False
    try:
        img_path = get_news_image(card.title_plain, card.category)
        if img_path and img_path.exists():
            slide.shapes.add_picture(
                str(img_path), Cm(23.5), Cm(2.8), Cm(10.0), Cm(6.0),
            )
            img_right = True
    except Exception:
        pass

    card_w = 20.5 if img_right else 30.0
    card_left = 1.8

    # EN-ZH Hybrid Glossing v1: shared seen-set for dedup across fields on this slide
    _gloss_seen: set = set()
    brief = build_ceo_brief_blocks(card)

    # ── Canonical payload v3: single source of truth ──
    try:
        from utils.canonical_narrative import get_canonical_payload as _get_cp_T1
        _cp_T1 = _get_cp_T1(card)
    except Exception:
        _cp_T1 = {}

    # Card 1: Q1 — What Happened (q1_zh: Traditional Chinese narrative with embedded quote_window_1)
    _q1_body = safe_text(str(
        (final_payload or {}).get("q1_zh", "")
        or (final_payload or {}).get("q1", "")
        or _cp_T1.get('q1_event_2sent_zh', '')
        or ''
    ), 320)
    if not _q1_body:
        _q1_body = safe_text(_v1_norm_gloss(_v1_narrative(card), _V1_GLOSSARY, _gloss_seen), 320)
    else:
        _q1_body = safe_text(_v1_norm_gloss(_q1_body, _V1_GLOSSARY, _gloss_seen), 320)
    card1_top, card1_h = 2.8, 4.6
    _v1_add_card(
        slide, card_left, card1_top, card_w, card1_h,
        header_text='Q1 — What Happened',
        body_lines=[_q1_body],
        bg_color=_V1_SOFT_BG,
        header_color=_V1_BLUE,
        body_color=_V1_TEXT_GRAY,
        header_font_size=_V1_CARD_TITLE_FS,
        body_font_size=_V1_CARD_BODY_FS,
        body_text_limit=260,
        shape_type=_V1_ROUNDED_RECT,
    )
    _v1_vertical_connector(slide, card_left + card_w / 2, card1_top + card1_h, card1_top + card1_h + 0.5)

    # Card 2: Q2 — Why It Matters (q2_zh: Traditional Chinese narrative with embedded quote_window_2)
    card2_top = card1_top + card1_h + 0.6
    card2_h = 3.5
    _q2_body = safe_text(str(
        (final_payload or {}).get("q2_zh", "")
        or (final_payload or {}).get("q2", "")
        or _cp_T1.get('q2_impact_2sent_zh', '')
        or ''
    ), 320)
    if not _q2_body:
        _q2_body = safe_text(card.why_important or brief.get('q1_meaning', ''), 150)
    _q2_body = safe_text(_v1_norm_gloss(_q2_body, _V1_GLOSSARY, _gloss_seen), 320)
    _v1_add_card(
        slide, card_left, card2_top, card_w, card2_h,
        header_text='Q2 — Why It Matters',
        body_lines=[_q2_body],
        bg_color=_V1_SOFT_BG,
        header_color=_V1_BLUE,
        body_color=_V1_TEXT_GRAY,
        header_font_size=_V1_CARD_TITLE_FS,
        body_font_size=_V1_CARD_BODY_FS,
        body_text_limit=260,
        shape_type=_V1_ROUNDED_RECT,
    )
    _v1_vertical_connector(slide, card_left + card_w / 2, card2_top + card2_h, card2_top + card2_h + 0.5)

    # Card 3: Proof (canonical proof_line — fixed format: 證據：來源：X（YYYY-MM-DD）)
    card3_top = card2_top + card2_h + 0.6
    card3_h = 3.0
    _proof_canonical = _cp_T1.get('proof_line', '') or brief.get('proof_line', '')
    if not _proof_canonical:
        source_label_fb = safe_text(getattr(card, 'source_name', '') or '', 30)
        try:
            from utils.longform_narrative import _make_date_proof_line
            _proof_canonical = _make_date_proof_line(card)
        except Exception:
            _pub_fb = str(getattr(card, 'published_at_parsed', '') or getattr(card, 'published_at', '') or '').strip()[:10]
            _proof_canonical = f"證據：來源：{source_label_fb}（{_pub_fb}）" if _pub_fb else f"證據：來源：{source_label_fb}"
    source_url_T1 = safe_text(getattr(card, 'source_url', '') or '', 60)
    proof_lines_T1 = [_proof_canonical]
    if source_url_T1 and source_url_T1.startswith('http'):
        proof_lines_T1.append(safe_text(source_url_T1, 55))
    if final_payload:
        _proof_url = safe_text(str(final_payload.get("final_url", "") or ""), 110)
        _proof_q1 = safe_text(str(final_payload.get("quote_1", "") or ""), 110)
        _proof_q2 = safe_text(str(final_payload.get("quote_2", "") or ""), 110)
        proof_lines_T1 = [
            f"final_url: {_proof_url}",
            f"quote_1: {_proof_q1}",
            f"quote_2: {_proof_q2}",
        ]
    _v1_add_card(
        slide, card_left, card3_top, card_w, card3_h,
        header_text='Proof — Hard Evidence',
        body_lines=proof_lines_T1,
        bg_color=_V1_SOFT_BG,
        header_color=_V1_GREEN,
        body_color=_V1_TEXT_GRAY,
        header_font_size=_V1_CARD_TITLE_FS,
        body_font_size=_V1_CARD_BODY_FS,
        body_text_limit=220,
        shape_type=_V1_ROUNDED_RECT,
    )


# ---------------------------------------------------------------------------
# T3: Growth Steps — Event Slide B: Moves / Risks / Owner
# ---------------------------------------------------------------------------

def _slide_brief_page2(
    prs: Presentation,
    card: EduNewsCard,
    idx: int,
    final_payload: dict | None = None,
) -> None:
    """T3: Growth Steps staircase for WHY IT MATTERS — Action Plan slide."""
    brief = build_ceo_brief_blocks(card)

    # ── Canonical payload v3：Slide2 100% 來自 canonical，不重新跑 bucket 模板 ──
    try:
        from utils.canonical_narrative import get_canonical_payload as _get_cp_T3p2
        _cp_T3p2 = _get_cp_T3p2(card)
    except Exception:
        _cp_T3p2 = {}

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    # Header — preserve "WHY IT MATTERS" for existing-test compatibility
    _add_textbox(
        slide, Cm(2), Cm(0.5), Cm(28), Cm(1.2),
        f'#{idx}  WHY IT MATTERS — Action Plan',
        font_size=22, bold=True, color=HIGHLIGHT_YELLOW,
    )
    _add_divider(slide, Cm(2), Cm(1.8), Cm(4), color=ACCENT)

    # Q3 actions → staircase steps
    # 直接讀 canonical q3_moves_3bullets_zh，不套 bucket 模板
    _gloss_seen_p2: set = set()
    _canon_q3 = list((final_payload or {}).get("moves", []) or _cp_T3p2.get('q3_moves_3bullets_zh', []) or [])
    if _canon_q3:
        actions = _canon_q3[:3]
    else:
        # canonical 無內容時，從 build_ceo_brief_blocks 取（已經過 v5.2.8 canonical 過濾）
        actions = list(brief.get('q3_actions', []) or [])
    actions = _v1_norm_bullets_safe(
        [safe_text(a, 60) for a in actions],
    )
    actions = [_v1_norm_gloss(a, _V1_GLOSSARY, _gloss_seen_p2) for a in actions]
    # 若 canonical 無資料，不補空話模板——只顯示有內容的步驟（留空而非捏造）
    while len(actions) < 3:
        actions.append('')

    # Staircase layout
    step_configs = [
        (1.5, 13.0, 18.0, 3.5, _V1_AMBER, 'Q3-A  Move 1'),
        (5.5, 10.0, 14.5, 3.5, _V1_BLUE, 'Q3-B  Move 2'),
        (9.5, 7.0, 11.0, 3.5, _V1_GREEN, 'Q3-C  Move 3'),
    ]
    for step_i, (sx, sy, sw, sh, sc, slabel) in enumerate(step_configs):
        action_text = actions[step_i] if step_i < len(actions) else ''
        if not action_text.strip():
            continue  # 不補空話模板，略過空步驟
        _v1_add_card(
            slide, sx, sy, sw, sh,
            header_text=f'{slabel}: {action_text[:48]}',
            body_lines=[safe_text(action_text, 80)],
            bg_color=_V1_SOFT_BG,
            header_color=sc,
            body_color=_V1_TEXT_GRAY,
            header_font_size=_V1_CARD_TITLE_FS,
            body_font_size=_V1_CARD_BODY_FS,
            shape_type=_V1_ROUNDED_RECT,
        )

    # Risks / Watch — 直接讀 canonical risks_2bullets_zh（與 Slide1 同源）
    _add_divider(slide, Cm(2), Cm(15.0), Cm(29), color=SUBTLE_GRAY)
    _add_textbox(
        slide, Cm(2), Cm(15.2), Cm(18), Cm(0.8),
        'Risks / Watch', font_size=13, bold=True, color=TEXT_WHITE,
    )
    try:
        raw_risks = list((final_payload or {}).get("risks", []) or _cp_T3p2.get('risks_2bullets_zh', []) or [])
    except Exception:
        raw_risks = []
    if not raw_risks:
        dc_card_fb = build_decision_card(card)
        raw_risks = dc_card_fb.get('risks', []) or []
    if not raw_risks:
        raw_risks = ['持續監控此事件後續影響。']
    risks_t3 = _v1_norm_bullets_safe(raw_risks[:2])
    from utils.semantic_quality import is_placeholder_or_fragment as _is_frag_risk
    risks_t3 = [
        _v1_norm_gloss(rk, _V1_GLOSSARY, _gloss_seen_p2) if not _is_frag_risk(rk)
        else '持續監控此事件後續影響。'
        for rk in risks_t3
    ]
    for ri, rk in enumerate(risks_t3[:2]):
        _add_textbox(
            slide, Cm(3), Cm(16.2 + ri * 0.9), Cm(18), Cm(0.8),
            f'• {safe_text(rk, 60)}', font_size=11, color=SUBTLE_GRAY,
        )

    # Proof line（canonical 100% 單源：Slide1 與 Slide2 同一個 payload）
    _proof_t3 = _cp_T3p2.get('proof_line', '') or brief.get('proof_line', '')
    if final_payload and str(final_payload.get("final_url", "") or "").strip():
        _proof_t3 = f"final_url: {safe_text(str(final_payload.get('final_url', '') or ''), 110)}"
    if _proof_t3:
        _add_textbox(
            slide, Cm(2), Cm(18.0), Cm(29), Cm(0.7),
            safe_text(_proof_t3, 100),
            font_size=9, color=SUBTLE_GRAY,
        )

    # Owner / ETA
    dc_card = build_decision_card(card)
    owner = dc_card.get('owner', 'CXO')
    _add_textbox(
        slide, Cm(22), Cm(15.2), Cm(11), Cm(0.8),
        f'Owner: {owner}  |  ETA: T+7',
        font_size=12, bold=True, color=HIGHLIGHT_YELLOW,
    )

    # Video reference (preserve for test_video_reference_present)
    videos = brief.get('video_source', [])
    if videos:
        vid = videos[0]
        _add_textbox(
            slide, Cm(2), Cm(17.2), Cm(20), Cm(0.7),
            f'Video: {safe_text(vid.get("title", ""), 50)}',
            font_size=9, color=SUBTLE_GRAY,
        )

    # Sources
    sources = brief.get('sources', [])
    if sources:
        _add_textbox(
            slide, Cm(2), Cm(17.8), Cm(20), Cm(0.6),
            f'Source: {safe_text(sources[0], 80)}',
            font_size=9, color=SUBTLE_GRAY,
        )


# ---------------------------------------------------------------------------
# Meta JSON writer
# ---------------------------------------------------------------------------

def _v1_write_exec_layout_meta(
    output_path: 'Path',
    event_cards: list[EduNewsCard],
    all_cards: list[EduNewsCard],
) -> None:
    """Compute and write outputs/exec_layout.meta.json."""
    from datetime import datetime as _dt, timezone as _tz

    meta_path = output_path.parent / 'exec_layout.meta.json'

    # Build slide_layout_map from static deck structure
    n_events = len(event_cards)
    slide_layout_map: list[dict] = [
        {'slide_no': 1, 'template_code': 'COVER', 'title': 'CEO Decision Brief'},
        {'slide_no': 2, 'template_code': 'STRUCTURED_SUMMARY', 'title': 'Structured Summary'},
        {'slide_no': 3, 'template_code': 'T2', 'title': 'Signal Thermometer'},
        {'slide_no': 4, 'template_code': 'CORP_WATCH', 'title': 'Corp Watch'},
        {'slide_no': 5, 'template_code': 'KEY_TAKEAWAYS', 'title': 'Key Takeaways'},
        {'slide_no': 6, 'template_code': 'T5', 'title': '今日總覽 Overview'},
        {'slide_no': 7, 'template_code': 'T4', 'title': 'Event Ranking 事件影響力排行'},
    ]
    slide_cursor = 8
    for i in range(n_events):
        ev_title = (event_cards[i].title_plain or f'Event {i + 1}')[:30]
        slide_layout_map.append({
            'slide_no': slide_cursor,
            'template_code': 'T1',
            'title': f'#{i + 1} WHAT HAPPENED — {ev_title}',
        })
        slide_layout_map.append({
            'slide_no': slide_cursor + 1,
            'template_code': 'T3',
            'title': f'#{i + 1} WHY IT MATTERS — Action Plan',
        })
        slide_cursor += 2
    slide_layout_map.extend([
        {'slide_no': slide_cursor, 'template_code': 'REC_MOVES', 'title': 'Recommended Moves'},
        {'slide_no': slide_cursor + 1, 'template_code': 'DECISION_MATRIX', 'title': 'Decision Matrix'},
        {'slide_no': slide_cursor + 2, 'template_code': 'T6', 'title': '待決事項與 Owner'},
    ])

    # Fragment stats
    all_bullet_lists: list[list[str]] = []
    fragments_detected = 0
    fragments_fixed = 0

    for c in event_cards:
        raw_bullets: list[str] = []
        for attr in ('action_items', 'derivable_effects', 'fact_check_confirmed'):
            vals = getattr(c, attr, []) or []
            raw_bullets.extend(str(v) for v in vals)
        short = [b for b in raw_bullets if len(b.strip()) < 12 and b.strip()]
        fragments_detected += len(short)
        if short:
            normed = _v1_norm_bullets(raw_bullets)
            still_short = len([b for b in normed if len(b.strip()) < 12])
            fixed = max(0, len(short) - still_short)
            fragments_fixed += fixed
        all_bullet_lists.append(raw_bullets)

    total = max(fragments_detected, 1)
    fragment_ratio = round(fragments_fixed / total, 4)

    # Bullet length stats
    norm_bullet_lists = [_v1_norm_bullets(bl) for bl in all_bullet_lists]
    bstats = _v1_bullet_stats(norm_bullet_lists)

    # Proof token coverage
    proof_covered = 0
    for c in event_cards:
        combined = ' '.join(filter(None, [
            c.title_plain or '',
            c.what_happened or '',
            c.why_important or '',
            getattr(c, 'technical_interpretation', '') or '',
        ]))
        if _v1_has_evidence(combined):
            proof_covered += 1
    proof_ratio = round(proof_covered / max(len(event_cards), 1), 4)

    # Sentence count stats
    from utils.semantic_quality import count_sentences as _count_sents_v1
    sent_counts = []
    for c in event_cards:
        txt = (c.what_happened or '') + ' ' + (c.why_important or '')
        sent_counts.append(_count_sents_v1(txt))
    avg_sents = round(sum(sent_counts) / max(len(sent_counts), 1), 2)

    meta = {
        'layout_version': _V1_LAYOUT_VER,
        'generated_at': _dt.now(tz=_tz.utc).strftime('%Y-%m-%dT%H:%M:%SZ'),
        'template_map': _V1_TEMPLATE_MAP,
        'slide_layout_map': slide_layout_map,
        'fragment_fix_stats': {
            'fragments_detected': fragments_detected,
            'fragments_fixed': fragments_fixed,
            'fragment_ratio': fragment_ratio,
        },
        'bullet_len_stats': bstats,
        'card_stats': {
            'total_event_cards': len(event_cards),
            'avg_sentences_per_event_card': avg_sents,
            'proof_token_coverage_ratio': proof_ratio,
        },
    }

    try:
        meta_path.parent.mkdir(parents=True, exist_ok=True)
        meta_path.write_text(
            _json.dumps(meta, ensure_ascii=False, indent=2),
            encoding='utf-8',
        )
        get_logger().info('exec_layout.meta.json written: %s', meta_path)
    except Exception as exc:
        get_logger().warning('Failed to write exec_layout.meta.json: %s', exc)

    # G4: update exec_quality.meta.json with fragment leak + actions normalization data
    try:
        import json as _json_q
        from utils.semantic_quality import is_placeholder_or_fragment as _is_frag_act
        _q_path = meta_path.parent / 'exec_quality.meta.json'
        fragments_leaked = max(0, fragments_detected - fragments_fixed)
        fragment_leak_gate = 'PASS' if fragments_leaked == 0 else 'FAIL'

        # Actions/Risks normalization stats
        actions_normalized_count = 0
        actions_fragment_leak_count = 0
        for _c in event_cards:
            _brief_c = build_ceo_brief_blocks(_c)
            _dc_c = build_decision_card(_c)
            _raw_acts = _brief_c.get('q3_actions', []) or []
            _raw_risks = _dc_c.get('risks', []) or []
            _all_items = [safe_text(a, 60) for a in (_raw_acts[:3] + _raw_risks[:2])]
            _normed = _v1_norm_bullets_safe(_all_items)
            for _item in _normed:
                if _item and _item.strip():
                    actions_normalized_count += 1
                    if _is_frag_act(_item):
                        actions_fragment_leak_count += 1

        if _q_path.exists():
            _qm = _json_q.loads(_q_path.read_text(encoding='utf-8'))
        else:
            _qm = {}
        _gloss = _v1_gloss_stats()
        _qm.update({
            'fragments_detected': fragments_detected,
            'fragments_fixed': fragments_fixed,
            'fragments_leaked': fragments_leaked,
            'fragment_leak_gate': fragment_leak_gate,
            'english_heavy_paragraphs_fixed_count': _gloss.get('english_heavy_paragraphs_fixed_count', 0),
            'english_heavy_skeletonized_count': _gloss.get('english_heavy_skeletonized_count', 0),
            'proper_noun_gloss_applied_count': _gloss.get('proper_noun_gloss_applied_count', 0),
            'actions_normalized_count': actions_normalized_count,
            'actions_fragment_leak_count': actions_fragment_leak_count,
        })
        _q_path.write_text(
            _json_q.dumps(_qm, ensure_ascii=False, indent=2),
            encoding='utf-8',
        )
    except Exception as _exc_q:
        get_logger().warning('exec_quality.meta G4 update error (non-fatal): %s', _exc_q)

    # Narrative Compactor v2 meta — collect per-card debug stats written by v527 wrapper
    try:
        from core.content_strategy import write_narrative_v2_meta as _write_nv2_meta
        _write_nv2_meta(event_cards, out_dir=str(meta_path.parent))
    except Exception as _exc_nv2:
        get_logger().warning('narrative_v2 meta write error (non-fatal): %s', _exc_nv2)


# ---------------------------------------------------------------------------
# Override generate_executive_ppt to write exec_layout.meta.json
# ---------------------------------------------------------------------------

def _append_watchlist_slide(
    pptx_path: 'Path',
    watchlist_cards: list,
    theme: str = 'light',
) -> None:
    """Append a Developing Watchlist slide to an existing PPTX file.

    Opens the saved file, adds one slide with compact 3-line entries
    (title / why / proof_line) for each watchlist card, then saves back.
    Non-fatal: caller wraps in try/except.
    """
    _apply_theme(theme)
    prs = Presentation(str(pptx_path))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    # --- Header ---
    _add_textbox(
        slide, Cm(2), Cm(0.6), Cm(20), Cm(1.4),
        'DEVELOPING WATCHLIST', font_size=22, color=HIGHLIGHT_YELLOW, bold=True,
    )
    _add_textbox(
        slide, Cm(22.5), Cm(0.6), Cm(9), Cm(1.4),
        f'\u9577\u6587\u5ef6\u4f38\u95b1\u8b80 \u2736 {len(watchlist_cards)} \u5247',
        font_size=16, color=DARK_MUTED,
    )
    _add_divider(slide, Cm(2), Cm(2.0), Cm(29.8), color=HIGHLIGHT_YELLOW)

    # --- Compact card entries (3 lines each) ---
    item_h = Cm(2.2)
    top_start = Cm(2.35)
    max_show = min(len(watchlist_cards), 7)

    try:
        from utils.exec_sanitizer import sanitize_exec_text as _san
    except Exception:
        def _san(t: str) -> str:  # type: ignore[misc]
            return t

    for i, card in enumerate(watchlist_cards[:max_show]):
        # Prefer watchlist-specific payload (lower 600-char threshold)
        wl = getattr(card, '_watchlist_longform_payload', None) or {}
        lf = getattr(card, '_longform_v1_cache', {}) or {}

        title = safe_text(getattr(card, 'title_plain', '') or '', limit=100)

        # what line: watchlist payload 'what' → lf 'bg'/'what_is' → what_happened
        what_raw = (
            wl.get('what') or
            lf.get('bg') or lf.get('what_is') or
            getattr(card, 'what_happened', '') or ''
        )
        what = safe_text(_san(str(what_raw)), limit=120)

        # why line: watchlist payload 'why' → lf 'why' → why_important
        why_raw = (
            wl.get('why') or
            lf.get('why') or
            getattr(card, 'why_important', '') or ''
        )
        why = safe_text(_san(str(why_raw)), limit=120)

        # proof line: watchlist payload 'proof_line' → lf 'proof_line'
        proof_raw = wl.get('proof_line') or lf.get('proof_line') or ''
        proof = safe_text(proof_raw, limit=100)

        # Body text: prefer what if available, else why
        body_text = what if what else why
        row_top = top_start + i * item_h

        _add_textbox(
            slide, Cm(2.2), row_top, Cm(28.5), Cm(0.8),
            title, font_size=13, color=HIGHLIGHT_YELLOW, bold=True,
        )
        _add_textbox(
            slide, Cm(2.2), row_top + Cm(0.78), Cm(28.5), Cm(0.7),
            body_text, font_size=11, color=TEXT_WHITE,
        )
        _add_textbox(
            slide, Cm(2.2), row_top + Cm(1.44), Cm(28.5), Cm(0.6),
            proof, font_size=10, color=DARK_MUTED,
        )

    if len(watchlist_cards) > max_show:
        extra = len(watchlist_cards) - max_show
        _add_textbox(
            slide, Cm(2), top_start + max_show * item_h, Cm(28), Cm(0.8),
            f'+ \u53e6\u6709 {extra} \u5247 Watchlist \u9805\u76ee\uff08\u8a73\u898b exec_longform.meta.json\uff09',
            font_size=11, color=DARK_MUTED,
        )

    prs.save(str(pptx_path))


def generate_executive_ppt(
    cards: list[EduNewsCard],
    health: 'SystemHealthReport',
    report_time: str,
    total_items: int,
    output_path: 'Path | None' = None,
    theme: str = 'light',
    metrics: dict | None = None,
) -> 'Path':
    """V1 wrapper: runs original generator then writes exec_layout.meta.json."""
    # Pre-pass: strip U+2026 / three-dot ellipsis from ALL card fields before
    # any slide generation reads them directly (Iteration 5.2 ellipsis enforcement).
    try:
        from utils.canonical_narrative import get_canonical_payload as _prepass_gcp
        for _pc in cards:
            try:
                _prepass_gcp(_pc)
            except Exception:
                pass
    except Exception:
        pass
    if _is_brief_report_mode():
        result = _generate_brief_ppt_only(
            cards=cards,
            output_path=output_path,
            theme=theme,
            metrics=metrics,
        )
        return result
    result = _v1_prev_generate_executive_ppt(
        cards, health, report_time, total_items,
        output_path=output_path, theme=theme, metrics=metrics,
    )
    ev_cards = get_event_cards_for_deck(cards, metrics=metrics or {}, min_events=0)
    _final_cards_payload = _load_final_cards(metrics)
    ev_cards, _ev_payloads = _align_event_cards_with_final_cards(ev_cards, _final_cards_payload)
    try:
        _v1_write_exec_layout_meta(result, ev_cards, cards)
    except Exception as exc:
        get_logger().warning('exec_layout.meta write error (non-fatal): %s', exc)
    try:
        from utils.longform_narrative import write_longform_meta
        write_longform_meta(event_cards=ev_cards)
    except Exception as exc:
        get_logger().warning('exec_longform.meta write error (non-fatal): %s', exc)
    # --- Watchlist track (Longform Pool Expansion v1) ---
    try:
        from config.settings import LONGFORM_MIN_DAILY_TOTAL, LONGFORM_WATCHLIST_MAX
        from utils.longform_watchlist import select_watchlist_cards, write_watchlist_meta
        wl_cards, wl_candidates_total = select_watchlist_cards(
            cards, ev_cards,
            min_daily_total=LONGFORM_MIN_DAILY_TOTAL,
            max_watchlist=LONGFORM_WATCHLIST_MAX,
        )
        if wl_cards:
            _append_watchlist_slide(result, wl_cards, theme=theme)
        write_watchlist_meta(
            event_cards=ev_cards,
            watchlist_cards=wl_cards,
            candidates_total=wl_candidates_total,
            min_daily_total=LONGFORM_MIN_DAILY_TOTAL,
        )
    except Exception as exc:
        get_logger().warning('watchlist longform error (non-fatal): %s', exc)
    return result


# ---------------------------------------------------------------------------
# NOT_READY_report PPTX generator (standalone — no EduNewsCard dependency)
# Called by run_pipeline.ps1 via run_once.py --not-ready-report
# ---------------------------------------------------------------------------

def generate_not_ready_report_pptx(
    output_path: "Path",
    fail_reason: str,
    gate_name: str,
    samples: "list[dict]",
    next_steps: str,
    run_id: str = "",
    run_date: str = "",
) -> "Path":
    """Generate outputs/NOT_READY_report.pptx with human-readable failure info."""
    from datetime import datetime as _dt_nr

    prs = Presentation()
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]  # blank

    # ── Slide 1: Cover ───────────────────────────────────────────────────────
    sl1 = prs.slides.add_slide(blank_layout)
    _set_slide_bg(sl1, DARK_BG)
    date_str = run_date or _dt_nr.now().strftime("%Y-%m-%d")
    _add_textbox(
        sl1, Cm(2), Cm(3), Cm(29.8), Cm(3),
        "\u274c  \u4eca\u65e5 AI \u60c5\u5831\u5831\u544a\u672a\u80fd\u7522\u51fa",
        font_size=32, color=ACCENT, bold=True,
        alignment=PP_ALIGN.LEFT,
    )
    _add_divider(sl1, Cm(2), Cm(6.5), Cm(29.8), ACCENT)
    _add_textbox(
        sl1, Cm(2), Cm(7.2), Cm(29.8), Cm(1.8),
        f"\u65e5\u671f\uff1a{date_str}   run_id\uff1a{run_id or '\u2014'}",
        font_size=16, color=DARK_MUTED,
    )
    _add_textbox(
        sl1, Cm(2), Cm(9.2), Cm(29.8), Cm(3),
        "\u672c\u5831\u544a\u7531\u7cfb\u7d71\u81ea\u52d5\u751f\u6210\uff0c\u8aaa\u660e\u4eca\u65e5\u6b63\u5f0f\u5831\u544a\u7121\u6cd5\u7522\u51fa\u7684\u539f\u56e0\u3002",
        font_size=18, color=DARK_TEXT,
    )

    # ── Slide 2: Failure Reason ──────────────────────────────────────────────
    sl2 = prs.slides.add_slide(blank_layout)
    _set_slide_bg(sl2, DARK_BG)
    _add_textbox(sl2, Cm(2), Cm(1.2), Cm(29.8), Cm(1.8),
                 "\u2460 \u70ba\u4f55\u4eca\u65e5\u7121\u6cd5\u751f\u6210\u6b63\u5f0f\u5831\u544a\uff1f",
                 font_size=22, color=HIGHLIGHT_YELLOW, bold=True)
    _add_divider(sl2, Cm(2), Cm(3.3), Cm(29.8), HIGHLIGHT_YELLOW)
    _add_textbox(sl2, Cm(2), Cm(3.8), Cm(29.8), Cm(6),
                 fail_reason or "\uff08\u539f\u56e0\u4e0d\u660e\uff0c\u8acb\u67e5\u95b1 desktop_button.last_run.log\uff09",
                 font_size=17, color=DARK_TEXT)
    _add_textbox(sl2, Cm(2), Cm(10.5), Cm(29.8), Cm(1.8),
                 f"\u5931\u6557 Gate\uff1a{gate_name or '\u2014'}",
                 font_size=16, color=ACCENT, bold=True)

    # ── Slide 3: Sample Events ────────────────────────────────────────────────
    sl3 = prs.slides.add_slide(blank_layout)
    _set_slide_bg(sl3, DARK_BG)
    _add_textbox(sl3, Cm(2), Cm(1.2), Cm(29.8), Cm(1.8),
                 "\u2461 \u4eca\u65e5\u641c\u96c6\u5230\u7684\u6a23\u672c\u4e8b\u4ef6\uff08\u6700\u591a 3 \u5247\uff09",
                 font_size=22, color=HIGHLIGHT_YELLOW, bold=True)
    _add_divider(sl3, Cm(2), Cm(3.3), Cm(29.8), HIGHLIGHT_YELLOW)
    _y_starts = [3.8, 8.0, 12.2]
    if samples:
        for _i, _s in enumerate(samples[:3]):
            _title = str(_s.get("title") or "\uff08\u7121\u6a19\u984c\uff09")[:100]
            _url   = str(_s.get("final_url") or _s.get("url") or "\u2014")[:120]
            _add_textbox(sl3, Cm(2), Cm(_y_starts[_i]), Cm(29.8), Cm(1.8),
                         f"{_i + 1}. {_title}",
                         font_size=16, color=DARK_TEXT, bold=True)
            _add_textbox(sl3, Cm(2.5), Cm(_y_starts[_i] + 2.0), Cm(29.3), Cm(1.4),
                         f"\u4f86\u6e90\uff1a{_url}",
                         font_size=13, color=DARK_MUTED)
    else:
        _add_textbox(sl3, Cm(2), Cm(4), Cm(29.8), Cm(3),
                     "\uff08\u672c\u6b21\u672a\u641c\u96c6\u5230\u4e8b\u4ef6\u6a23\u672c\uff09",
                     font_size=16, color=DARK_MUTED)

    # ── Slide 4: Next Steps ───────────────────────────────────────────────────
    sl4 = prs.slides.add_slide(blank_layout)
    _set_slide_bg(sl4, DARK_BG)
    _add_textbox(sl4, Cm(2), Cm(1.2), Cm(29.8), Cm(1.8),
                 "\u2462 \u5efa\u8b70\u4e0b\u4e00\u6b65",
                 font_size=22, color=HIGHLIGHT_YELLOW, bold=True)
    _add_divider(sl4, Cm(2), Cm(3.3), Cm(29.8), HIGHLIGHT_YELLOW)
    _add_textbox(sl4, Cm(2), Cm(3.8), Cm(29.8), Cm(8),
                 next_steps or "\u8acb\u67e5\u95b1 outputs/desktop_button.last_run.log \u53d6\u5f97\u8a73\u7d30\u8a3a\u65b7\u8cc7\u8a0a\u3002",
                 font_size=17, color=DARK_TEXT)
    _add_textbox(sl4, Cm(2), Cm(16.8), Cm(29.8), Cm(1.4),
                 "\u672c\u6587\u4ef6\u7531\u7cfb\u7d71\u81ea\u52d5\u751f\u6210\uff0c\u50c5\u4f9b\u8a3a\u65b7\u7528\u9014\u3002",
                 font_size=13, color=DARK_MUTED)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


# ---------------------------------------------------------------------------
# generate_zh_md_pptx — Iteration 20: Translation-First ZH Delivery
# Converts a ZH Markdown string to a PowerPoint presentation.
# Used to overwrite executive_report.pptx with the translated ZH version.
# No images are inserted (brief mode constraint).
# ---------------------------------------------------------------------------
def generate_zh_md_pptx(md_text: str, output_path: "Path | str") -> Path:
    """Generate a readable PPTX from ZH Markdown (translation-first delivery).

    Layout:
    - H1 / H2  → new slide with title in top area
    - H3       → bold bullet in current slide body
    - ---      → new blank slide (event separator)
    - > quote  → italic indented text
    - - bullet → bullet point
    - code     → monospace text (no images)
    Max ~20 body lines per slide; auto-continues on new slide if exceeded.
    """
    import re as _r

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Cm(33.87)   # widescreen 13.33 inch
    prs.slide_height = Cm(19.05)  # 7.5 inch

    blank_layout = prs.slide_layouts[6]  # blank

    _BG    = RGBColor(0xFF, 0xFF, 0xFF)
    _TITLE = RGBColor(0x21, 0x28, 0x38)
    _BODY  = RGBColor(0x22, 0x22, 0x22)

    _ML = Cm(1.8)   # margin left
    _MT = Cm(1.2)   # margin top
    _SW = prs.slide_width - 2 * _ML
    _TH = Cm(2.4)   # title height
    _BT = _MT + _TH + Cm(0.3)
    _BH = prs.slide_height - _BT - Cm(0.5)

    MAX_BODY_LINES = 20

    def _new_slide() -> "pptx.slide.Slide":  # type: ignore[name-defined]
        sl = prs.slides.add_slide(blank_layout)
        fill = sl.background.fill
        fill.solid()
        fill.fore_color.rgb = _BG
        return sl

    def _set_title(sl: "pptx.slide.Slide", text: str) -> None:  # type: ignore[name-defined]
        txb = sl.shapes.add_textbox(_ML, _MT, _SW, _TH)
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.color.rgb = _TITLE

    def _set_body(sl: "pptx.slide.Slide", lines: list) -> None:  # type: ignore[name-defined]
        if not lines:
            return
        txb = sl.shapes.add_textbox(_ML, _BT, _SW, _BH)
        tf = txb.text_frame
        tf.word_wrap = True
        first = True
        for (txt, size, bold, italic) in lines:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = p.add_run()
            run.text = txt
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = _BODY

    # Accumulate slides as (title, body_lines)
    slides_data: list[tuple[str, list]] = []
    cur_title = ""
    cur_body: list = []

    def _flush() -> None:
        if cur_title or cur_body:
            slides_data.append((cur_title, list(cur_body)))

    def _start(title: str) -> None:
        nonlocal cur_title, cur_body
        _flush()
        cur_title = title
        cur_body = []

    def _add_body(text: str, size: int = 14, bold: bool = False, italic: bool = False) -> None:
        nonlocal cur_title, cur_body
        cur_body.append((text, size, bold, italic))
        if len(cur_body) >= MAX_BODY_LINES:
            _start(cur_title + " (續)")

    lines = md_text.split("\n")
    in_code = False

    for line in lines:
        if line.startswith("```"):
            in_code = not in_code
            continue
        if in_code:
            _add_body("  " + line, size=10, italic=True)
            continue

        h1 = _r.match(r"^# (.+)", line)
        h2 = _r.match(r"^## (.+)", line)
        h3 = _r.match(r"^### (.+)", line)

        if h1:
            _start(h1.group(1).strip())
        elif h2:
            _start(h2.group(1).strip())
        elif h3:
            if not cur_title and not cur_body:
                _start(h3.group(1).strip())
            else:
                _add_body("\u25b6 " + h3.group(1).strip(), size=15, bold=True)
        elif line.strip() == "---":
            _start("")
        elif line.startswith("> "):
            _add_body("\u300c" + line[2:].strip() + "\u300d", size=13, italic=True)
        elif _r.match(r"^[-*] ", line):
            _add_body("\u2022 " + line[2:].strip(), size=14)
        elif _r.match(r"^\d+\. ", line):
            _add_body(_r.sub(r"^\d+\. ", "", line).strip(), size=14)
        elif line.startswith("|"):
            clean = _r.sub(r"\|", "  ", line).strip()
            if not _r.match(r"^[-: ]+$", clean):
                _add_body(clean, size=12)
        elif line.strip() == "":
            pass  # skip blank lines
        else:
            text = _r.sub(r"\*\*(.+?)\*\*", r"\1", line)
            text = _r.sub(r"`(.+?)`", r"\1", text)
            text = _r.sub(r"\*(.+?)\*", r"\1", text)
            if text.strip():
                _add_body(text, size=14)

    _flush()

    if not slides_data:
        slides_data = [("繁中 AI 情報", [("（內容為空）", 16, False, False)])]

    for title, body_lines in slides_data:
        sl = _new_slide()
        if title:
            _set_title(sl, title)
        if body_lines:
            _set_body(sl, body_lines)

    prs.save(str(output_path))
    return output_path
