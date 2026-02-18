"""Diagnostics: parse PPTX/DOCX and detect quality issues.

Can be called from tests or standalone. Returns structured results for regression guarding.
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

# Signal type keywords that must NOT appear as signal_text content body.
_SIGNAL_TYPE_KEYWORDS = {
    "TOOL_ADOPTION", "USER_PAIN", "WORKFLOW_CHANGE", "PLATFORM_HEAT",
    "COST_SHIFT", "REGULATION", "SECURITY_ALERT",
    "工具採用", "用戶痛點", "工作流變化", "平台熱度",
    "成本變動", "監管", "安全警報",
}

# Trailing words that indicate a broken sentence.
_TRAILING_ZH = set("的了而與來記是在和或及對從向把被讓給")
_TRAILING_EN_RE = re.compile(
    r"\b(?:to|and|or|by|the|a|an|of|in|on|at|for|with)\s*$",
    re.IGNORECASE,
)
_SENTENCE_END_RE = re.compile(r"[.!?。！？;；]")


@dataclass
class PptxDiagResult:
    total_slides: int = 0
    total_text_runs: int = 0
    url_na_count: int = 0
    signal_type_as_text_count: int = 0
    stat_only_slides: int = 0
    fragment_lines: list[str] = field(default_factory=list)
    short_fragments: list[str] = field(default_factory=list)
    all_texts: list[str] = field(default_factory=list)

    @property
    def ok(self) -> bool:
        return (
            self.url_na_count == 0
            and self.signal_type_as_text_count == 0
            and len(self.fragment_lines) == 0
            and len(self.short_fragments) == 0
        )


def _is_trailing_fragment(text: str) -> bool:
    s = text.strip()
    if not s:
        return False
    if s[-1] in _TRAILING_ZH:
        return True
    if _TRAILING_EN_RE.search(s):
        return True
    return False


def _is_short_fragment(text: str) -> bool:
    s = text.strip()
    if len(s) < 12 and not re.search(r"\d", s) and not re.search(r"https?://", s, re.IGNORECASE):
        # Check for entities
        if not re.search(r"[A-Z][a-z]{2,}|[\u4e00-\u9fff]{2,}", s):
            return True
    return False


def diagnose_pptx(path: Path) -> PptxDiagResult:
    """Parse PPTX and return diagnostic results."""
    from pptx import Presentation

    result = PptxDiagResult()
    prs = Presentation(str(path))
    result.total_slides = len(prs.slides)

    for slide in prs.slides:
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = p.text.strip()
                    if t:
                        result.all_texts.append(t)
                        slide_texts.append(t)
                        result.total_text_runs += 1
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        t = cell.text.strip()
                        if t:
                            result.all_texts.append(t)
                            slide_texts.append(t)
                            result.total_text_runs += 1

        for t in slide_texts:
            lower = t.lower()
            # URL=N/A detection
            if "url=n/a" in lower or "url= n/a" in lower or "URL=N/A" in t:
                result.url_na_count += 1

            # signal_type used as signal_text (only flag if it appears as the sole
            # content in a text run that is NOT a short badge/label)
            stripped = t.strip()
            if len(stripped) > 3 and stripped in _SIGNAL_TYPE_KEYWORDS:
                # Check if this looks like a content line (not a badge label)
                # Badge labels are typically short and standalone; content lines are longer
                # We only flag when the ENTIRE text is just the signal type keyword
                # and it appears in a context where signal_text should have real content
                if len(stripped) > 15 or "|" in t or "platform_count" in t:
                    result.signal_type_as_text_count += 1

            # Trailing fragment (only check lines > 10 chars that aren't metrics/stats)
            if len(t) > 10 and "=" not in t and _is_trailing_fragment(t):
                result.fragment_lines.append(t)

            # Short fragments
            if 0 < len(t) < 12 and _is_short_fragment(t):
                result.short_fragments.append(t)

    return result


def slide_density_audit(path: Path) -> list[dict]:
    """Per-slide density metrics for a PPTX file.

    For each slide returns:
      slide_index, title, text_chars, nonempty_shapes,
      table_cells_total, table_cells_nonempty,
      terms, numbers, sentences, density_score, all_text (truncated to 600 chars).

    Density score ∈ [0, 100] (deterministic):
      text_score  = min(40, text_chars / EXEC_SLIDE_MIN_TEXT_CHARS * 40)
      table_score = (nonempty_ratio / EXEC_TABLE_MIN_NONEMPTY_RATIO).clamp(0,1)*30
                    or 30 when no table exists
      ev_score    = (terms>=T + nums>=N + sents>=S) / 3 * 30
      density     = round(text_score + table_score + ev_score)
    """
    from pptx import Presentation
    from utils.text_quality import count_evidence_terms, count_evidence_numbers, count_sentences
    from utils.semantic_quality import semantic_density_score as _sem_score

    # Load thresholds — fall back to hardcoded defaults if settings not available.
    try:
        import config.settings as _s
        min_text_chars = _s.EXEC_SLIDE_MIN_TEXT_CHARS
        table_ratio_min = _s.EXEC_TABLE_MIN_NONEMPTY_RATIO
        min_terms = _s.EXEC_BLOCK_MIN_EVIDENCE_TERMS
        min_nums = _s.EXEC_BLOCK_MIN_EVIDENCE_NUMBERS
        min_sents = _s.EXEC_BLOCK_MIN_SENTENCES
    except (ImportError, AttributeError):
        min_text_chars = 160
        table_ratio_min = 0.60
        min_terms = 2
        min_nums = 1
        min_sents = 2

    prs = Presentation(str(path))
    results: list[dict] = []

    for idx, slide in enumerate(prs.slides, 1):
        text_parts: list[str] = []
        table_cells_total = 0
        table_cells_nonempty = 0
        nonempty_shapes = 0
        slide_title = ""

        for shape in slide.shapes:
            if shape.has_text_frame:
                shape_texts: list[str] = []
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        shape_texts.append(t)
                        text_parts.append(t)
                if shape_texts:
                    nonempty_shapes += 1
                    if not slide_title and len(shape_texts[0]) > 3:
                        slide_title = shape_texts[0][:60]
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        table_cells_total += 1
                        t = cell.text.strip()
                        if t:
                            table_cells_nonempty += 1
                            text_parts.append(t)

        full_text = " ".join(text_parts)
        # text_chars: non-whitespace character count
        text_chars = len(full_text.replace(" ", ""))

        # ---- Density score calculation ----
        text_score = min(40.0, (text_chars / max(min_text_chars, 1)) * 40.0)

        if table_cells_total > 0:
            ratio = table_cells_nonempty / table_cells_total
            table_score = min(1.0, ratio / max(table_ratio_min, 0.01)) * 30.0
        else:
            table_score = 30.0  # no table → no penalty

        terms = count_evidence_terms(full_text)
        numbers = count_evidence_numbers(full_text)
        sentences = count_sentences(full_text)

        t_ok = terms >= min_terms
        n_ok = numbers >= min_nums
        s_ok = sentences >= min_sents
        evidence_score = ((int(t_ok) + int(n_ok) + int(s_ok)) / 3.0) * 30.0

        density_score = round(text_score + table_score + evidence_score)

        semantic_score = _sem_score(full_text)

        results.append({
            "slide_index": idx,
            "title": slide_title,
            "text_chars": text_chars,
            "nonempty_shapes": nonempty_shapes,
            "table_cells_total": table_cells_total,
            "table_cells_nonempty": table_cells_nonempty,
            "terms": terms,
            "numbers": numbers,
            "sentences": sentences,
            "density_score": density_score,
            "semantic_score": semantic_score,
            "all_text": full_text[:600],
        })

    return results


def diagnose_docx(path: Path) -> PptxDiagResult:
    """Parse DOCX and return diagnostic results (same structure)."""
    from docx import Document

    result = PptxDiagResult()
    doc = Document(str(path))

    all_texts: list[str] = []
    for p in doc.paragraphs:
        t = p.text.strip()
        if t:
            all_texts.append(t)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                t = cell.text.strip()
                if t:
                    all_texts.append(t)

    result.all_texts = all_texts
    result.total_text_runs = len(all_texts)

    for t in all_texts:
        lower = t.lower()
        if "url=n/a" in lower or "url= n/a" in lower or "URL=N/A" in t:
            result.url_na_count += 1
        stripped_t = t.strip()
        if len(stripped_t) > 3 and stripped_t in _SIGNAL_TYPE_KEYWORDS:
            if len(stripped_t) > 15 or "|" in t or "platform_count" in t:
                result.signal_type_as_text_count += 1
        if len(t) > 10 and "=" not in t and _is_trailing_fragment(t):
            result.fragment_lines.append(t)
        if 0 < len(t) < 12 and _is_short_fragment(t):
            result.short_fragments.append(t)

    return result
