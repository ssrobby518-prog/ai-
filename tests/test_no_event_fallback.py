from __future__ import annotations

from pathlib import Path
import re
from unittest.mock import patch

from pptx import Presentation

from core.content_strategy import build_corp_watch_summary, build_signal_summary
from core.ppt_generator import generate_executive_ppt
from schemas.education_models import EduNewsCard, SystemHealthReport


def _no_event_cards() -> list[EduNewsCard]:
    return [
        EduNewsCard(
            item_id="index-001",
            is_valid_news=True,
            title_plain="As part of its mission to preserve the web",
            what_happened="Curated links and archive index",
            why_important="Reference page, not an event",
            source_name="TechCrunch",
            source_url="https://example.com/index",
            final_score=2.0,
        ),
        EduNewsCard(
            item_id="invalid-001",
            is_valid_news=False,
            title_plain="Login page captured",
            what_happened="Please sign in to continue",
            invalid_reason="system banner",
            invalid_cause="blocked",
            source_name="HackerNews",
            source_url="https://example.com/login",
            final_score=0.0,
        ),
    ]


def _all_text(prs: Presentation) -> str:
    buf: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        buf.append(p.text.strip())
    return "\n".join(buf)


def test_no_event_signal_summary_top3() -> None:
    signals = build_signal_summary(_no_event_cards())
    assert len(signals) >= 3
    for sig in signals[:3]:
        assert "signal_name" in sig
        assert "signal_text" in sig
        assert "platform_count" in sig
        assert "heat_score" in sig
        assert "example_snippet" in sig
        assert len(sig["example_snippet"]) <= 120
        assert len(str(sig["example_snippet"]).strip()) >= 30
        assert int(sig["platform_count"]) >= 1
        assert int(sig["heat_score"]) >= 30
        assert str(sig.get("source_name", "")).strip()
        assert str(sig.get("source_name", "")).strip().lower() != "unknown"
        assert "evidence_tokens" in sig
        assert isinstance(sig["evidence_tokens"], list)
        assert len(sig["evidence_tokens"]) >= 2
        assert "fallback monitoring signal" not in str(sig["signal_text"]).lower()
        assert "fallback monitoring signal" not in str(sig["example_snippet"]).lower()
        assert "smoke" not in str(sig["signal_text"]).lower()
        assert "smoke" not in str(sig["example_snippet"]).lower()
        assert "source=unknown" not in str(sig["example_snippet"]).lower()
        assert re.search(r"[\u4e00-\u9fff]", str(sig.get("signal_text", "")))


def test_no_event_corp_watch_includes_scan_stats() -> None:
    corp = build_corp_watch_summary(_no_event_cards())
    assert corp.get("updates", -1) == 0
    assert corp.get("mentions_count", -1) == 0
    assert corp.get("trend_direction") == "STABLE"
    status = str(corp.get("status_message", ""))
    assert "掃描統計" in status
    assert "sources_total=" in status
    assert "success_count=" in status
    assert "fail_count=" in status
    assert "sources_total" in corp
    assert "success_count" in corp
    assert "fail_count" in corp
    assert "top_fail_reasons" in corp
    assert "top_sources" in corp
    assert corp["sources_total"] >= 1
    assert corp["top_fail_reasons"]
    assert corp["top_sources"]
    assert str(corp["top_fail_reasons"][0].get("reason", "")).strip()


def test_no_event_still_generates_complete_deck(tmp_path: Path) -> None:
    out = tmp_path / "no_event.pptx"
    health = SystemHealthReport(success_rate=0.0, p50_latency=0.0, p95_latency=0.0)
    cards = _no_event_cards()

    with patch("core.ppt_generator.get_news_image", return_value=None):
        generate_executive_ppt(
            cards=cards,
            health=health,
            report_time="2026-02-15 09:00",
            total_items=len(cards),
            output_path=out,
            metrics={
                "fetched_total": 14,
                "gate_pass_total": 5,
                "after_filter_total": 5,
                "sources_total": 3,
            },
        )

    prs = Presentation(str(out))
    assert len(prs.slides) == 12

    text = _all_text(prs)
    assert "Signal Thermometer" in text
    assert "Corp Watch" in text
    assert "Structured Summary" in text
    assert "Key Takeaways" in text
    assert "Overview" in text
    assert "Event Ranking" in text
    assert "Recommended Moves" in text
    assert "Decision Matrix" in text
    assert "待決事項與 Owner" in text
    assert "sources_total" in text
    assert "fetched_total" in text
    assert "gate_pass_total" in text
    assert "sources_total" in text.lower()
    assert "success_count" in text.lower()
    assert "fallback monitoring signal" not in text.lower()
    assert "desktop smoke signal" not in text.lower()
    assert "signals_insufficient=true" not in text.lower()
    assert "source=unknown" not in text.lower()


def test_empty_passed_signals_mark_insufficient() -> None:
    signals = build_signal_summary([])
    assert len(signals) >= 3
    for sig in signals[:3]:
        assert sig.get("signals_insufficient") is True
        assert int(sig.get("passed_total_count", -1)) == 0
        assert str(sig.get("signal_text", "")).strip()


def test_no_event_summary_is_stats(tmp_path: Path) -> None:
    cards = _no_event_cards()
    with patch("core.ppt_generator.get_news_image", return_value=None):
        out = generate_executive_ppt(
            cards=cards,
            health=SystemHealthReport(success_rate=0.0, p50_latency=0.0, p95_latency=0.0),
            report_time="2026-02-15 09:00",
            total_items=len(cards),
            output_path=tmp_path / "no_event_summary_stats.pptx",
            metrics={
                "fetched_total": 20,
                "gate_pass_total": 7,
                "after_filter_total": 7,
                "sources_total": 4,
            },
        )
    prs = Presentation(str(out))
    text = _all_text(prs).lower()
    assert "fetched_total" in text
    assert "gate_pass_total" in text
    assert "sources_total" in text
