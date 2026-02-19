"""Executive PPT/DOCX semantic guard tests.

Validates that hollow content (fragments, placeholders, template remnants)
cannot reach the rendered PPT or DOCX output.

Test classes:
  TestPlaceholderFragmentDetection      — is_placeholder_or_fragment() unit tests
  TestSemanticGuardBackfill             — semantic_guard_text() backfill behaviour
  TestPptDocNoAirInKeySlides            — integration: no air in Overview/Ranking/Pending
  TestVerifyRunDensityAuditOutput       — contract: verify_run.ps1 prints expected lines
"""

from __future__ import annotations

import re
from pathlib import Path
from unittest.mock import patch

import pytest

import config.settings as settings
from core.content_strategy import semantic_guard_text
from schemas.education_models import EduNewsCard, SystemHealthReport
from utils.semantic_quality import (
    count_evidence_numbers,
    count_evidence_terms,
    count_sentences,
    is_placeholder_or_fragment,
    semantic_density_score,
)


# ---------------------------------------------------------------------------
# Shared fixtures
# ---------------------------------------------------------------------------

def _make_card(
    title: str = "OpenAI releases GPT-5",
    what: str = "OpenAI announced GPT-5 with 1M context window for enterprise use.",
    why: str = "Enables new LLM use cases; impacts AI training cost globally.",
    source: str = "TechCrunch",
    url: str = "https://techcrunch.com/gpt5",
    score: float = 8.5,
) -> EduNewsCard:
    return EduNewsCard(
        item_id="guard-test-001",
        is_valid_news=True,
        title_plain=title,
        what_happened=what,
        why_important=why,
        source_name=source,
        source_url=url,
        final_score=score,
        category="AI",
    )


def _health() -> SystemHealthReport:
    return SystemHealthReport(success_rate=80.0, p50_latency=1.0, p95_latency=3.0)


def _rich_cards() -> list[EduNewsCard]:
    return [
        EduNewsCard(
            item_id=f"air-{i:03d}",
            is_valid_news=True,
            title_plain=title,
            what_happened=what,
            why_important=why,
            source_name=src,
            source_url=url,
            final_score=score,
            category="AI",
        )
        for i, (title, what, why, src, url, score) in enumerate([
            (
                "NVIDIA H200 GPU launched at $30k",
                "NVIDIA launched the H200 GPU with 141GB HBM3e memory for $30k per unit.",
                "Impacts AI training costs globally for LLM model development.",
                "TechCrunch",
                "https://techcrunch.com/nvidia-h200",
                9.0,
            ),
            (
                "OpenAI releases GPT-5 with 1M context",
                "OpenAI announced GPT-5 with 1 million token context window for enterprise RAG.",
                "Enables new LLM use cases with large AI model context windows.",
                "TheVerge",
                "https://theverge.com/openai-gpt5",
                8.5,
            ),
            (
                "Anthropic Claude reaches 100M users",
                "Anthropic reported Claude AI model reached 100 million monthly active users.",
                "Claude becomes a major AI competitor challenging ChatGPT and GPT models.",
                "Reuters",
                "https://reuters.com/anthropic-claude-100m",
                8.0,
            ),
        ])
    ]


def _hollow_cards() -> list[EduNewsCard]:
    """Cards with hollow / placeholder field values — the pipeline stress test."""
    return [
        EduNewsCard(
            item_id=f"hollow-{i:03d}",
            is_valid_news=True,
            title_plain=title,
            what_happened=what,
            why_important=why,
            source_name=src,
            source_url=url,
            final_score=score,
            category="AI",
        )
        for i, (title, what, why, src, url, score) in enumerate([
            (
                "的趨勢，解決方 記",          # hollow title — truncation artifact
                "的趨勢",                      # hollow what
                "",                            # empty why
                "HollowSource",
                "https://example.com/hollow1",
                6.0,
            ),
            (
                "Last July was a turning point",  # template remnant title
                "Last July was an interesting time",
                "",
                "TplSource",
                "https://example.com/hollow2",
                5.5,
            ),
            (
                "OpenAI GPT-5 Released",           # valid title (needed for fallback)
                "",                                # empty what
                "",                                # empty why
                "ValidSource",
                "https://example.com/hollow3",
                7.0,
            ),
        ])
    ]


def _gen_pptx(tmp_path: Path, cards: list[EduNewsCard]) -> Path:
    from core.ppt_generator import generate_executive_ppt

    pptx_path = tmp_path / "test_air.pptx"
    with patch("core.ppt_generator.get_news_image", return_value=None):
        generate_executive_ppt(
            cards=cards,
            health=_health(),
            report_time="2026-01-01 09:00",
            total_items=len(cards),
            output_path=pptx_path,
        )
    return pptx_path


def _gen_docx(tmp_path: Path, cards: list[EduNewsCard]) -> Path:
    from core.doc_generator import generate_executive_docx

    docx_path = tmp_path / "test_air.docx"
    generate_executive_docx(
        cards=cards,
        health=_health(),
        report_time="2026-01-01 09:00",
        total_items=len(cards),
        output_path=docx_path,
    )
    return docx_path


# ---------------------------------------------------------------------------
# 1) Placeholder / fragment detection unit tests
# ---------------------------------------------------------------------------

class TestPlaceholderFragmentDetection:
    """is_placeholder_or_fragment() must catch all hollow text cases."""

    def test_empty_string_is_fragment(self) -> None:
        assert is_placeholder_or_fragment("") is True

    def test_whitespace_only_is_fragment(self) -> None:
        assert is_placeholder_or_fragment("   ") is True

    def test_lone_sequence_number_is_fragment(self) -> None:
        assert is_placeholder_or_fragment("2. ") is True
        assert is_placeholder_or_fragment("3)") is True
        assert is_placeholder_or_fragment("10.") is True

    def test_lone_bullet_is_fragment(self) -> None:
        assert is_placeholder_or_fragment("•") is True
        assert is_placeholder_or_fragment("—") is True
        assert is_placeholder_or_fragment("→") is True

    def test_last_july_was_is_fragment(self) -> None:
        assert is_placeholder_or_fragment("Last July was a turning point") is True

    def test_truncation_artifact_is_fragment(self) -> None:
        assert is_placeholder_or_fragment("的趨勢，解決方 記") is True

    def test_dangling_zh_particle_phrase_is_fragment(self) -> None:
        # "的趨勢" starts with ZH particle, short, no sentence end
        assert is_placeholder_or_fragment("的趨勢") is True

    def test_trailing_connector_is_fragment(self) -> None:
        assert is_placeholder_or_fragment("something but") is True
        assert is_placeholder_or_fragment("analysis of,") is True

    def test_real_sentence_is_not_fragment(self) -> None:
        assert is_placeholder_or_fragment(
            "NVIDIA launched H200 GPU with 141GB HBM3e for $30k."
        ) is False

    def test_chinese_sentence_is_not_fragment(self) -> None:
        assert is_placeholder_or_fragment(
            "OpenAI 發布 GPT-5，支援 100 萬 token 上下文，對企業 RAG 場景影響深遠。"
        ) is False

    def test_number_in_short_text_not_fragment(self) -> None:
        # Short but has a number → not a fragment
        assert is_placeholder_or_fragment("v3.5") is False

    def test_impact_ratio_not_fragment(self) -> None:
        assert is_placeholder_or_fragment("impact=4/5") is False

    def test_settings_placeholder_patterns_are_detectable(self) -> None:
        """All PLACEHOLDER_PATTERNS from settings must be caught."""
        for pat in settings.PLACEHOLDER_PATTERNS:
            # Build a minimal string matching the pattern
            if pat == r"^\s*[0-9]+[.)]\s*$":
                sample = "2."
            elif r"Last" in pat:
                sample = "Last July was a turning point"
            elif "解決方" in pat:
                sample = "解決方 記"
            elif "WHY IT MATTERS" in pat:
                sample = "WHY IT MATTERS:"
            else:
                continue
            assert re.search(pat, sample), (
                f"Pattern {pat!r} didn't match sample {sample!r}"
            )


# ---------------------------------------------------------------------------
# 2) Semantic guard backfill tests
# ---------------------------------------------------------------------------

class TestSemanticGuardBackfill:
    """semantic_guard_text() must always return substantial, readable content."""

    def test_good_text_passes_through(self) -> None:
        card = _make_card()
        good = "NVIDIA launched H200 GPU at $30k for AI training workloads."
        result = semantic_guard_text(good, card)
        assert result == good

    def test_fragment_triggers_backfill(self) -> None:
        card = _make_card()
        result = semantic_guard_text("的趨勢，解決方 記", card)
        assert not is_placeholder_or_fragment(result), (
            f"Guard returned fragment: {result!r}"
        )

    def test_empty_triggers_backfill(self) -> None:
        card = _make_card()
        result = semantic_guard_text("", card)
        assert result  # non-empty
        assert not is_placeholder_or_fragment(result)

    def test_lone_number_triggers_backfill(self) -> None:
        card = _make_card()
        result = semantic_guard_text("2. ", card)
        assert not is_placeholder_or_fragment(result)

    def test_backfill_has_min_sentences(self) -> None:
        card = _make_card(
            what="OpenAI launched GPT-5 with 1 million token context for enterprise.",
            why="Impacts AI training globally for LLM development.",
        )
        result = semantic_guard_text("的趨勢", card, context="action")
        assert count_sentences(result) >= 1, (
            f"Backfill has no sentence end: {result!r}"
        )

    def test_backfill_has_min_terms(self) -> None:
        card = _make_card(
            what="OpenAI launched GPT-5 with 1 million token context for enterprise.",
            why="Impacts AI training globally for LLM development.",
        )
        result = semantic_guard_text("Last July was interesting", card)
        assert count_evidence_terms(result) >= 1, (
            f"Backfill has no evidence terms: {result!r}"
        )

    def test_backfill_has_min_numbers(self) -> None:
        card = _make_card(
            what="NVIDIA H200 launched at $30k with 141GB HBM3e memory.",
        )
        result = semantic_guard_text("", card)
        assert count_evidence_numbers(result) >= 1, (
            f"Backfill has no numbers: {result!r}"
        )

    def test_hollow_card_uses_structured_fallback(self) -> None:
        """When all card fields are hollow, fallback must still be a full sentence."""
        hollow = EduNewsCard(
            item_id="hollow-001",
            is_valid_news=True,
            title_plain="Hollow card title",
            what_happened="",
            why_important="",
            source_name="TestSource",
            source_url="https://example.com",
            final_score=5.0,
            category="AI",
        )
        result = semantic_guard_text("的趨勢", hollow)
        assert "Hollow card title" in result or "TestSource" in result, (
            f"Fallback missing title/source: {result!r}"
        )
        assert not is_placeholder_or_fragment(result)

    def test_semantic_density_score_above_threshold_after_backfill(self) -> None:
        card = _make_card()
        result = semantic_guard_text("Last July was", card)
        score = semantic_density_score(result)
        assert score >= 70, f"Backfill semantic score too low ({score}): {result!r}"

    def test_hollow_card_fallback_has_number(self) -> None:
        """Guaranteed fallback must contain a number so density scoring finds one."""
        hollow = EduNewsCard(
            item_id="hollow-002",
            is_valid_news=True,
            title_plain="的趨勢，解決方 記",
            what_happened="的趨勢",
            why_important="",
            source_name="TestSource",
            source_url="https://example.com",
            final_score=7.0,
            category="AI",
        )
        result = semantic_guard_text("的趨勢", hollow)
        assert count_evidence_numbers(result) >= 1, (
            f"Fallback has no numeric evidence: {result!r}"
        )

    def test_hollow_card_fallback_not_fragment(self) -> None:
        """When title is hollow, fallback must still not be a fragment."""
        hollow = EduNewsCard(
            item_id="hollow-003",
            is_valid_news=True,
            title_plain="的趨勢，解決方 記",
            what_happened="",
            why_important="",
            source_name="TestSource",
            source_url="https://example.com",
            final_score=5.0,
            category="AI",
        )
        result = semantic_guard_text("", hollow)
        assert not is_placeholder_or_fragment(result), (
            f"Fallback with hollow title is still a fragment: {result!r}"
        )
        assert count_sentences(result) >= 1, (
            f"Fallback has no sentence: {result!r}"
        )


# ---------------------------------------------------------------------------
# 3) Integration: no air in key slides (PPT + DOCX)
# ---------------------------------------------------------------------------

class TestPptDocNoAirInKeySlides:
    """PPT / DOCX must not emit hollow content in Overview / Ranking / Pending slides."""

    # Placeholder patterns to detect (as raw strings)
    _BAD_PATTERNS = [
        r"Last\s+\w+\s+was\b",
        r"解決方\s*[記表]",
        r"^\s*[0-9]+[.)]\s*$",      # lone number like "2."
        r"的趨勢，解決方",
    ]

    def _extract_pptx_texts(self, pptx_path: Path) -> list[str]:
        from pptx import Presentation

        prs = Presentation(str(pptx_path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            t = cell.text.strip()
                            if t:
                                texts.append(t)
        return texts

    def _extract_docx_texts(self, docx_path: Path) -> list[str]:
        from docx import Document

        doc = Document(str(docx_path))
        texts = []
        for p in doc.paragraphs:
            t = p.text.strip()
            if t:
                texts.append(t)
        for tbl in doc.tables:
            for row in tbl.rows:
                for cell in row.cells:
                    t = cell.text.strip()
                    if t:
                        texts.append(t)
        return texts

    def _assert_no_bad_patterns(self, texts: list[str], source: str) -> None:
        for text in texts:
            for pat in self._BAD_PATTERNS:
                if re.search(pat, text, re.IGNORECASE | re.MULTILINE):
                    pytest.fail(
                        f"Placeholder pattern {pat!r} found in {source}: {text!r}"
                    )

    def _assert_no_forbidden_fragments(self, texts: list[str], source: str) -> None:
        for frag in settings.EXEC_FORBIDDEN_FRAGMENTS:
            for text in texts:
                assert frag not in text, (
                    f"Forbidden fragment {frag!r} found in {source}: {text!r}"
                )

    def test_pptx_no_placeholder_patterns(self, tmp_path: Path) -> None:
        pptx = _gen_pptx(tmp_path, _rich_cards())
        texts = self._extract_pptx_texts(pptx)
        assert texts, "PPTX produced no text at all"
        self._assert_no_bad_patterns(texts, "PPTX")

    def test_pptx_no_forbidden_fragments(self, tmp_path: Path) -> None:
        pptx = _gen_pptx(tmp_path, _rich_cards())
        texts = self._extract_pptx_texts(pptx)
        self._assert_no_forbidden_fragments(texts, "PPTX")

    def test_docx_no_placeholder_patterns(self, tmp_path: Path) -> None:
        docx = _gen_docx(tmp_path, _rich_cards())
        texts = self._extract_docx_texts(docx)
        assert texts, "DOCX produced no text at all"
        self._assert_no_bad_patterns(texts, "DOCX")

    def test_docx_no_forbidden_fragments(self, tmp_path: Path) -> None:
        docx = _gen_docx(tmp_path, _rich_cards())
        texts = self._extract_docx_texts(docx)
        self._assert_no_forbidden_fragments(texts, "DOCX")

    def test_pptx_key_slides_density(self, tmp_path: Path) -> None:
        """Overview / Ranking / Pending slides must pass semantic density gate."""
        from scripts.diagnostics_pptx import slide_density_audit

        pptx = _gen_pptx(tmp_path, _rich_cards())
        results = slide_density_audit(pptx)

        key_patterns = [
            ("overview", ["overview", "總覽"]),
            ("ranking", ["event ranking", "排行"]),
            ("pending", ["pending", "待決"]),
        ]
        for slide_key, search_patterns in key_patterns:
            slide = next(
                (r for r in results
                 if any(p in r["title"].lower() for p in search_patterns)),
                None,
            )
            assert slide is not None, f"Key slide '{slide_key}' not found"
            threshold = settings.EXEC_DENSITY_THRESHOLDS.get(slide_key, 80)
            assert slide["density_score"] >= threshold, (
                f"{slide_key} density={slide['density_score']} < {threshold}"
            )
            sem_threshold = settings.EXEC_SEMANTIC_THRESHOLDS.get(slide_key, 40)
            assert slide.get("semantic_score", 0) >= sem_threshold, (
                f"{slide_key} sem_score={slide.get('semantic_score', 0)} < {sem_threshold}"
            )

    def test_pptx_table_nonempty_ratio(self, tmp_path: Path) -> None:
        """Table cells in key slides must meet PER_CELL_MIN_NONEMPTY_RATIO."""
        from scripts.diagnostics_pptx import slide_density_audit

        pptx = _gen_pptx(tmp_path, _rich_cards())
        results = slide_density_audit(pptx)

        for slide in results:
            if any(p in slide["title"].lower() for p in ["overview", "總覽", "event ranking", "排行"]):
                total = slide["table_cells_total"]
                if total > 0:
                    ratio = slide["table_cells_nonempty"] / total
                    assert ratio >= settings.EXEC_TABLE_MIN_NONEMPTY_RATIO, (
                        f"{slide['title']!r} table ratio {ratio:.2%} < "
                        f"{settings.EXEC_TABLE_MIN_NONEMPTY_RATIO:.2%}"
                    )


# ---------------------------------------------------------------------------
# 4) Contract: verify_run.ps1 prints expected density audit lines
# ---------------------------------------------------------------------------

class TestVerifyRunDensityAuditOutput:
    """verify_run.ps1 must print [DENSITY] lines and density audit result."""

    def _read_script(self) -> str:
        return Path("scripts/verify_run.ps1").read_text(encoding="utf-8")

    def test_density_format_string_present(self) -> None:
        text = self._read_script()
        assert "[DENSITY] slide=" in text, (
            "verify_run.ps1 missing '[DENSITY] slide=' format string"
        )

    def test_density_audit_label_present(self) -> None:
        text = self._read_script()
        assert "Executive Slide Density Audit" in text

    def test_density_audit_passed_label_present(self) -> None:
        text = self._read_script()
        assert "Executive Slide Density Audit PASSED" in text

    def test_no_ansi_escape_in_density_lines(self) -> None:
        """[DENSITY] lines must not contain ANSI escape sequences."""
        text = self._read_script()
        # Extract lines that contain [DENSITY] format strings
        for line in text.splitlines():
            if "[DENSITY]" in line:
                assert "\x1b" not in line and "\\x1b" not in line, (
                    f"ANSI escape in [DENSITY] line: {line!r}"
                )

    def test_no_text_chars_skip_guard_in_density_gate(self) -> None:
        """No text_chars-based skip guard in density gate section."""
        text = self._read_script()
        # The gate used to have: if ($isKey -and $s.text_chars -ge 60)
        # Build pattern at runtime to avoid self-match
        guard = "text_chars" + " -ge"
        assert guard not in text, (
            "text_chars skip guard detected in verify_run.ps1 density gate"
        )

    def test_semantic_score_in_density_output(self) -> None:
        """verify_run.ps1 must output sem_score= in density lines."""
        text = self._read_script()
        assert "sem_score=" in text, (
            "verify_run.ps1 missing 'sem_score=' in density audit output"
        )

    def test_semantic_gate_is_primary_hard_gate(self) -> None:
        """verify_run.ps1: semantic gate is hard (FAIL), formal density is warn-only."""
        text = self._read_script()
        assert "requiredSemanticDensity" in text
        # Semantic FAIL must be present
        assert "DENSITY FAIL" in text
        # Formal density downgraded to WARN (not hard FAIL) when semantic is OK
        assert "DENSITY WARN" in text

    def test_notebase64_not_present(self) -> None:
        """Base64-encoded note must be removed."""
        text = self._read_script()
        assert "noteBase64" not in text


# ---------------------------------------------------------------------------
# 5) Hollow-card regression: placeholder text must not reach PPT or DOCX
# ---------------------------------------------------------------------------

class TestHollowCardNoAir:
    """PPT / DOCX rendered from hollow/template-remnant cards must not leak placeholders."""

    _BAD_PATTERNS = [
        r"Last\s+\w+\s+was\b",
        r"解決方\s*[記表]",
        r"^\s*[0-9]+[.)]\s*$",
        r"的趨勢，解決方",
    ]

    def _extract_pptx_texts(self, pptx_path: Path) -> list[str]:
        from pptx import Presentation

        prs = Presentation(str(pptx_path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            t = cell.text.strip()
                            if t:
                                texts.append(t)
        return texts

    def _extract_docx_texts(self, docx_path: Path) -> list[str]:
        from docx import Document

        doc = Document(str(docx_path))
        texts = []
        for p in doc.paragraphs:
            t = p.text.strip()
            if t:
                texts.append(t)
        for tbl in doc.tables:
            for row in tbl.rows:
                for cell in row.cells:
                    t = cell.text.strip()
                    if t:
                        texts.append(t)
        return texts

    def _assert_no_bad_patterns(self, texts: list[str], source: str) -> None:
        for text in texts:
            for pat in self._BAD_PATTERNS:
                if re.search(pat, text, re.IGNORECASE | re.MULTILINE):
                    pytest.fail(
                        f"Placeholder pattern {pat!r} leaked in {source}: {text!r}"
                    )

    def test_placeholder_fragment_examples_rejected(self) -> None:
        """All hollow card field values must be flagged as placeholders/fragments."""
        bad_examples = [
            "的趨勢，解決方 記",
            "Last July was a turning point",
            "的趨勢",
            "",
            "2. ",
            "WHY IT MATTERS:",
        ]
        for example in bad_examples:
            assert is_placeholder_or_fragment(example), (
                f"Expected is_placeholder_or_fragment({example!r}) == True"
            )

    def test_semantic_guard_backfills_to_full_sentence(self) -> None:
        """semantic_guard_text with hollow card must backfill to sentence with evidence."""
        hollow_card = EduNewsCard(
            item_id="hollow-regr-001",
            is_valid_news=True,
            title_plain="的趨勢，解決方 記",
            what_happened="的趨勢",
            why_important="",
            source_name="HollowSrc",
            source_url="https://example.com/hollow",
            final_score=6.5,
            category="AI",
        )
        result = semantic_guard_text("的趨勢，解決方 記", hollow_card)
        assert not is_placeholder_or_fragment(result), (
            f"Guard returned fragment: {result!r}"
        )
        assert count_evidence_terms(result) >= 2 or count_evidence_numbers(result) >= 1, (
            f"Backfill has insufficient evidence: terms={count_evidence_terms(result)}, "
            f"nums={count_evidence_numbers(result)}, text={result!r}"
        )
        assert count_sentences(result) >= 1, (
            f"Backfill has no sentence boundary: {result!r}"
        )

    def test_pptx_no_placeholder_with_hollow_cards(self, tmp_path: Path) -> None:
        """PPT generated from hollow cards must not contain placeholder patterns."""
        pptx = _gen_pptx(tmp_path, _hollow_cards())
        texts = self._extract_pptx_texts(pptx)
        assert texts, "PPTX with hollow cards produced no text at all"
        self._assert_no_bad_patterns(texts, "PPTX[hollow]")

    def test_docx_no_placeholder_with_hollow_cards(self, tmp_path: Path) -> None:
        """DOCX generated from hollow cards must not contain placeholder patterns."""
        docx = _gen_docx(tmp_path, _hollow_cards())
        texts = self._extract_docx_texts(docx)
        assert texts, "DOCX with hollow cards produced no text at all"
        self._assert_no_bad_patterns(texts, "DOCX[hollow]")
