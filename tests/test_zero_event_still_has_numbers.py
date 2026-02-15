from __future__ import annotations

from pathlib import Path
from unittest.mock import patch

from docx import Document
from pptx import Presentation

from core.doc_generator import generate_executive_docx
from core.ppt_generator import generate_executive_ppt
from schemas.education_models import EduNewsCard, SystemHealthReport


def _cards() -> list[EduNewsCard]:
    return [
        EduNewsCard(
            item_id="zero-001",
            is_valid_news=True,
            title_plain="Archive overview page",
            what_happened="Curated links and historical index only",
            why_important="No single event signal available",
            source_name="Reddit",
            source_url="https://example.com/reddit",
            final_score=2.3,
        ),
        EduNewsCard(
            item_id="zero-002",
            is_valid_news=False,
            title_plain="Sign in required",
            what_happened="Please login to continue",
            invalid_reason="login",
            invalid_cause="blocked",
            source_name="XHS",
            source_url="https://example.com/xhs",
            final_score=0.0,
        ),
    ]


def _ppt_text(path: Path) -> str:
    prs = Presentation(str(path))
    out: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        out.append(p.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            out.append(cell.text.strip())
    return "\n".join(out)


def _doc_text(path: Path) -> str:
    doc = Document(str(path))
    out: list[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            out.append(p.text.strip())
    for t in doc.tables:
        for row in t.rows:
            for cell in row.cells:
                if cell.text.strip():
                    out.append(cell.text.strip())
    return "\n".join(out)


def test_zero_event_still_has_numeric_signal_and_corp_stats(tmp_path: Path) -> None:
    cards = _cards()
    health = SystemHealthReport(success_rate=0.0, p50_latency=0.0, p95_latency=0.0)
    pptx_path = tmp_path / "zero_event.pptx"
    docx_path = tmp_path / "zero_event.docx"
    metrics = {
        "fetched_total": 11,
        "gate_pass_total": 3,
        "after_filter_total": 3,
        "sources_total": 2,
        "sources_success": 1,
        "sources_failed": 1,
    }

    with patch("core.ppt_generator.get_news_image", return_value=None), patch(
        "core.doc_generator.get_news_image", return_value=None
    ):
        generate_executive_ppt(
            cards=cards,
            health=health,
            report_time="2026-02-15 09:00",
            total_items=len(cards),
            output_path=pptx_path,
            metrics=metrics,
        )
        generate_executive_docx(
            cards=cards,
            health=health,
            report_time="2026-02-15 09:00",
            total_items=len(cards),
            output_path=docx_path,
            metrics=metrics,
        )

    ppt_text = _ppt_text(pptx_path).lower()
    doc_text = _doc_text(docx_path).lower()
    merged = ppt_text + "\n" + doc_text

    assert "heat_score=" in merged
    assert "platform_count=" in merged
    assert "sources_total" in merged
    assert "success_count" in merged
    assert "fail_count" in merged
    assert "top_fail_reasons" in merged
