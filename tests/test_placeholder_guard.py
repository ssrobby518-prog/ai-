from __future__ import annotations

from pathlib import Path
import re
from unittest.mock import patch

from docx import Document
from pptx import Presentation

from core.doc_generator import generate_executive_docx
from core.ppt_generator import generate_executive_ppt
from schemas.education_models import EduNewsCard, SystemHealthReport


def _cards() -> list[EduNewsCard]:
    return [
        EduNewsCard(
            item_id="no-event-001",
            is_valid_news=True,
            title_plain="Archive index update",
            what_happened="Last July was...",
            why_important="Reference-only page without concrete event.",
            source_name="ExampleSource",
            source_url="https://example.com/archive",
            final_score=1.0,
        ),
        EduNewsCard(
            item_id="invalid-001",
            is_valid_news=False,
            title_plain="Sign in",
            what_happened="Please sign in to continue",
            invalid_reason="blocked",
            invalid_cause="login",
            source_name="ExampleSource2",
            source_url="https://example.com/login",
            final_score=0.0,
        ),
    ]


def _extract_ppt_text(path: Path) -> str:
    prs = Presentation(str(path))
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            texts.append(cell.text.strip())
    return "\n".join(texts)


def _extract_doc_text(path: Path) -> str:
    doc = Document(str(path))
    texts: list[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            texts.append(p.text.strip())
    for t in doc.tables:
        for row in t.rows:
            for cell in row.cells:
                if cell.text.strip():
                    texts.append(cell.text.strip())
    return "\n".join(texts)


def test_pptx_and_docx_have_no_placeholder_terms(tmp_path: Path) -> None:
    cards = _cards()
    health = SystemHealthReport(success_rate=50.0, p50_latency=1.0, p95_latency=2.0)
    pptx_path = tmp_path / "guard.pptx"
    docx_path = tmp_path / "guard.docx"

    with patch("core.ppt_generator.get_news_image", return_value=None), patch(
        "core.doc_generator.get_news_image",
        return_value=None,
    ):
        generate_executive_ppt(
            cards=cards,
            health=health,
            report_time="2026-02-15 09:00",
            total_items=len(cards),
            output_path=pptx_path,
        )
        generate_executive_docx(
            cards=cards,
            health=health,
            report_time="2026-02-15 09:00",
            total_items=len(cards),
            output_path=docx_path,
        )

    ppt_text = _extract_ppt_text(pptx_path).lower()
    doc_text = _extract_doc_text(docx_path).lower()

    banned_phrases = [
        "fallback monitoring signal",
        "desktop smoke signal",
        "signals_insufficient=true",
        "source=unknown",
        "last july was",
        "stay tuned",
        "this shows",
        "was...",
        "is...",
    ]
    for phrase in banned_phrases:
        assert phrase not in ppt_text
        assert phrase not in doc_text

    fragment_patterns = [
        re.compile(r"last\\s+\\w+\\s+was", re.IGNORECASE),
        re.compile(r"this\\s+\\w+\\s+was", re.IGNORECASE),
        re.compile(r"\\b(?:was|is|are)\\s*\\.\\.\\.", re.IGNORECASE),
        re.compile(r"\\b(?:this|that|it)\\s+\\w+\\s+(?:was|is|are)\\s*,\\s*$", re.IGNORECASE),
    ]
    for pattern in fragment_patterns:
        assert pattern.search(ppt_text) is None
        assert pattern.search(doc_text) is None

    # Density floor: enforce minimum numeric and proper-noun presence.
    combined = f"{ppt_text}\n{doc_text}"
    numeric_count = len(re.findall(r"\b\d+(?:\.\d+)?\b", combined))
    proper_noun_count = len(set(re.findall(r"\b[A-Za-z0-9][A-Za-z0-9\-_]{2,}\b", combined)))
    assert numeric_count >= 8
    assert proper_noun_count >= 10
