from __future__ import annotations

import re
from pathlib import Path
from unittest.mock import patch

from docx import Document
from pptx import Presentation

from core.doc_generator import generate_executive_docx
from core.ppt_generator import generate_executive_ppt
from schemas.education_models import EduNewsCard, SystemHealthReport


def _cards() -> list[EduNewsCard]:
    # Intentionally no strong event anchors -> should still produce non-air fallback content.
    return [
        EduNewsCard(
            item_id="no-event-01",
            is_valid_news=True,
            title_plain="As part of its mission to preserve the web",
            what_happened="Curated links and archive index.",
            why_important="Reference page, not a single verifiable event.",
            source_name="TechCrunch",
            source_url="https://example.com/index",
            final_score=2.0,
        ),
        EduNewsCard(
            item_id="no-event-02",
            is_valid_news=True,
            title_plain="Weekly AI roundup links",
            what_happened="Weekly roundup index with links only.",
            why_important="Useful as coverage signal, not event card.",
            source_name="HackerNews",
            source_url="https://example.com/roundup",
            final_score=2.2,
        ),
    ]


def _extract_ppt_text(path: Path) -> str:
    prs = Presentation(str(path))
    chunks: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        chunks.append(p.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            chunks.append(cell.text.strip())
    return "\n".join(chunks)


def _extract_doc_text(path: Path) -> str:
    doc = Document(str(path))
    chunks: list[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            chunks.append(p.text.strip())
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    chunks.append(cell.text.strip())
    return "\n".join(chunks)


def test_report_not_air_density_floor(tmp_path: Path) -> None:
    cards = _cards()
    health = SystemHealthReport(success_rate=0.0, p50_latency=0.0, p95_latency=0.0)
    pptx_path = tmp_path / "not_air.pptx"
    docx_path = tmp_path / "not_air.docx"
    metrics = {
        "fetched_total": 157,
        "event_gate_pass_total": 0,
        "signal_gate_pass_total": 9,
        "gate_pass_total": 4,
        "hard_pass_total": 4,
        "soft_pass_total": 5,
        "gate_reject_total": 22,
        "sources_total": 6,
        "sources_success": 4,
        "sources_failed": 2,
        "rejected_reason_top": [("content_too_short", 8), ("rejected_keyword:index", 4)],
    }

    with patch("core.ppt_generator.get_news_image", return_value=None), patch(
        "core.doc_generator.get_news_image",
        return_value=None,
    ):
        generate_executive_ppt(
            cards=cards,
            health=health,
            report_time="2026-02-16 09:00",
            total_items=len(cards),
            output_path=pptx_path,
            metrics=metrics,
        )
        generate_executive_docx(
            cards=cards,
            health=health,
            report_time="2026-02-16 09:00",
            total_items=len(cards),
            output_path=docx_path,
            metrics=metrics,
        )

    ppt_text = _extract_ppt_text(pptx_path)
    doc_text = _extract_doc_text(docx_path)
    merged = f"{ppt_text}\n{doc_text}"

    # Non-air floor.
    assert len(ppt_text) >= 2000
    assert len(re.findall(r"\d", merged)) >= 12
    assert len(set(re.findall(r"\b[A-Za-z][A-Za-z0-9\-_]{2,}\b", merged))) >= 12

    # Placeholder + fragment guard.
    lowered = merged.lower()
    banned = [
        "last july was",
        "desktop smoke signal",
        "fallback monitoring signal",
        "signals_insufficient=true",
        "source=platform",
    ]
    for phrase in banned:
        assert phrase not in lowered

    fragment_tail_re = re.compile(r"\b(?:was|is|are|the|and|to|of)\s*$", re.IGNORECASE)
    for line in merged.splitlines():
        text = line.strip()
        if not text or len(text) < 10:
            continue
        if "=" in text:
            continue
        assert fragment_tail_re.search(text) is None
