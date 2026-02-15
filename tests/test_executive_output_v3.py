"""Executive Output v3 tests — PPTX/DOCX/Notion/XMind generation.

Covers:
- File existence and minimum sizes (docx/pptx/xmind >10KB, notion >1KB)
- DOCX has embedded images (word/media/)
- PPTX has embedded images (ppt/media/)
- Notion page has required sections (今日重點/決策卡/風險/待決問題)
- XMind is valid zip with content.json + metadata.json, JSON parseable
- Banned words scan across all outputs
- Empty cards edge case
"""

import json
import re
import sys
import zipfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from schemas.education_models import EduNewsCard, SystemHealthReport
from schemas.models import (
    DeepAnalysisReport,
    ItemDeepDive,
    MergedResult,
    SchemaA,
    SchemaB,
    SchemaC,
)

# ---------------------------------------------------------------------------
# Banned words list
# ---------------------------------------------------------------------------

BANNED_WORDS = [
    "ai捕捉", "AI Intel", "Z1", "Z2", "Z3", "Z4", "Z5",
    "pipeline", "ETL", "verify_run", "ingestion", "ai_core",
]

# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

SAMPLE_METRICS: dict = {
    "run_id": "test_run_v3",
    "total_items": 2,
    "enrich_success_rate": 80.0,
    "enrich_latency_p50": 4.0,
    "enrich_latency_p95": 11.0,
    "entity_noise_removed": 2,
    "total_runtime_seconds": 30.0,
    "enrich_fail_reasons": {"blocked": 1},
}


def _make_health() -> SystemHealthReport:
    return SystemHealthReport(
        success_rate=80.0,
        p50_latency=4.0,
        p95_latency=11.0,
        entity_noise_removed=2,
        total_runtime=30.0,
        run_id="test_run_v3",
        fail_reasons={"blocked": 1},
    )


def _make_valid_card(idx: int = 1) -> EduNewsCard:
    return EduNewsCard(
        item_id=f"news_{idx:03d}",
        is_valid_news=True,
        title_plain=f"Test News Title {idx}",
        what_happened="Something important happened in the tech world.",
        why_important="This affects the entire industry supply chain.",
        focus_action="Monitor subsequent announcements from major players.",
        metaphor="Like upgrading a factory assembly line.",
        fact_check_confirmed=["Fact A confirmed", "Fact B confirmed"],
        fact_check_unverified=["Claim X needs verification"],
        evidence_lines=["Evidence: original text excerpt here"],
        technical_interpretation="Technical analysis of the event.",
        derivable_effects=["Direct effect on market pricing"],
        speculative_effects=["Possible regulatory response"],
        observation_metrics=["Watch industry adoption rate"],
        action_items=["Research latest reports", "Assess impact"],
        image_suggestions=["Image suggestion placeholder"],
        video_suggestions=["YouTube search: test topic"],
        reading_suggestions=["Google search: test analysis"],
        source_url="https://example.com/news",
        category="tech",
        signal_strength=0.8,
        final_score=8.0,
        source_name="TestSource",
    )


def _make_invalid_card() -> EduNewsCard:
    return EduNewsCard(
        item_id="invalid_001",
        is_valid_news=False,
        invalid_reason="System banner, not real news",
        title_plain="Invalid Content",
        what_happened="Scraper captured a login page.",
        why_important="Recognizing invalid content is important.",
        invalid_cause="blocked / anti-scraping",
        invalid_fix="Adjust scraping strategy",
        evidence_lines=["Original: You signed in with another tab"],
        source_url="N/A",
        signal_strength=0.0,
    )


def _make_cards() -> list[EduNewsCard]:
    return [_make_valid_card(1), _make_valid_card(2), _make_invalid_card()]


# ---------------------------------------------------------------------------
# DOCX Generator Tests
# ---------------------------------------------------------------------------


class TestExecutiveDocx:
    def test_generates_docx_file(self, tmp_path):
        from core.doc_generator import generate_executive_docx

        out = tmp_path / "test.docx"
        result = generate_executive_docx(
            cards=_make_cards(),
            health=_make_health(),
            report_time="2026-02-14 10:00",
            total_items=3,
            output_path=out,
        )
        assert result.exists()
        assert result.suffix == ".docx"
        assert result.stat().st_size > 10000  # >10KB

    def test_docx_has_embedded_image(self, tmp_path):
        from core.doc_generator import generate_executive_docx

        out = tmp_path / "test.docx"
        generate_executive_docx(
            cards=_make_cards(),
            health=_make_health(),
            report_time="2026-02-14 10:00",
            total_items=3,
            output_path=out,
        )
        with zipfile.ZipFile(str(out)) as z:
            media_files = [n for n in z.namelist() if n.startswith("word/media/")]
            assert len(media_files) >= 1, "DOCX must contain at least 1 embedded image"

    def test_docx_has_tables(self, tmp_path):
        from core.doc_generator import generate_executive_docx
        from docx import Document

        out = tmp_path / "test.docx"
        generate_executive_docx(
            cards=_make_cards(),
            health=_make_health(),
            report_time="2026-02-14 10:00",
            total_items=3,
            output_path=out,
        )
        doc = Document(str(out))
        assert len(doc.tables) >= 1  # overview table

    def test_docx_no_banned_words(self, tmp_path):
        from core.doc_generator import generate_executive_docx
        from docx import Document

        out = tmp_path / "test.docx"
        generate_executive_docx(
            cards=_make_cards(),
            health=_make_health(),
            report_time="2026-02-14 10:00",
            total_items=3,
            output_path=out,
        )
        doc = Document(str(out))
        full_text = " ".join(p.text for p in doc.paragraphs)
        for t in doc.tables:
            for row in t.rows:
                for cell in row.cells:
                    full_text += " " + cell.text
        for bw in BANNED_WORDS:
            assert bw not in full_text, f"Banned word '{bw}' found in DOCX"


# ---------------------------------------------------------------------------
# PPTX Generator Tests
# ---------------------------------------------------------------------------


class TestExecutivePptx:
    def test_generates_pptx_file(self, tmp_path):
        from core.ppt_generator import generate_executive_ppt

        out = tmp_path / "test.pptx"
        result = generate_executive_ppt(
            cards=_make_cards(),
            health=_make_health(),
            report_time="2026-02-14 10:00",
            total_items=3,
            output_path=out,
        )
        assert result.exists()
        assert result.suffix == ".pptx"
        assert result.stat().st_size > 10000  # >10KB

    def test_pptx_has_embedded_image(self, tmp_path):
        from core.ppt_generator import generate_executive_ppt

        out = tmp_path / "test.pptx"
        generate_executive_ppt(
            cards=_make_cards(),
            health=_make_health(),
            report_time="2026-02-14 10:00",
            total_items=3,
            output_path=out,
        )
        with zipfile.ZipFile(str(out)) as z:
            media_files = [n for n in z.namelist() if n.startswith("ppt/media/")]
            assert len(media_files) >= 1, "PPTX must contain at least 1 embedded image"

    def test_pptx_no_banned_words(self, tmp_path):
        from pptx import Presentation

        from core.ppt_generator import generate_executive_ppt

        out = tmp_path / "test.pptx"
        generate_executive_ppt(
            cards=_make_cards(),
            health=_make_health(),
            report_time="2026-02-14 10:00",
            total_items=3,
            output_path=out,
        )
        prs = Presentation(str(out))
        full_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        full_text += " " + p.text
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            full_text += " " + cell.text
        for bw in BANNED_WORDS:
            assert bw not in full_text, f"Banned word '{bw}' found in PPTX"


# ---------------------------------------------------------------------------
# Notion Page Tests
# ---------------------------------------------------------------------------


class TestNotionPage:
    def test_generates_notion_md(self, tmp_path):
        from core.notion_generator import generate_notion_page

        out = tmp_path / "notion_page.md"
        result = generate_notion_page(
            cards=_make_cards(),
            health=_make_health(),
            report_time="2026-02-14 10:00",
            total_items=3,
            output_path=out,
        )
        assert result.exists()
        assert result.stat().st_size > 1000  # >1KB

    def test_notion_has_required_sections(self, tmp_path):
        from core.notion_generator import generate_notion_page

        out = tmp_path / "notion_page.md"
        generate_notion_page(
            cards=_make_cards(),
            health=_make_health(),
            report_time="2026-02-14 10:00",
            total_items=3,
            output_path=out,
        )
        content = out.read_text(encoding="utf-8")
        assert "今日重點" in content
        assert "決策卡" in content
        assert "風險清單" in content or "風險" in content
        assert "待決問題" in content

    def test_notion_no_banned_words(self, tmp_path):
        from core.notion_generator import generate_notion_page

        out = tmp_path / "notion_page.md"
        generate_notion_page(
            cards=_make_cards(),
            health=_make_health(),
            report_time="2026-02-14 10:00",
            total_items=3,
            output_path=out,
        )
        content = out.read_text(encoding="utf-8")
        for bw in BANNED_WORDS:
            assert bw not in content, f"Banned word '{bw}' found in Notion page"


# ---------------------------------------------------------------------------
# XMind Tests
# ---------------------------------------------------------------------------


class TestXmindGenerator:
    def test_generates_xmind_file(self, tmp_path):
        from core.xmind_generator import generate_xmind

        out = tmp_path / "mindmap.xmind"
        result = generate_xmind(
            cards=_make_cards(),
            health=_make_health(),
            report_time="2026-02-14 10:00",
            output_path=out,
        )
        assert result.exists()
        assert result.suffix == ".xmind"
        assert result.stat().st_size > 10_000  # >10KB

    def test_xmind_is_valid_zip(self, tmp_path):
        from core.xmind_generator import generate_xmind

        out = tmp_path / "mindmap.xmind"
        generate_xmind(
            cards=_make_cards(),
            health=_make_health(),
            report_time="2026-02-14 10:00",
            output_path=out,
        )
        assert zipfile.is_zipfile(str(out))

    def test_xmind_contains_required_files(self, tmp_path):
        from core.xmind_generator import generate_xmind

        out = tmp_path / "mindmap.xmind"
        generate_xmind(
            cards=_make_cards(),
            health=_make_health(),
            report_time="2026-02-14 10:00",
            output_path=out,
        )
        with zipfile.ZipFile(str(out)) as z:
            names = z.namelist()
            assert "content.json" in names
            assert "metadata.json" in names

    def test_xmind_json_parseable(self, tmp_path):
        from core.xmind_generator import generate_xmind

        out = tmp_path / "mindmap.xmind"
        generate_xmind(
            cards=_make_cards(),
            health=_make_health(),
            report_time="2026-02-14 10:00",
            output_path=out,
        )
        with zipfile.ZipFile(str(out)) as z:
            content = json.loads(z.read("content.json"))
            metadata = json.loads(z.read("metadata.json"))
            assert isinstance(content, list)
            assert len(content) >= 1
            assert "rootTopic" in content[0]
            assert "creator" in metadata
            assert "timestamp" in metadata
            assert "theme" in metadata
            assert "zoom" in metadata

    def test_xmind_no_banned_words(self, tmp_path):
        from core.xmind_generator import generate_xmind

        out = tmp_path / "mindmap.xmind"
        generate_xmind(
            cards=_make_cards(),
            health=_make_health(),
            report_time="2026-02-14 10:00",
            output_path=out,
        )
        with zipfile.ZipFile(str(out)) as z:
            content_text = z.read("content.json").decode("utf-8")
            for bw in BANNED_WORDS:
                assert bw not in content_text, f"Banned word '{bw}' found in XMind content"


# ---------------------------------------------------------------------------
# Integration: generate_executive_reports
# ---------------------------------------------------------------------------


class TestExecutiveReportsIntegration:
    def test_generate_all_four_files(self, tmp_path):
        from core.education_renderer import generate_executive_reports

        pptx_path, docx_path, notion_path, xmind_path = generate_executive_reports(
            results=None,
            metrics=SAMPLE_METRICS,
            project_root=tmp_path,
        )

        assert pptx_path.exists()
        assert docx_path.exists()
        assert notion_path.exists()
        assert xmind_path.exists()
        assert pptx_path.suffix == ".pptx"
        assert docx_path.suffix == ".docx"
        assert notion_path.name == "notion_page.md"
        assert xmind_path.name == "mindmap.xmind"

    def test_empty_cards_still_generates(self, tmp_path):
        from core.education_renderer import generate_executive_reports

        pptx_path, docx_path, notion_path, xmind_path = generate_executive_reports(
            results=None,
            metrics=SAMPLE_METRICS,
            project_root=tmp_path,
        )

        assert pptx_path.exists()
        assert docx_path.exists()
        assert notion_path.exists()
        assert xmind_path.exists()
        # All should be non-empty
        assert pptx_path.stat().st_size > 500
        assert docx_path.stat().st_size > 500
        assert notion_path.stat().st_size > 100
        assert xmind_path.stat().st_size > 100
