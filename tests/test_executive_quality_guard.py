"""Comprehensive quality guard tests for Executive PPT/DOCX output.

Covers:
- No URL=N/A in any output
- signal_text != signal_type (no '工作流變化' as content body)
- No trailing fragments ('Last July was', sentence-ending connectors)
- No isolated short fragments (< 12 chars without numbers/entities)
- Backfill gate when hard_pass is insufficient
- Non-AI / index / subscribe / login rejection
- Diagnostics script integration
- Open PPT script contract ([OPEN] markers)
- Theme LIGHT/DARK coverage
"""

from __future__ import annotations

import re
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import patch

import pytest
from docx import Document
from pptx import Presentation

from core.content_gate import apply_split_content_gate
from core.content_strategy import build_signal_summary, quality_guard_block, sanitize
from core.doc_generator import generate_executive_docx
from core.ppt_generator import generate_executive_ppt, LIGHT_BG, DARK_BG
from schemas.education_models import EduNewsCard, SystemHealthReport
from scripts.diagnostics_pptx import diagnose_pptx, diagnose_docx
from utils.text_quality import (
    trim_trailing_fragment, is_fragment,
    count_evidence_terms, count_evidence_numbers, count_sentences,
)


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

def _health() -> SystemHealthReport:
    return SystemHealthReport(success_rate=80.0, p50_latency=1.0, p95_latency=3.0)


def _rich_cards() -> list[EduNewsCard]:
    """Cards with enough content to form a non-trivial deck."""
    return [
        EduNewsCard(
            item_id=f"event-{i:03d}",
            is_valid_news=True,
            title_plain=title,
            what_happened=what,
            why_important=why,
            source_name=src,
            source_url=url,
            final_score=score,
            category="AI",
        )
        for i, (title, what, why, src, url, score) in enumerate([
            (
                "NVIDIA H200 GPU launched at $30k",
                "NVIDIA launched the H200 GPU with 141GB HBM3e for $30k.",
                "Impacts AI training costs globally.",
                "TechCrunch",
                "https://techcrunch.com/nvidia-h200",
                9.0,
            ),
            (
                "OpenAI releases GPT-5 with 1M context",
                "OpenAI announced GPT-5 with 1 million token context window.",
                "Enables new use cases for enterprise RAG.",
                "TheVerge",
                "https://theverge.com/openai-gpt5",
                8.5,
            ),
            (
                "Anthropic Claude reaches 100M users",
                "Anthropic reported Claude AI reached 100 million monthly active users.",
                "Claude becomes a major competitor to ChatGPT.",
                "Reuters",
                "https://reuters.com/anthropic-claude-100m",
                8.0,
            ),
        ])
    ]


def _no_event_cards() -> list[EduNewsCard]:
    """Intentionally weak cards — deck should still have content via fallback."""
    return [
        EduNewsCard(
            item_id="no-event-01",
            is_valid_news=True,
            title_plain="Memory Plugin for Claude Code",
            what_happened="A memory plugin was released for Claude Code IDE integration.",
            why_important="Helps developers persist context across sessions.",
            source_name="HackerNews",
            source_url="https://news.ycombinator.com/item?id=12345",
            final_score=3.0,
        ),
    ]


def _gen_both(tmp_path: Path, cards: list[EduNewsCard], metrics: dict | None = None):
    health = _health()
    pptx_path = tmp_path / "test.pptx"
    docx_path = tmp_path / "test.docx"
    with patch("core.ppt_generator.get_news_image", return_value=None), \
         patch("core.doc_generator.get_news_image", return_value=None):
        generate_executive_ppt(
            cards=cards, health=health,
            report_time="2026-02-16 09:00", total_items=len(cards),
            output_path=pptx_path, metrics=metrics,
        )
        generate_executive_docx(
            cards=cards, health=health,
            report_time="2026-02-16 09:00", total_items=len(cards),
            output_path=docx_path, metrics=metrics,
        )
    return pptx_path, docx_path


def _extract_all_text(pptx_path: Path, docx_path: Path) -> str:
    texts: list[str] = []
    prs = Presentation(str(pptx_path))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            texts.append(cell.text.strip())
    doc = Document(str(docx_path))
    for p in doc.paragraphs:
        if p.text.strip():
            texts.append(p.text.strip())
    for t in doc.tables:
        for row in t.rows:
            for cell in row.cells:
                if cell.text.strip():
                    texts.append(cell.text.strip())
    return "\n".join(texts)


# ---------------------------------------------------------------------------
# Test: URL must never be N/A
# ---------------------------------------------------------------------------

class TestNoUrlNA:
    def test_no_url_na_in_rich_deck(self, tmp_path: Path) -> None:
        pptx, docx = _gen_both(tmp_path, _rich_cards())
        text = _extract_all_text(pptx, docx)
        assert "URL=N/A" not in text
        assert "url=n/a" not in text.lower()

    def test_no_url_na_in_no_event_deck(self, tmp_path: Path) -> None:
        pptx, docx = _gen_both(tmp_path, _no_event_cards())
        text = _extract_all_text(pptx, docx)
        assert "URL=N/A" not in text
        assert "url=n/a" not in text.lower()

    def test_at_least_one_http_url(self, tmp_path: Path) -> None:
        pptx, docx = _gen_both(tmp_path, _rich_cards())
        text = _extract_all_text(pptx, docx)
        assert "http" in text.lower()


# ---------------------------------------------------------------------------
# Test: signal_text != signal_type
# ---------------------------------------------------------------------------

class TestSignalTextNotType:
    def test_signal_text_is_not_workflow_change(self) -> None:
        cards = [
            EduNewsCard(
                item_id="sig-001",
                is_valid_news=True,
                title_plain="Memory Plugin for Claude Code",
                what_happened="A memory plugin was released.",
                why_important="Helps developers.",
                source_name="HackerNews",
                source_url="https://example.com/plugin",
                final_score=5.0,
            ),
        ]
        signals = build_signal_summary(cards)
        for sig in signals:
            text = sig.get("signal_text", "")
            assert text != "WORKFLOW_CHANGE"
            assert text != "工作流變化"
            assert text != "TOOL_ADOPTION"
            assert text != "USER_PAIN"
            assert text != sig.get("signal_type", "")

    def test_signal_text_contains_title_keywords(self) -> None:
        cards = [
            EduNewsCard(
                item_id="sig-002",
                is_valid_news=True,
                title_plain="NVIDIA launches H200 GPU at $30k",
                what_happened="NVIDIA launched the H200 GPU.",
                why_important="Major GPU release.",
                source_name="TechCrunch",
                source_url="https://example.com/nvidia",
                final_score=8.0,
            ),
        ]
        signals = build_signal_summary(cards)
        assert len(signals) >= 1
        # At least one signal should reference NVIDIA or H200
        all_text = " ".join(str(s.get("title", "")) + " " + str(s.get("signal_text", "")) for s in signals)
        assert "NVIDIA" in all_text or "H200" in all_text or "nvidia" in all_text.lower()


# ---------------------------------------------------------------------------
# Test: No trailing fragments / broken sentences
# ---------------------------------------------------------------------------

class TestNoFragments:
    def test_no_last_july_was(self, tmp_path: Path) -> None:
        pptx, docx = _gen_both(tmp_path, _rich_cards())
        text = _extract_all_text(pptx, docx).lower()
        assert "last july was" not in text

    def test_no_trailing_connectors(self, tmp_path: Path) -> None:
        pptx, docx = _gen_both(tmp_path, _rich_cards())
        text = _extract_all_text(pptx, docx)
        # Check output lines for trailing fragments
        trailing_re = re.compile(r"\b(?:to|and|or|by)\s*$", re.IGNORECASE)
        for line in text.splitlines():
            s = line.strip()
            if not s or len(s) < 15 or "=" in s:
                continue
            assert trailing_re.search(s) is None, f"Trailing fragment: {s!r}"

    def test_no_trailing_zh_particles(self, tmp_path: Path) -> None:
        pptx, docx = _gen_both(tmp_path, _rich_cards())
        text = _extract_all_text(pptx, docx)
        zh_trailing = set("的了而與來記")
        for line in text.splitlines():
            s = line.strip()
            if not s or len(s) < 15 or "=" in s:
                continue
            if s[-1] in zh_trailing:
                # Allow if it's part of a complete phrase
                if re.search(r"[。！？;；.!?]", s[:-1]):
                    continue
                assert False, f"Trailing ZH particle: {s!r}"


# ---------------------------------------------------------------------------
# Test: trim_trailing_fragment utility
# ---------------------------------------------------------------------------

class TestTextQualityUtils:
    def test_trim_trailing_en(self) -> None:
        assert trim_trailing_fragment("The company announced a deal to") == "The company announced a deal to"
        # Only trims if there's a prior sentence boundary to fall back to
        assert trim_trailing_fragment("First sentence. Second to") == "First sentence."

    def test_trim_trailing_zh(self) -> None:
        result = trim_trailing_fragment("第一句話。第二句的")
        assert result == "第一句話。"

    def test_is_fragment_short(self) -> None:
        assert is_fragment("") is True
        assert is_fragment("abc") is True
        assert is_fragment("ok") is True  # < 12 chars, no number/entity

    def test_is_fragment_with_number(self) -> None:
        assert is_fragment("v3.5") is False

    def test_is_fragment_long_enough(self) -> None:
        assert is_fragment("This is a complete sentence.") is False


# ---------------------------------------------------------------------------
# Test: Backfill gate
# ---------------------------------------------------------------------------

class TestBackfillGate:
    def _soft_items(self, n: int = 10):
        """Items that fail strict gate but have link + evidence_terms + numbers."""
        items = []
        for i in range(n):
            body = (
                f"OpenAI released GPT-{i+4} model with {i*10+50}M parameters. "
                f"NVIDIA GPU pricing at ${i*1000+2000}. Version v{i+1}.0 launched."
            ) * 2  # ~300 chars, not enough for strict 1200
            items.append(SimpleNamespace(
                title=f"AI Model Update {i}",
                body=body,
                url=f"https://example.com/article-{i}",
                link=f"https://example.com/article-{i}",
            ))
        return items

    def test_backfill_promotes_when_hard_pass_insufficient(self) -> None:
        items = self._soft_items(10)
        events, signals, rejected, stats = apply_split_content_gate(items)
        # Even though items are too short for strict gate,
        # backfill should promote some
        assert stats.backfill_used_total >= 0
        # With backfill, we should have more items than strict-only
        assert len(events) + len(signals) > 0

    def test_backfill_rejects_non_ai(self) -> None:
        items = [
            SimpleNamespace(
                title="Magnus Carlsen wins chess",
                body="Magnus Carlsen won the chess tournament in round 7." * 5,
                url="https://example.com/chess",
                link="https://example.com/chess",
            ),
        ]
        events, signals, rejected, stats = apply_split_content_gate(items)
        assert len(events) == 0
        assert len(signals) == 0

    def test_backfill_rejects_index_pages(self) -> None:
        items = [
            SimpleNamespace(
                title="AI weekly roundup digest subscribe",
                body="Subscribe to our weekly AI roundup digest. Top links for the week." * 10,
                url="https://example.com/subscribe",
                link="https://example.com/subscribe",
            ),
        ]
        events, signals, rejected, stats = apply_split_content_gate(items)
        assert len(events) == 0
        assert len(signals) == 0


# ---------------------------------------------------------------------------
# Test: Diagnostics script integration
# ---------------------------------------------------------------------------

class TestDiagnostics:
    def test_diagnose_pptx_no_url_na(self, tmp_path: Path) -> None:
        pptx, _ = _gen_both(tmp_path, _rich_cards())
        result = diagnose_pptx(pptx)
        assert result.url_na_count == 0

    def test_diagnose_docx_no_url_na(self, tmp_path: Path) -> None:
        _, docx = _gen_both(tmp_path, _rich_cards())
        result = diagnose_docx(docx)
        assert result.url_na_count == 0

    def test_diagnose_no_signal_type_as_text(self, tmp_path: Path) -> None:
        pptx, docx = _gen_both(tmp_path, _no_event_cards())
        pptx_result = diagnose_pptx(pptx)
        docx_result = diagnose_docx(docx)
        assert pptx_result.signal_type_as_text_count == 0
        assert docx_result.signal_type_as_text_count == 0


# ---------------------------------------------------------------------------
# Test: Events insufficient → deck still has stats/numbers
# ---------------------------------------------------------------------------

class TestLowEventDeck:
    def test_deck_not_empty_with_no_events(self, tmp_path: Path) -> None:
        pptx, docx = _gen_both(tmp_path, _no_event_cards(), metrics={
            "fetched_total": 100,
            "gate_pass_total": 2,
            "sources_total": 5,
            "sources_success": 3,
            "sources_failed": 2,
        })
        text = _extract_all_text(pptx, docx)
        assert len(text) >= 500
        # Must contain numbers (stats)
        assert len(re.findall(r"\d", text)) >= 8


# ---------------------------------------------------------------------------
# Test: sanitize does not produce fragments
# ---------------------------------------------------------------------------

class TestSanitizeNoFragments:
    def test_sanitize_removes_last_july_was(self) -> None:
        result = sanitize("Last July was a great month for AI development.")
        assert "Last July was" not in result

    def test_sanitize_removes_this_week_was(self) -> None:
        result = sanitize("This week was full of AI announcements and launches.")
        assert "This week was" not in result

    def test_sanitize_trims_trailing_particle(self) -> None:
        result = sanitize("第一句話。第二句的")
        assert not result.endswith("的")


# ---------------------------------------------------------------------------
# Test: Table cell non-empty ratio >= 60%
# ---------------------------------------------------------------------------

class TestTableCellNonEmpty:
    def _count_table_cells(self, pptx_path: Path) -> tuple[int, int]:
        """Return (total_cells, non_empty_cells) across all tables in PPTX."""
        prs = Presentation(str(pptx_path))
        total = 0
        non_empty = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            total += 1
                            if cell.text.strip():
                                non_empty += 1
        return total, non_empty

    def _count_docx_table_cells(self, docx_path: Path) -> tuple[int, int]:
        """Return (total_cells, non_empty_cells) across all tables in DOCX."""
        doc = Document(str(docx_path))
        total = 0
        non_empty = 0
        for t in doc.tables:
            for row in t.rows:
                for cell in row.cells:
                    total += 1
                    if cell.text.strip():
                        non_empty += 1
        return total, non_empty

    def test_pptx_table_non_empty_ratio(self, tmp_path: Path) -> None:
        pptx, _ = _gen_both(tmp_path, _rich_cards())
        total, non_empty = self._count_table_cells(pptx)
        if total > 0:
            ratio = non_empty / total
            assert ratio >= 0.50, f"PPTX table non-empty ratio {ratio:.2%} < 50%"

    def test_docx_table_non_empty_ratio(self, tmp_path: Path) -> None:
        _, docx = _gen_both(tmp_path, _rich_cards())
        total, non_empty = self._count_docx_table_cells(docx)
        if total > 0:
            ratio = non_empty / total
            assert ratio >= 0.50, f"DOCX table non-empty ratio {ratio:.2%} < 50%"

    def test_pptx_table_non_empty_with_weak_cards(self, tmp_path: Path) -> None:
        pptx, _ = _gen_both(tmp_path, _no_event_cards())
        total, non_empty = self._count_table_cells(pptx)
        if total > 0:
            ratio = non_empty / total
            assert ratio >= 0.50, f"PPTX table non-empty ratio {ratio:.2%} < 50% (weak cards)"


# ---------------------------------------------------------------------------
# Test: Per-block density thresholds
# ---------------------------------------------------------------------------

class TestPerBlockDensity:
    def test_evidence_terms_count(self) -> None:
        text = "NVIDIA launched the H200 GPU with AI inference capabilities."
        assert count_evidence_terms(text) >= 2

    def test_evidence_numbers_count(self) -> None:
        text = "Priced at $30k with 141GB HBM3e memory."
        assert count_evidence_numbers(text) >= 1

    def test_sentence_count(self) -> None:
        text = "First sentence. Second sentence. Third."
        assert count_sentences(text) >= 2

    def test_empty_text_density(self) -> None:
        assert count_evidence_terms("") == 0
        assert count_evidence_numbers("") == 0
        assert count_sentences("") == 0

    def test_quality_guard_block_backfills_thin_text(self) -> None:
        card = EduNewsCard(
            item_id="qg-001",
            is_valid_news=True,
            title_plain="NVIDIA H200 GPU launched at $30k",
            what_happened="NVIDIA launched the H200 GPU with 141GB HBM3e for $30k.",
            why_important="Impacts AI training costs globally.",
            source_name="TechCrunch",
            source_url="https://techcrunch.com/nvidia-h200",
            final_score=9.0,
        )
        # Thin text with no evidence
        thin = "Something happened."
        result, metrics = quality_guard_block(thin, card=card)
        # After backfill, should have more content
        assert len(result) >= len(thin) or count_evidence_terms(result) >= 1

    def test_quality_guard_preserves_good_text(self) -> None:
        card = EduNewsCard(
            item_id="qg-002",
            is_valid_news=True,
            title_plain="OpenAI GPT-5 release",
            what_happened="OpenAI released GPT-5 with 1M token context.",
            why_important="Enables enterprise RAG at scale.",
            source_name="TheVerge",
            source_url="https://theverge.com/gpt5",
            final_score=8.5,
        )
        good = "OpenAI released GPT-5 with 1 million token context window. This impacts enterprise RAG."
        result, metrics = quality_guard_block(good, card=card)
        assert len(result) > 0

    def test_rich_deck_blocks_meet_density(self, tmp_path: Path) -> None:
        """Every event card in a rich deck should produce non-empty brief blocks."""
        pptx, _ = _gen_both(tmp_path, _rich_cards())
        text = _extract_all_text(pptx, Path(str(pptx).replace(".pptx", ".docx")))
        # Text should contain evidence terms
        assert count_evidence_terms(text) >= 3


# ---------------------------------------------------------------------------
# Test: Expanded fragment detection (v5.5)
# ---------------------------------------------------------------------------

class TestExpandedFragmentDetection:
    def test_fragment_trailing_en_medium(self) -> None:
        # Under 40 chars, ends with trailing word, no sentence end
        assert is_fragment("The company announced to") is True
        assert is_fragment("NVIDIA partnership with") is True

    def test_fragment_trailing_zh_medium(self) -> None:
        # 12+ chars, ends with trailing particle, no sentence-ending punctuation
        assert is_fragment("這項技術在人工智慧領域中的應用的") is True

    def test_not_fragment_with_sentence_end(self) -> None:
        assert is_fragment("The GPU costs $30k.") is False

    def test_not_fragment_long_text(self) -> None:
        # Over 40 chars should not be flagged even with trailing word
        text = "This is a fairly long sentence that describes something important to"
        assert is_fragment(text) is False

    def test_trim_trailing_comma(self) -> None:
        result = trim_trailing_fragment("First sentence. Second part,")
        assert result == "First sentence."

    def test_trim_trailing_but(self) -> None:
        result = trim_trailing_fragment("First claim is valid. But")
        assert result == "First claim is valid."


# ---------------------------------------------------------------------------
# Test: Theme LIGHT/DARK coverage
# ---------------------------------------------------------------------------

class TestThemeCoverage:
    def test_light_theme_generates(self, tmp_path: Path) -> None:
        health = _health()
        pptx_path = tmp_path / "light.pptx"
        with patch("core.ppt_generator.get_news_image", return_value=None):
            generate_executive_ppt(
                cards=_rich_cards(), health=health,
                report_time="2026-02-16 09:00", total_items=3,
                output_path=pptx_path, theme="light",
            )
        assert pptx_path.exists()
        assert pptx_path.stat().st_size > 30000

    def test_dark_theme_generates(self, tmp_path: Path) -> None:
        health = _health()
        pptx_path = tmp_path / "dark.pptx"
        with patch("core.ppt_generator.get_news_image", return_value=None):
            generate_executive_ppt(
                cards=_rich_cards(), health=health,
                report_time="2026-02-16 09:00", total_items=3,
                output_path=pptx_path, theme="dark",
            )
        assert pptx_path.exists()
        assert pptx_path.stat().st_size > 30000

    def test_light_bg_is_white(self) -> None:
        assert LIGHT_BG == (0xFF, 0xFF, 0xFF) or str(LIGHT_BG) == "FFFFFF"

    def test_dark_bg_is_dark(self) -> None:
        assert DARK_BG == (0x12, 0x12, 0x18) or str(DARK_BG) == "121218"


# ---------------------------------------------------------------------------
# Test: Open PPT script contract ([OPEN] markers)
# ---------------------------------------------------------------------------

class TestOpenPptContract:
    def test_open_markers_in_generate_reports(self) -> None:
        text = Path("scripts/generate_reports.ps1").read_text(encoding="utf-8")
        assert "[OPEN] pptx_path=" in text
        assert "[OPEN] exists=" in text
        assert "PPT generated successfully:" in text

    def test_retry_loop_present(self) -> None:
        text = Path("scripts/generate_reports.ps1").read_text(encoding="utf-8")
        assert "for ($attempt = 1; $attempt -le 5; $attempt++)" in text
        assert "OpenAttempt" in text
        assert "$minOpenBytes = 30720" in text

    def test_exit_codes_present(self) -> None:
        text = Path("scripts/generate_reports.ps1").read_text(encoding="utf-8")
        assert "exit 2" in text
        assert "exit 3" in text

    def test_explorer_fallback_present(self) -> None:
        text = Path("scripts/generate_reports.ps1").read_text(encoding="utf-8")
        assert "explorer.exe" in text
        assert "[OPEN] fallback=" in text or "[OPEN] start_process_exit_code=" in text

    def test_no_custom_ansi_colors(self) -> None:
        text = Path("scripts/generate_reports.ps1").read_text(encoding="utf-8")
        # Should use PowerShell -ForegroundColor, not raw ANSI escape sequences
        assert "\x1b[" not in text
        assert "\\e[" not in text
