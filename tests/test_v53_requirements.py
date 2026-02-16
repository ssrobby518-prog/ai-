"""v5.3 requirement tests.

Covers:
1. test_non_ai_topic_rejected — chess/Magnus Carlsen rejected as non_ai_topic
2. test_soft_pass_nonzero_and_observable — soft_pass_total > 0, top_rejected_reasons present
3. test_no_event_fallback_density — Event=0 deck has Top3 signals with evidence_terms/numbers
4. test_placeholder_guard_extended — "Last July was" + trailing fragment words banned
5. test_theme_default_light_and_dark_explicit — default=light, dark=explicit
"""

from __future__ import annotations

import re
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import patch

from pptx import Presentation

from core.content_gate import (
    apply_adaptive_content_gate,
    apply_split_content_gate,
    is_ai_relevant,
)
from core.content_strategy import build_signal_summary, build_corp_watch_summary
from core.ppt_generator import (
    DARK_BG,
    LIGHT_BG,
    generate_executive_ppt,
)
from schemas.education_models import EduNewsCard, SystemHealthReport


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _ai_body(n: int = 20) -> str:
    """Generate a realistic AI-related article body."""
    return (
        "OpenAI GPT-5 rollout reached 12 enterprise teams in 2026 with 35% latency reduction. "
        "Microsoft Azure benchmark confirmed 90ms response and $2.1M annual cost savings. "
        "Google and NVIDIA jointly reported production stability across 3 regions. "
    ) * n


def _chess_body(n: int = 20) -> str:
    """Generate a non-AI chess article body."""
    return (
        "Magnus Carlsen won the chess championship in round 7 with a brilliant endgame. "
        "The tournament featured 14 grandmasters from 8 countries competing over 3 weeks. "
        "Viewers praised the dramatic queen sacrifice that sealed the victory. "
    ) * n


def _no_event_cards() -> list[EduNewsCard]:
    return [
        EduNewsCard(
            item_id="idx-001",
            is_valid_news=True,
            title_plain="OpenAI releases GPT-5 with improved reasoning",
            what_happened="OpenAI announced GPT-5 launch with 40% benchmark improvement in 2026.",
            why_important="Impacts enterprise AI adoption and cloud inference costs.",
            source_name="TechCrunch",
            source_url="https://techcrunch.com/openai-gpt5",
            final_score=7.5,
            fact_check_confirmed=["GPT-5 benchmark score 92.1%", "Pricing at $0.03/1k tokens"],
            evidence_lines=["Azure integration confirmed", "3 region rollout"],
        ),
        EduNewsCard(
            item_id="idx-002",
            is_valid_news=True,
            title_plain="NVIDIA H200 GPU production ramp",
            what_happened="NVIDIA confirmed H200 GPU mass production starting Q2 2026.",
            why_important="Affects AI infrastructure supply and pricing globally.",
            source_name="Reuters",
            source_url="https://reuters.com/nvidia-h200",
            final_score=6.8,
            fact_check_confirmed=["H200 specs: 141GB HBM3e", "$30k per unit"],
            evidence_lines=["TSMC 4nm process", "2M units planned for 2026"],
        ),
    ]


# ---------------------------------------------------------------------------
# 1. test_non_ai_topic_rejected
# ---------------------------------------------------------------------------


class TestNonAiTopicRejected:
    """Chess / Magnus Carlsen must be rejected as non_ai_topic."""

    def test_chess_article_rejected_by_split_gate(self):
        items = [
            SimpleNamespace(
                title="Magnus Carlsen wins chess championship",
                body=_chess_body(20),
            ),
        ]
        _event, _signal, rejected_map, stats = apply_split_content_gate(items)
        assert len(_event) == 0
        assert len(_signal) == 0
        assert stats.rejected_total >= 1
        assert "non_ai_topic" in stats.rejected_by_reason

    def test_chess_article_rejected_by_adaptive_gate(self):
        items = [
            SimpleNamespace(
                title="Magnus Carlsen wins chess championship",
                body=_chess_body(20),
            ),
        ]
        kept, _rej, stats = apply_adaptive_content_gate(items, min_keep_items=1)
        assert len(kept) == 0
        assert "non_ai_topic" in stats.rejected_by_reason

    def test_ai_article_passes(self):
        items = [
            SimpleNamespace(
                title="OpenAI launches GPT-5 for enterprise",
                body=_ai_body(20),
            ),
        ]
        _event, signal, _rej, stats = apply_split_content_gate(items)
        assert stats.event_gate_pass_total >= 1 or stats.signal_gate_pass_total >= 1

    def test_is_ai_relevant_rejects_chess(self):
        assert is_ai_relevant("Magnus Carlsen wins", _chess_body(1)) is False

    def test_is_ai_relevant_accepts_ai(self):
        assert is_ai_relevant("OpenAI GPT-5", _ai_body(1)) is True


# ---------------------------------------------------------------------------
# 2. test_soft_pass_nonzero_and_observable
# ---------------------------------------------------------------------------


class TestSoftPassNonzeroAndObservable:
    """soft_pass_total > 0 and top_rejected_reasons has values."""

    def test_soft_pass_observable(self):
        short_ai_body = (
            "Google announced Gemini 2.0 with multimodal reasoning improvements. "
            "The model achieves 85% on MMLU benchmark. "
        ) * 4  # ~400 chars, passes soft but not hard (1200)
        long_ai_body = _ai_body(20)  # passes hard

        items = [
            SimpleNamespace(title="Gemini 2.0 update", body=short_ai_body),
            SimpleNamespace(title="OpenAI enterprise rollout", body=long_ai_body),
            SimpleNamespace(title="Chess tournament", body=_chess_body(20)),
        ]
        kept, _rej, stats = apply_adaptive_content_gate(items, min_keep_items=2)
        # At least one should be soft pass (relaxed or density)
        assert stats.soft_pass_total >= 1
        # rejected_reason_top should have entries (chess rejected)
        assert len(stats.rejected_reason_top) >= 1
        assert stats.rejected_by_reason.get("non_ai_topic", 0) >= 1


# ---------------------------------------------------------------------------
# 3. test_no_event_fallback_density
# ---------------------------------------------------------------------------


class TestNoEventFallbackDensity:
    """Event=0 deck still has Top3 signals with evidence_terms/numbers."""

    def test_signal_evidence_fields(self):
        signals = build_signal_summary(_no_event_cards())
        assert len(signals) >= 3
        at_least_one_has_numbers = False
        for sig in signals[:3]:
            assert "evidence_terms" in sig
            assert isinstance(sig["evidence_terms"], list)
            assert len(sig["evidence_terms"]) >= 2, (
                f"evidence_terms should have >=2 items, got {sig['evidence_terms']}"
            )
            assert "evidence_numbers" in sig
            assert isinstance(sig["evidence_numbers"], list)
            if len(sig["evidence_numbers"]) >= 1:
                at_least_one_has_numbers = True
        assert at_least_one_has_numbers, "At least one signal should have evidence_numbers >= 1"

    def test_corp_watch_update_type_counts(self):
        corp = build_corp_watch_summary(_no_event_cards())
        assert "update_type_counts" in corp
        expected_keys = {"model", "product", "cloud_infra", "ecosystem", "risk_policy"}
        assert set(corp["update_type_counts"].keys()) == expected_keys
        assert "top_fail_reasons" in corp
        assert len(corp["top_fail_reasons"]) >= 1


# ---------------------------------------------------------------------------
# 4. test_placeholder_guard_extended
# ---------------------------------------------------------------------------


class TestPlaceholderGuardExtended:
    """Banned: 'Last July was' + trailing fragment words (was/the/of/to)."""

    def _gen_ppt(self, tmp_path: Path, cards: list[EduNewsCard]) -> str:
        out = tmp_path / "guard_ext.pptx"
        health = SystemHealthReport(success_rate=50.0, p50_latency=1.0, p95_latency=2.0)
        with patch("core.ppt_generator.get_news_image", return_value=None):
            generate_executive_ppt(
                cards=cards,
                health=health,
                report_time="2026-02-16 09:00",
                total_items=len(cards),
                output_path=out,
            )
        prs = Presentation(str(out))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        if p.text.strip():
                            texts.append(p.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                texts.append(cell.text.strip())
        return "\n".join(texts)

    def test_no_last_july_was(self, tmp_path: Path):
        cards = [
            EduNewsCard(
                item_id="frag-001",
                is_valid_news=True,
                title_plain="Archive index update",
                what_happened="Last July was...",
                why_important="Reference-only page without concrete event.",
                source_name="ExampleSource",
                source_url="https://example.com/archive",
                final_score=1.0,
            ),
        ]
        text = self._gen_ppt(tmp_path, cards).lower()
        assert "last july was" not in text

    def test_no_trailing_fragment_words(self, tmp_path: Path):
        cards = _no_event_cards()
        text = self._gen_ppt(tmp_path, cards)
        # Check each non-empty line doesn't end with dangling connector words
        fragment_end_re = re.compile(
            r"\b(?:was|is|are|the|of|to)\s*$",
            re.IGNORECASE,
        )
        for line in text.split("\n"):
            line = line.strip()
            if not line or len(line) < 10:
                continue
            # Skip lines that are clearly data/stats (contain = signs)
            if "=" in line:
                continue
            assert fragment_end_re.search(line) is None, (
                f"Fragment ending detected: '{line}'"
            )


# ---------------------------------------------------------------------------
# 5. test_theme_default_light_and_dark_explicit
# ---------------------------------------------------------------------------


class TestThemeDefaultLightAndDarkExplicit:
    """Default theme=light, dark requires explicit theme='dark'."""

    def test_default_is_light(self, tmp_path: Path):
        out = tmp_path / "default_light.pptx"
        health = SystemHealthReport(success_rate=80.0, p50_latency=2.0, p95_latency=5.0)
        cards = [
            EduNewsCard(
                item_id="theme-001",
                is_valid_news=True,
                title_plain="NVIDIA launches new GPU architecture",
                what_happened="Launch event with architecture updates and benchmarks.",
                why_important="Impacts AI infrastructure pricing and capacity.",
                source_name="TechCrunch",
                source_url="https://example.com/news",
                final_score=8.0,
            )
        ]
        with patch("core.ppt_generator.get_news_image", return_value=None):
            generate_executive_ppt(
                cards=cards,
                health=health,
                report_time="2026-02-16 09:00",
                total_items=1,
                output_path=out,
                # No theme parameter — should default to light
            )
        prs = Presentation(str(out))
        for slide in prs.slides:
            bg_color = slide.background.fill.fore_color.rgb
            assert bg_color != DARK_BG, f"Default theme should be LIGHT, got dark bg on slide"
            assert bg_color == LIGHT_BG

    def test_dark_explicit(self, tmp_path: Path):
        out = tmp_path / "dark_explicit.pptx"
        health = SystemHealthReport(success_rate=80.0, p50_latency=2.0, p95_latency=5.0)
        cards = [
            EduNewsCard(
                item_id="theme-002",
                is_valid_news=True,
                title_plain="NVIDIA launches new GPU architecture",
                what_happened="Launch event with architecture updates and benchmarks.",
                why_important="Impacts AI infrastructure pricing and capacity.",
                source_name="TechCrunch",
                source_url="https://example.com/news",
                final_score=8.0,
            )
        ]
        with patch("core.ppt_generator.get_news_image", return_value=None):
            generate_executive_ppt(
                cards=cards,
                health=health,
                report_time="2026-02-16 09:00",
                total_items=1,
                output_path=out,
                theme="dark",
            )
        prs = Presentation(str(out))
        for slide in prs.slides:
            assert slide.background.fill.fore_color.rgb == DARK_BG
