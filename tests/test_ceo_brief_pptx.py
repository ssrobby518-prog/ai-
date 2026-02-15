"""Tests for CEO Motion Slides Dark Theme — ppt_generator.py upgrade.

Validates:
- PPT can be generated without errors
- Dark theme (BG_DARK) is applied to slide backgrounds
- Each event card produces exactly 2 slides (brief_page1 + brief_page2)
- Q&A content, Data card, Video reference, Sources are present
"""

from __future__ import annotations

import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from unittest.mock import patch

import pytest

from core.ppt_generator import (
    BG_DARK,
    CARD_BG,
    HIGHLIGHT_YELLOW,
    TEXT_WHITE,
    generate_executive_ppt,
)

# ---------------------------------------------------------------------------
# Minimal test fixtures (avoid importing full schemas for unit tests)
# ---------------------------------------------------------------------------


@dataclass
class _FakeHealth:
    success_rate: float = 0.95
    latency_p95: float = 1.2
    noise_filtered: int = 3
    traffic_light: str = "green"


@dataclass
class _FakeCard:
    item_id: str = "test-001"
    is_valid_news: bool = True
    invalid_reason: str = ""
    title_plain: str = "Nvidia launches new GPU chip today"
    what_happened: str = "Nvidia 今日推出新款 GPU 晶片，效能提升 40%"
    why_important: str = "影響 AI 訓練成本與效能"
    focus_action: str = "評估是否採購"
    metaphor: str = "就像汽車引擎升級，所有車都能跑更快"
    fact_check_confirmed: list = field(default_factory=lambda: [
        "NVIDIA announced the chip at CES 2026",
        "Performance benchmarks show 40% improvement",
    ])
    fact_check_unverified: list = field(default_factory=list)
    evidence_lines: list = field(default_factory=lambda: [
        "NVIDIA CEO Jensen Huang presented the new architecture",
    ])
    technical_interpretation: str = "新架構採用 4nm 製程"
    derivable_effects: list = field(default_factory=lambda: [
        "AI 訓練成本可能下降 20-30%",
        "競爭對手 AMD 壓力增加",
    ])
    speculative_effects: list = field(default_factory=lambda: [
        "可能引發新一輪 AI 軍備競賽",
    ])
    observation_metrics: list = field(default_factory=list)
    action_items: list = field(default_factory=lambda: [
        "評估新 GPU 對現有 AI 專案的影響",
    ])
    image_suggestions: list = field(default_factory=list)
    video_suggestions: list = field(default_factory=lambda: [
        "NVIDIA GPU chip launch 2026",
    ])
    reading_suggestions: list = field(default_factory=list)
    source_url: str = "https://example.com/nvidia-gpu"
    invalid_cause: str = ""
    invalid_fix: str = ""
    category: str = "人工智慧"
    final_score: float = 8.5
    one_liner: str = ""


def _make_event_card(**overrides) -> _FakeCard:
    return _FakeCard(**overrides)


def _make_non_event_card() -> _FakeCard:
    return _FakeCard(
        item_id="non-event-001",
        title_plain="As part of its mission to preserve the web",
        what_happened="Curated list of resources",
        category="綜合",
        final_score=2.0,
    )


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------


class TestPPTGeneration:
    """Test that the PPT can be generated without errors."""

    def test_generate_with_event_cards(self, tmp_path: Path):
        cards = [_make_event_card(), _make_event_card(item_id="test-002",
                 title_plain="Google releases Gemini 3.0 today")]
        health = _FakeHealth()
        out = tmp_path / "test_report.pptx"

        with patch("core.ppt_generator.get_news_image", return_value=None):
            result = generate_executive_ppt(
                cards, health, "2026-02-15 09:00", 10, out)

        assert result.exists()
        assert result.suffix == ".pptx"
        assert result.stat().st_size > 0

    def test_generate_with_no_event_cards(self, tmp_path: Path):
        cards = [_make_non_event_card()]
        health = _FakeHealth()
        out = tmp_path / "empty_report.pptx"

        with patch("core.ppt_generator.get_news_image", return_value=None):
            result = generate_executive_ppt(
                cards, health, "2026-02-15 09:00", 1, out)

        assert result.exists()

    def test_generate_with_empty_cards(self, tmp_path: Path):
        health = _FakeHealth()
        out = tmp_path / "no_cards.pptx"

        with patch("core.ppt_generator.get_news_image", return_value=None):
            result = generate_executive_ppt(
                [], health, "2026-02-15 09:00", 0, out)

        assert result.exists()


class TestDarkTheme:
    """Test that dark theme colours are correctly applied."""

    def test_dark_theme_constants_exist(self):
        from core.ppt_generator import (
            ACCENT, BG_DARK, CARD_BG, HIGHLIGHT_YELLOW, SUBTLE_GRAY, TEXT_WHITE,
        )
        assert BG_DARK is not None
        assert TEXT_WHITE is not None
        assert HIGHLIGHT_YELLOW is not None
        assert CARD_BG is not None

    def test_slide_backgrounds_are_dark(self, tmp_path: Path):
        cards = [_make_event_card()]
        health = _FakeHealth()
        out = tmp_path / "dark_theme.pptx"

        with patch("core.ppt_generator.get_news_image", return_value=None):
            generate_executive_ppt(cards, health, "2026-02-15 09:00", 5, out, theme="dark")

        from pptx import Presentation
        prs = Presentation(str(out))
        for slide in prs.slides:
            bg_fill = slide.background.fill
            assert bg_fill.fore_color.rgb == BG_DARK, (
                f"Slide background should be BG_DARK (#121218), "
                f"got #{bg_fill.fore_color.rgb}"
            )


class TestEventSlideCount:
    """Test that each event card produces exactly 2 slides."""

    def test_two_slides_per_event(self, tmp_path: Path):
        n_events = 3
        cards = [
            _make_event_card(item_id=f"ev-{i}", title_plain=f"Microsoft launched Event {i} today")
            for i in range(n_events)
        ]
        health = _FakeHealth()
        out = tmp_path / "slide_count.pptx"

        with patch("core.ppt_generator.get_news_image", return_value=None):
            generate_executive_ppt(cards, health, "2026-02-15 09:00", n_events, out)

        from pptx import Presentation
        prs = Presentation(str(out))
        total_slides = len(prs.slides)

        # Expected: Cover(1) + StructuredSummary(1) + SignalThermometer(1) +
        #           CorpWatch(1) + KeyTakeaways(1) + OverviewTable(1) +
        #           EventRanking(1) + events(3*2=6) + RecommendedMoves(1) +
        #           DecisionMatrix(1) + PendingDecisions(1) = 16
        expected = 1 + 1 + 1 + 1 + 1 + 1 + 1 + (n_events * 2) + 1 + 1 + 1
        assert total_slides == expected, (
            f"Expected {expected} slides, got {total_slides}"
        )


class TestSlideContent:
    """Test that Q&A, Data card, Video, Sources content exists in slides."""

    def _get_all_text(self, prs) -> str:
        """Extract all text from presentation."""
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                texts.append(run.text.strip())
                        if para.text.strip():
                            texts.append(para.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                texts.append(cell.text.strip())
        return "\n".join(texts)

    def test_qa_content_present(self, tmp_path: Path):
        cards = [_make_event_card()]
        health = _FakeHealth()
        out = tmp_path / "qa_test.pptx"

        with patch("core.ppt_generator.get_news_image", return_value=None):
            generate_executive_ppt(cards, health, "2026-02-15 09:00", 5, out)

        from pptx import Presentation
        prs = Presentation(str(out))
        all_text = self._get_all_text(prs)

        assert "Q1" in all_text, "Q1 should be present in slides"
        assert "Q2" in all_text, "Q2 should be present in slides"
        assert "Q3" in all_text, "Q3 should be present in slides"
        assert "WHY IT MATTERS" in all_text

    def test_data_card_present(self, tmp_path: Path):
        cards = [_make_event_card()]
        health = _FakeHealth()
        out = tmp_path / "data_card_test.pptx"

        with patch("core.ppt_generator.get_news_image", return_value=None):
            generate_executive_ppt(cards, health, "2026-02-15 09:00", 5, out)

        from pptx import Presentation
        prs = Presentation(str(out))
        all_text = self._get_all_text(prs)

        # Data card should have numeric content (from build_data_card)
        assert any(c.isdigit() for c in all_text), (
            "Data card numeric values should be present"
        )

    def test_video_reference_present(self, tmp_path: Path):
        cards = [_make_event_card()]
        health = _FakeHealth()
        out = tmp_path / "video_test.pptx"

        with patch("core.ppt_generator.get_news_image", return_value=None):
            generate_executive_ppt(cards, health, "2026-02-15 09:00", 5, out)

        from pptx import Presentation
        prs = Presentation(str(out))
        all_text = self._get_all_text(prs)

        assert "Video" in all_text or "video" in all_text.lower(), (
            "Video reference should be present"
        )

    def test_sources_present(self, tmp_path: Path):
        cards = [_make_event_card()]
        health = _FakeHealth()
        out = tmp_path / "sources_test.pptx"

        with patch("core.ppt_generator.get_news_image", return_value=None):
            generate_executive_ppt(cards, health, "2026-02-15 09:00", 5, out)

        from pptx import Presentation
        prs = Presentation(str(out))
        all_text = self._get_all_text(prs)

        assert "Source" in all_text or "source" in all_text.lower(), (
            "Sources should be present"
        )

    def test_structured_summary_sections(self, tmp_path: Path):
        cards = [_make_event_card()]
        health = _FakeHealth()
        out = tmp_path / "summary_test.pptx"

        with patch("core.ppt_generator.get_news_image", return_value=None):
            generate_executive_ppt(cards, health, "2026-02-15 09:00", 5, out)

        from pptx import Presentation
        prs = Presentation(str(out))
        all_text = self._get_all_text(prs)

        for section in ["AI Trends", "Tech Landing", "Market Competition",
                        "Opportunities & Risks", "Recommended Actions"]:
            assert section in all_text, f"Section '{section}' should be present"

    def test_decision_matrix_present(self, tmp_path: Path):
        cards = [_make_event_card()]
        health = _FakeHealth()
        out = tmp_path / "matrix_test.pptx"

        with patch("core.ppt_generator.get_news_image", return_value=None):
            generate_executive_ppt(cards, health, "2026-02-15 09:00", 5, out)

        from pptx import Presentation
        prs = Presentation(str(out))
        all_text = self._get_all_text(prs)

        assert "Decision Matrix" in all_text, "Decision Matrix should be present"

    def test_signal_thermometer_present(self, tmp_path: Path):
        cards = [_make_event_card()]
        health = _FakeHealth()
        out = tmp_path / "signal_test.pptx"

        with patch("core.ppt_generator.get_news_image", return_value=None):
            generate_executive_ppt(cards, health, "2026-02-15 09:00", 5, out)

        from pptx import Presentation
        prs = Presentation(str(out))
        all_text = self._get_all_text(prs)

        assert "Signal Thermometer" in all_text, "Signal Thermometer should be present"
        assert "Market Heat Index" in all_text, "Market Heat Index should be present"

    def test_corp_watch_present(self, tmp_path: Path):
        cards = [_make_event_card()]
        health = _FakeHealth()
        out = tmp_path / "corp_test.pptx"

        with patch("core.ppt_generator.get_news_image", return_value=None):
            generate_executive_ppt(cards, health, "2026-02-15 09:00", 5, out)

        from pptx import Presentation
        prs = Presentation(str(out))
        all_text = self._get_all_text(prs)

        assert "Corp Watch" in all_text, "Corp Watch should be present"

    def test_event_ranking_present(self, tmp_path: Path):
        cards = [_make_event_card()]
        health = _FakeHealth()
        out = tmp_path / "ranking_test.pptx"

        with patch("core.ppt_generator.get_news_image", return_value=None):
            generate_executive_ppt(cards, health, "2026-02-15 09:00", 5, out)

        from pptx import Presentation
        prs = Presentation(str(out))
        all_text = self._get_all_text(prs)

        assert "Event Ranking" in all_text, "Event Ranking should be present"

    def test_recommended_moves_present(self, tmp_path: Path):
        cards = [_make_event_card()]
        health = _FakeHealth()
        out = tmp_path / "moves_test.pptx"

        with patch("core.ppt_generator.get_news_image", return_value=None):
            generate_executive_ppt(cards, health, "2026-02-15 09:00", 5, out)

        from pptx import Presentation
        prs = Presentation(str(out))
        all_text = self._get_all_text(prs)

        assert "Recommended Moves" in all_text, "Recommended Moves should be present"
