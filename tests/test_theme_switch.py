from __future__ import annotations

from pathlib import Path
from unittest.mock import patch

from pptx import Presentation

from core.ppt_generator import (
    DARK_BG,
    LIGHT_BG,
    MIN_FONT_SIZE_PT,
    MIN_LINE_SPACING,
    generate_executive_ppt,
)
from schemas.education_models import EduNewsCard, SystemHealthReport


def _cards() -> list[EduNewsCard]:
    return [
        EduNewsCard(
            item_id="event-001",
            is_valid_news=True,
            title_plain="NVIDIA launches new GPU architecture",
            what_happened="Launch event with architecture updates.",
            why_important="Impacts AI infrastructure pricing and capacity.",
            source_name="TechCrunch",
            source_url="https://example.com/news",
            final_score=8.0,
        )
    ]


def test_default_theme_is_light(tmp_path: Path) -> None:
    out = tmp_path / "light_default.pptx"
    health = SystemHealthReport(success_rate=80.0, p50_latency=2.0, p95_latency=5.0)

    with patch("core.ppt_generator.get_news_image", return_value=None):
        generate_executive_ppt(
            cards=_cards(),
            health=health,
            report_time="2026-02-15 09:00",
            total_items=1,
            output_path=out,
        )

    prs = Presentation(str(out))
    for slide in prs.slides:
        assert slide.background.fill.fore_color.rgb == LIGHT_BG


def test_dark_theme_can_generate(tmp_path: Path) -> None:
    out = tmp_path / "dark_theme.pptx"
    health = SystemHealthReport(success_rate=80.0, p50_latency=2.0, p95_latency=5.0)

    with patch("core.ppt_generator.get_news_image", return_value=None):
        generate_executive_ppt(
            cards=_cards(),
            health=health,
            report_time="2026-02-15 09:00",
            total_items=1,
            output_path=out,
            theme="dark",
        )

    prs = Presentation(str(out))
    for slide in prs.slides:
        assert slide.background.fill.fore_color.rgb == DARK_BG


def test_light_theme_readability_floor(tmp_path: Path) -> None:
    out = tmp_path / "light_readability.pptx"
    health = SystemHealthReport(success_rate=80.0, p50_latency=2.0, p95_latency=5.0)

    with patch("core.ppt_generator.get_news_image", return_value=None):
        generate_executive_ppt(
            cards=_cards(),
            health=health,
            report_time="2026-02-15 09:00",
            total_items=1,
            output_path=out,
        )

    prs = Presentation(str(out))
    font_sizes: list[float] = []
    line_spacings: list[float] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.font.size is not None:
                    font_sizes.append(float(paragraph.font.size.pt))
                if isinstance(paragraph.line_spacing, (int, float)):
                    line_spacings.append(float(paragraph.line_spacing))

    assert font_sizes, "Expected explicit font sizes in generated PPT"
    assert min(font_sizes) >= MIN_FONT_SIZE_PT
    assert line_spacings, "Expected explicit line spacing in generated PPT"
    assert min(line_spacings) >= MIN_LINE_SPACING
