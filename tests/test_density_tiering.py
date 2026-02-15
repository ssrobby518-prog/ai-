from __future__ import annotations

from pathlib import Path
from unittest.mock import patch

from docx import Document
from pptx import Presentation

from core.content_strategy import build_signal_summary
from core.doc_generator import generate_executive_docx
from core.info_density import apply_density_tiering
from core.ppt_generator import generate_executive_ppt
from schemas.education_models import EduNewsCard, SystemHealthReport


def _tier_b_cards() -> list[EduNewsCard]:
    return [
        EduNewsCard(
            item_id="tierb-001",
            is_valid_news=True,
            title_plain="OpenAI workflow update",
            what_happened="OpenAI copilots were adopted by support teams in Taipei.",
            why_important="Integration improved ticket triage consistency.",
            source_name="HackerNews",
            source_url="https://news.ycombinator.com/item?id=1",
            final_score=8.2,
        ),
        EduNewsCard(
            item_id="tierb-002",
            is_valid_news=True,
            title_plain="Google model operations shift",
            what_happened="Google Gemini rollout reached product squads in one week.",
            why_important="Benchmarking improved inference stability across teams.",
            source_name="TechCrunch",
            source_url="https://techcrunch.com/example",
            final_score=7.6,
        ),
        EduNewsCard(
            item_id="noise-001",
            is_valid_news=False,
            title_plain="Login required",
            what_happened="Please sign in to continue",
            invalid_reason="login_page",
            invalid_cause="blocked",
            source_name="Example",
            source_url="https://example.com/login",
            final_score=0.0,
        ),
    ]


def _extract_ppt_text(path: Path) -> str:
    prs = Presentation(str(path))
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            texts.append(cell.text.strip())
    return "\n".join(texts)


def _extract_doc_text(path: Path) -> str:
    doc = Document(str(path))
    texts: list[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            texts.append(p.text.strip())
    for t in doc.tables:
        for row in t.rows:
            for cell in row.cells:
                if cell.text.strip():
                    texts.append(cell.text.strip())
    return "\n".join(texts)


def test_density_tiering_routes_medium_density_into_tier_b() -> None:
    cards = [c for c in _tier_b_cards() if c.is_valid_news]
    tier_a, tier_b, tier_c, _ = apply_density_tiering(cards, "event")
    assert len(tier_a) == 0
    assert len(tier_b) >= 2
    assert len(tier_c) == 0
    assert all(getattr(c, "density_tier", "") == "B" for c in tier_b)


def test_event_zero_signal_top3_has_evidence_tokens() -> None:
    signals = build_signal_summary(_tier_b_cards())
    assert len(signals) >= 3
    for sig in signals[:3]:
        assert int(sig.get("platform_count", 0)) >= 1
        assert int(sig.get("heat_score", 0)) >= 30
        snippet = str(sig.get("example_snippet", "")).strip()
        assert 30 <= len(snippet) <= 120
        assert "source=unknown" not in snippet.lower()
        evidence_tokens = sig.get("evidence_tokens", [])
        assert isinstance(evidence_tokens, list)
        assert len(evidence_tokens) >= 2
        assert all(str(t).strip() for t in evidence_tokens)


def test_event_zero_deck_has_signal_content_without_unknown_source(tmp_path: Path) -> None:
    cards = _tier_b_cards()
    health = SystemHealthReport(success_rate=0.0, p50_latency=0.0, p95_latency=0.0)
    pptx_path = tmp_path / "tiering_no_event.pptx"
    docx_path = tmp_path / "tiering_no_event.docx"

    with patch("core.ppt_generator.get_news_image", return_value=None), patch(
        "core.doc_generator.get_news_image",
        return_value=None,
    ):
        generate_executive_ppt(
            cards=cards,
            health=health,
            report_time="2026-02-15 09:00",
            total_items=len(cards),
            output_path=pptx_path,
            metrics={
                "fetched_total": 42,
                "gate_pass_total": 2,
                "after_filter_total": 2,
                "sources_total": 2,
            },
        )
        generate_executive_docx(
            cards=cards,
            health=health,
            report_time="2026-02-15 09:00",
            total_items=len(cards),
            output_path=docx_path,
            metrics={
                "fetched_total": 42,
                "gate_pass_total": 2,
                "after_filter_total": 2,
                "sources_total": 2,
            },
        )

    ppt_text = _extract_ppt_text(pptx_path).lower()
    doc_text = _extract_doc_text(docx_path).lower()
    assert "signal thermometer" in ppt_text
    assert "corp watch" in ppt_text
    assert "source=unknown" not in ppt_text
    assert "source=unknown" not in doc_text
    assert "last july was" not in ppt_text
    assert "last july was" not in doc_text
